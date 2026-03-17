"""Perspective Combination Dialog - Multi-image comparison for defect detection.

This dialog allows users to:
1. Select a base image and multiple compare images
2. Align and blend compare images
3. Subtract from base to reveal defect hints
4. View SNR map highlighting strong signal areas
"""
from __future__ import annotations

from typing import Dict, List, Optional, Tuple
from pathlib import Path

from PySide6 import QtWidgets, QtCore, QtGui
from PySide6.QtCore import Qt, Signal

import numpy as np
import cv2

from matplotlib.backends.backend_qtagg import FigureCanvasQTAgg as FigureCanvas
from matplotlib.figure import Figure

from ..core.ebeam_snr import EbeamCondition, load_image_gray
from ..core.perspective_combine import (
    align_and_subtract,
    compute_single_pair,
    compute_multi_pairs,
    compute_quadrant_fusion,
    compute_roi_full_stats,
    colorize_snr_map,
    CombineResult,
    SinglePairResult,
    QuadrantFusionResult,
)
from ..core.roi_set import MultiROISet, NamedROI, ROIFullResult

# 使用 design_tokens 統一配色（遵循 AGENTS.md 規範）
from .design_tokens import Colors, Typography, Spacing, BorderRadius
from .welcome_tutorial import WelcomeTutorialOverlay, should_show_tutorial, mark_tutorial_completed

UI_PRIMARY = Colors.BRAND_PRIMARY
UI_PRIMARY_HOVER = Colors.BRAND_PRIMARY_HOVER
UI_PRIMARY_PRESSED = Colors.BRAND_PRIMARY_PRESSED
UI_PRIMARY_SOFT = Colors.BRAND_PRIMARY_SOFT

UI_BG_WINDOW = Colors.BG_WINDOW
UI_BG_PANEL = Colors.BG_PANEL
UI_BG_CARD = Colors.BG_CARD
UI_BG_SUBTLE = Colors.BG_SUBTLE
UI_BG_INPUT = Colors.BG_INPUT
UI_BG_VIEWER = Colors.BG_VIEWER

UI_TEXT = Colors.TEXT_PRIMARY
UI_TEXT_SECONDARY = Colors.TEXT_SECONDARY
UI_TEXT_MUTED = Colors.TEXT_MUTED
UI_TEXT_INVERSE = Colors.TEXT_INVERSE
UI_TEXT_ON_PRIMARY = Colors.TEXT_ON_PRIMARY

UI_BORDER = Colors.BORDER_DEFAULT
UI_BORDER_HOVER = Colors.BORDER_HOVER
UI_BORDER_ACTIVE = Colors.BORDER_ACTIVE

UI_SUCCESS = Colors.SUCCESS
UI_WARNING = Colors.WARNING
UI_INFO = Colors.INFO

UI_WINDOW_BG_LIGHT = "#F3F4F6"
UI_LEFT_PANEL_BG = "#F9FAFB"
UI_VIEWER_BG = "#FFFFFF"
UI_TEXT_PRIMARY_STRONG = "#111827"
UI_TEXT_SECONDARY_MUTED = "#6B7280"
UI_ACCENT_HOVER = "#FB923C"
UI_ACCENT_LIGHT = "#FFF4E5"


def _hex_to_bgr(color: str) -> Tuple[int, int, int]:
    """Convert #RRGGBB color to BGR tuple for OpenCV drawing."""
    c = color.lstrip("#")
    return (int(c[4:6], 16), int(c[2:4], 16), int(c[0:2], 16))


DIALOG_STYLE = f"""
    QDialog {{
        background-color: {UI_WINDOW_BG_LIGHT};
        color: {UI_TEXT};
        font-family: {Typography.FONT_FAMILY};
        font-size: {Typography.FONT_SIZE_BODY};
    }}

    QGroupBox {{
        font-weight: {Typography.FONT_WEIGHT_SEMIBOLD};
        font-size: {Typography.FONT_SIZE_BODY};
        border: 1px solid {UI_BORDER};
        border-radius: {BorderRadius.MD};
        background-color: {UI_BG_PANEL};
        margin-top: 14px;
        padding: 16px 12px 12px 12px;
    }}
    QGroupBox::title {{
        subcontrol-origin: margin;
        left: 10px;
        top: 2px;
        padding: 0 6px;
        color: {UI_TEXT};
        background-color: transparent;
    }}
    QGroupBox#AdvancedSettings {{
        background-color: {UI_BG_CARD};
    }}
    QGroupBox#AdvancedSettings::title {{
        color: {UI_TEXT_SECONDARY};
    }}
    QFrame#TopToolbar {{
        background-color: {UI_BG_PANEL};
        border: 1px solid {UI_BORDER};
        border-radius: {BorderRadius.MD};
    }}
    QFrame#MainToolbar {{
        background-color: {UI_BG_PANEL};
        border: 1px solid {UI_BORDER};
        border-radius: {BorderRadius.MD};
        padding: 4px 8px;
    }}
    QFrame#BottomCard {{
        background-color: {UI_BG_PANEL};
        border: 1px solid {UI_BORDER};
        border-radius: {BorderRadius.MD};
    }}
    QLabel#SectionTitle {{
        color: {UI_TEXT};
        font-size: {Typography.FONT_SIZE_H3};
        font-weight: {Typography.FONT_WEIGHT_BOLD};
        padding: 0px;
        border: none;
        background: transparent;
    }}
    QLabel#SectionSeparator {{
        background-color: {UI_PRIMARY};
        min-height: 2px;
        max-height: 2px;
        border: none;
    }}
    QPushButton#ToolbarAction {{
        background-color: {UI_BG_PANEL};
        color: {UI_TEXT};
        border: 1px solid #D1D5DB;
        border-radius: {BorderRadius.SM};
        padding: 6px 14px;
        min-height: 32px;
        min-width: 80px;
        font-weight: {Typography.FONT_WEIGHT_MEDIUM};
        font-size: {Typography.FONT_SIZE_SMALL};
    }}
    QPushButton#ToolbarAction:hover {{
        background-color: #FFF8ED;
        border-color: {UI_PRIMARY};
    }}
    QPushButton#ToolbarSecondary {{
        background-color: #FEF3C7;
        color: #92400E;
        border: 1.5px solid {UI_PRIMARY};
        border-radius: {BorderRadius.SM};
        padding: 6px 14px;
        min-height: 32px;
        min-width: 80px;
        font-weight: {Typography.FONT_WEIGHT_SEMIBOLD};
        font-size: {Typography.FONT_SIZE_SMALL};
    }}
    QPushButton#ToolbarSecondary:hover {{
        background-color: #FDE68A;
        border-color: #D97706;
    }}
    QPushButton#ToolbarSecondary:disabled {{
        background-color: {UI_BG_SUBTLE};
        color: {UI_TEXT_MUTED};
        border-color: {UI_BORDER};
    }}
    QPushButton#ToolbarPrimary {{
        background-color: {UI_PRIMARY};
        color: {UI_TEXT_ON_PRIMARY};
        border: 1px solid {UI_PRIMARY};
        border-radius: {BorderRadius.SM};
        padding: 10px 24px;
        min-height: 44px;
        font-weight: {Typography.FONT_WEIGHT_BOLD};
        font-size: {Typography.FONT_SIZE_BODY};
        letter-spacing: 0.3px;
    }}
    QPushButton#ToolbarPrimary:hover {{
        background-color: {UI_PRIMARY_HOVER};
        border-color: {UI_PRIMARY_HOVER};
    }}
    QPushButton[viewerMode="true"] {{
        background-color: {UI_BG_PANEL};
        color: {UI_TEXT_SECONDARY};
        border: 1px solid {UI_BORDER};
        border-radius: 0px;
        padding: 6px 16px;
        min-height: 30px;
        min-width: 72px;
        font-size: {Typography.FONT_SIZE_SMALL};
        font-weight: {Typography.FONT_WEIGHT_MEDIUM};
    }}
    QPushButton[viewerMode="true"]:hover {{
        background-color: #FFF8ED;
        color: {UI_TEXT};
    }}
    QPushButton[viewerMode="true"]:checked {{
        background-color: {UI_PRIMARY};
        color: {UI_TEXT_ON_PRIMARY};
        border-color: {UI_PRIMARY};
        font-weight: {Typography.FONT_WEIGHT_BOLD};
    }}
    QPushButton#ViewerModeFirst {{
        border-top-left-radius: {BorderRadius.SM};
        border-bottom-left-radius: {BorderRadius.SM};
    }}
    QPushButton#ViewerModeLast {{
        border-top-right-radius: {BorderRadius.SM};
        border-bottom-right-radius: {BorderRadius.SM};
    }}

    QLabel {{
        color: {UI_TEXT};
        font-size: {Typography.FONT_SIZE_SMALL};
    }}
    QLabel[toolbarLabel="true"] {{
        color: {UI_TEXT_SECONDARY};
        font-size: {Typography.FONT_SIZE_SMALL};
        font-weight: {Typography.FONT_WEIGHT_MEDIUM};
    }}
    QLabel[secondary="true"] {{
        color: {UI_TEXT_SECONDARY};
        font-size: {Typography.FONT_SIZE_CAPTION};
    }}
    QLabel[statLabel="true"] {{
        color: {UI_TEXT_SECONDARY};
        font-size: {Typography.FONT_SIZE_SMALL};
    }}
    QLabel[statValue="true"] {{
        color: {UI_TEXT};
        font-size: {Typography.FONT_SIZE_BODY};
        font-weight: {Typography.FONT_WEIGHT_BOLD};
    }}
    QLabel[statValueLong="true"] {{
        color: {UI_TEXT_SECONDARY};
        font-size: {Typography.FONT_SIZE_SMALL};
        font-weight: {Typography.FONT_WEIGHT_SEMIBOLD};
    }}

    QPushButton {{
        background-color: {UI_BG_PANEL};
        color: {UI_TEXT};
        border: 1px solid {UI_BORDER};
        border-radius: {BorderRadius.MD};
        padding: 8px 14px;
        min-height: 34px;
        min-width: 92px;
        font-weight: {Typography.FONT_WEIGHT_MEDIUM};
        font-size: {Typography.FONT_SIZE_SMALL};
    }}
    QPushButton:hover {{
        background-color: #FFF8ED;
        border-color: {UI_PRIMARY};
    }}
    QPushButton:pressed {{
        background-color: #FDE7C2;
        border-color: {UI_PRIMARY_PRESSED};
    }}
    QPushButton:disabled {{
        background-color: {UI_BG_SUBTLE};
        color: {UI_TEXT_MUTED};
        border-color: {UI_BORDER};
    }}
    QPushButton[variant="primary"] {{
        background-color: {UI_PRIMARY};
        color: {UI_TEXT_ON_PRIMARY};
        border: 1px solid {UI_PRIMARY};
        font-weight: {Typography.FONT_WEIGHT_BOLD};
    }}
    QPushButton[variant="primary"]:hover {{
        background-color: {UI_PRIMARY_HOVER};
        border-color: {UI_PRIMARY_HOVER};
    }}
    QPushButton[variant="primary"]:pressed {{
        background-color: {UI_PRIMARY_PRESSED};
        border-color: {UI_PRIMARY_PRESSED};
    }}
    QPushButton[variant="ghost"] {{
        background-color: transparent;
        border-color: transparent;
    }}
    QPushButton[variant="ghost"]:hover {{
        background-color: #FFF8ED;
        border-color: {UI_BORDER};
    }}
    QPushButton#ViewerToolToggle {{
        background-color: {UI_BG_PANEL};
        color: {UI_TEXT_SECONDARY};
        border: 1px solid {UI_BORDER};
        border-radius: {BorderRadius.SM};
        padding: 6px 12px;
        min-height: 30px;
        min-width: 96px;
        font-size: {Typography.FONT_SIZE_SMALL};
        font-weight: {Typography.FONT_WEIGHT_MEDIUM};
    }}
    QPushButton#ViewerToolToggle:hover {{
        background-color: #FFF8ED;
        border-color: {UI_PRIMARY};
    }}
    QPushButton#ViewerToolToggle:checked {{
        background-color: {UI_PRIMARY};
        color: {UI_TEXT_ON_PRIMARY};
        border-color: {UI_PRIMARY};
        font-weight: {Typography.FONT_WEIGHT_BOLD};
    }}
    QLabel[toolbarGroupLabel="true"] {{
        color: {UI_TEXT_MUTED};
        font-size: {Typography.FONT_SIZE_CAPTION};
        font-weight: {Typography.FONT_WEIGHT_BOLD};
        letter-spacing: 0.35px;
    }}
    QPushButton[toolbarToggle="true"] {{
        background-color: {UI_BG_SUBTLE};
        border: 1px solid {UI_BORDER};
        border-radius: {BorderRadius.SM};
        padding: 6px 12px;
        min-height: 32px;
        min-width: 88px;
        font-size: {Typography.FONT_SIZE_SMALL};
    }}
    QPushButton[toolbarToggle="true"]:hover {{
        background-color: #FFF8ED;
        border-color: {UI_PRIMARY};
    }}
    QPushButton[toolbarToggle="true"]:checked {{
        background-color: {UI_PRIMARY};
        color: {UI_TEXT_ON_PRIMARY};
        border-color: {UI_PRIMARY};
        font-weight: {Typography.FONT_WEIGHT_BOLD};
    }}

    QComboBox {{
        background-color: {UI_BG_INPUT};
        color: {UI_TEXT};
        border: 1px solid {UI_BORDER};
        border-radius: {BorderRadius.SM};
        padding: {Spacing.INPUT_PADDING};
        font-size: {Typography.FONT_SIZE_SMALL};
        min-height: 28px;
    }}
    QComboBox:hover, QComboBox:focus {{
        border-color: {UI_PRIMARY};
    }}
    QComboBox::drop-down {{
        border: none;
        padding-right: 8px;
    }}
    QComboBox QAbstractItemView {{
        background-color: {UI_BG_PANEL};
        color: {UI_TEXT};
        selection-background-color: {UI_PRIMARY};
        selection-color: {UI_TEXT_ON_PRIMARY};
        border: 1px solid {UI_BORDER};
    }}

    QCheckBox {{
        color: {UI_TEXT};
        spacing: {Spacing.SM};
        font-size: {Typography.FONT_SIZE_SMALL};
    }}
    QCheckBox::indicator {{
        width: 18px;
        height: 18px;
        border: 2px solid #F7C267;
        border-radius: 4px;
        background-color: #FFF8ED;
    }}
    QCheckBox::indicator:checked {{
        background-color: {UI_PRIMARY};
        border-color: {UI_PRIMARY};
    }}

    QSpinBox, QDoubleSpinBox, QLineEdit {{
        background-color: {UI_BG_INPUT};
        color: {UI_TEXT};
        border: 1px solid {UI_BORDER};
        border-radius: {BorderRadius.SM};
        padding: 5px 26px 5px 10px;
        min-height: 24px;
        font-size: {Typography.FONT_SIZE_SMALL};
    }}
    QSpinBox:hover, QDoubleSpinBox:hover, QLineEdit:hover,
    QSpinBox:focus, QDoubleSpinBox:focus, QLineEdit:focus {{
        border-color: {UI_PRIMARY};
    }}
    QSpinBox::up-button, QDoubleSpinBox::up-button,
    QSpinBox::down-button, QDoubleSpinBox::down-button {{
        subcontrol-origin: border;
        width: 14px;
        border-left: 1px solid {UI_BORDER};
        background-color: {UI_BG_CARD};
    }}
    QSpinBox::up-button, QDoubleSpinBox::up-button {{
        subcontrol-position: top right;
        border-top-right-radius: {BorderRadius.SM};
    }}
    QSpinBox::down-button, QDoubleSpinBox::down-button {{
        subcontrol-position: bottom right;
        border-bottom-right-radius: {BorderRadius.SM};
        border-top: 1px solid {UI_BORDER};
    }}
    QSpinBox::up-button:hover, QDoubleSpinBox::up-button:hover,
    QSpinBox::down-button:hover, QDoubleSpinBox::down-button:hover {{
        background-color: #FFF8ED;
    }}
    QSpinBox::up-arrow, QDoubleSpinBox::up-arrow,
    QSpinBox::down-arrow, QDoubleSpinBox::down-arrow {{
        width: 7px;
        height: 7px;
    }}

    QSlider::groove:horizontal {{
        border: none;
        height: 8px;
        background: {UI_BG_SUBTLE};
        border-radius: 4px;
    }}
    QSlider::handle:horizontal {{
        background: {UI_PRIMARY};
        width: 16px;
        height: 16px;
        margin: -5px 0;
        border-radius: 8px;
    }}
    QSlider::handle:horizontal:hover {{
        background: {UI_PRIMARY_HOVER};
    }}
    QSlider::sub-page:horizontal {{
        background: {UI_PRIMARY};
        border-radius: 4px;
    }}

    QScrollArea {{
        border: none;
        background-color: transparent;
    }}
    QScrollBar:vertical {{
        background-color: {UI_BG_SUBTLE};
        width: 8px;
        margin: 0;
        border-radius: 4px;
    }}
    QScrollBar::handle:vertical {{
        background-color: #C7D0DC;
        min-height: 30px;
        border-radius: 4px;
    }}
    QScrollBar::handle:vertical:hover {{
        background-color: {UI_PRIMARY};
    }}
    QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical {{
        height: 0;
    }}

    QFrame {{
        border: none;
    }}
"""


class SyncZoomImageWidget(QtWidgets.QWidget):
    """Image widget with synchronized zoom on hover.

    Signals:
        cursor_moved: Emitted when cursor moves, sends (norm_x, norm_y)
        cursor_left: Emitted when cursor leaves the widget
        roi_selected: Emitted when user finishes drawing ROI (norm_x, norm_y, norm_w, norm_h)
    """

    cursor_moved = Signal(float, float)  # Normalized position (0-1)
    cursor_left = Signal()
    roi_selected = Signal(float, float, float, float)  # norm x, y, w, h (legacy single-ROI)
    clicked = Signal()
    # Multi-ROI signals
    multi_roi_drawn = Signal(float, float, float, float)  # norm x, y, w, h — drag or single-add
    multi_roi_anchor = Signal(str, float, float)          # 'tl'/'br', norm cx, cy

    ZOOM_FACTOR = 2.0  # Magnification factor (lower = more FOV)
    ZOOM_SIZE = 220  # Larger circular FOV for easier inspection

    def __init__(self, title: str = "", parent=None):
        super().__init__(parent)
        self._title = title
        self._image: Optional[np.ndarray] = None
        self._hint_rect: Optional[tuple] = None
        self._hint_info: str = ""
        self._show_hint: bool = False
        self._cursor_pos: Optional[tuple] = None  # (norm_x, norm_y)
        self._show_zoom: bool = False
        self._partner: Optional['SyncZoomImageWidget'] = None  # Linked partner
        # Legacy single-ROI drawing state (used by ROI Nulling)
        self._roi_mode: bool = False
        self._roi_start: Optional[tuple] = None  # (norm_x, norm_y)
        self._roi_current: Optional[tuple] = None  # (norm_x, norm_y)
        self._roi_dragging: bool = False
        # Active ROI overlay (drawn permanently until cleared, legacy single-ROI)
        self._active_roi: Optional[tuple] = None  # (norm_x, norm_y, norm_w, norm_h)
        # Multi-ROI state
        self._multi_roi_set: Optional[MultiROISet] = None
        # Draw mode: 'idle' | 'drag' | 'single_add' | 'multi_add'
        self._multi_draw_mode: str = 'idle'
        self._multi_roi_start: Optional[tuple] = None   # drag start (norm x, y)
        self._multi_roi_current: Optional[tuple] = None # drag current (norm x, y)
        self._multi_roi_dragging: bool = False
        self._add_size_norm: Tuple[float, float] = (0.05, 0.05)  # w, h for single_add
        self._grid_anchor_tl: Optional[Tuple[float, float]] = None  # norm cx, cy
        self._grid_anchor_br: Optional[Tuple[float, float]] = None  # norm cx, cy
        self._grid_preview_rects: List[Tuple] = []   # norm_rects for preview (yellow dashed)

        self.setMinimumSize(350, 350)  # Larger minimum for better visibility
        self.setMouseTracking(True)
        self.setStyleSheet(f"""
            QWidget {{
                background-color: {UI_BG_PANEL};
                border: 1px solid {UI_BORDER};
                border-radius: 12px;
            }}
            QLabel {{
                background-color: {UI_BG_VIEWER};
                border: none;
                border-radius: 8px;
                color: #D1D5DB;
            }}
        """)

        # Image label
        self._label = QtWidgets.QLabel(self)
        self._label.setAlignment(Qt.AlignCenter)
        self._label.setStyleSheet(f"border: none; background: {UI_BG_VIEWER}; border-radius: 8px;")
        self._label.setMouseTracking(True)  # Enable hover without clicking

        layout = QtWidgets.QVBoxLayout(self)
        layout.setContentsMargins(4, 4, 4, 4)
        layout.addWidget(self._label)

    def linkPartner(self, partner: 'SyncZoomImageWidget'):
        """Link to another widget for synchronized cursor display."""
        self._partner = partner

    def setImage(self, image: np.ndarray):
        """Set image to display."""
        self._image = image
        self._update_display()

    def setHint(self, rect: tuple = None, info: str = "", show: bool = True):
        """Set hint rectangle and info."""
        self._hint_rect = rect
        self._hint_info = info
        self._show_hint = show
        self._update_display()

    def setShowHint(self, show: bool):
        """Toggle hint visibility."""
        self._show_hint = show
        self._update_display()

    def set_roi_mode(self, enabled: bool):
        """Enter or exit ROI drawing mode."""
        self._roi_mode = enabled
        self._roi_start = None
        self._roi_current = None
        self._roi_dragging = False
        self.setCursor(Qt.CrossCursor if enabled else Qt.ArrowCursor)
        self._update_display()

    def set_active_roi(self, norm_rect: Optional[tuple]):
        """Set the persistent ROI overlay (norm_x, norm_y, norm_w, norm_h), or None to clear."""
        self._active_roi = norm_rect
        self._update_display()

    # ------------------------------------------------------------------
    # Multi-ROI public API
    # ------------------------------------------------------------------

    def set_multi_roi_set(self, roi_set: Optional[MultiROISet]) -> None:
        """Attach a MultiROISet whose ROIs are drawn on this widget."""
        self._multi_roi_set = roi_set
        self._update_display()

    def set_multi_draw_mode(
        self,
        mode: str,
        add_size_norm: Optional[Tuple[float, float]] = None,
    ) -> None:
        """Set the active drawing mode.

        Parameters
        ----------
        mode : 'idle' | 'drag' | 'single_add' | 'multi_add'
        add_size_norm : (norm_w, norm_h) used when mode == 'single_add'.
        """
        self._multi_draw_mode = mode
        if add_size_norm is not None:
            self._add_size_norm = add_size_norm
        self._multi_roi_start = None
        self._multi_roi_current = None
        self._multi_roi_dragging = False
        self.setCursor(
            Qt.CrossCursor if mode != 'idle' else Qt.ArrowCursor
        )
        self._update_display()

    def clear_multi_draw_preview(self) -> None:
        """Clear transient ROI-draw previews and hover state."""
        self._multi_roi_start = None
        self._multi_roi_current = None
        self._multi_roi_dragging = False
        self._cursor_pos = None
        self._show_zoom = False
        self._update_display()

    def set_grid_preview(self, rects: List[Tuple]) -> None:
        """Show yellow dashed preview rectangles (norm_rects) for Multi-Add grid."""
        self._grid_preview_rects = rects
        self._update_display()

    def clear_grid_preview(self) -> None:
        self._grid_preview_rects = []
        self._update_display()

    def set_grid_anchors(
        self,
        tl: Optional[Tuple[float, float]],
        br: Optional[Tuple[float, float]],
    ) -> None:
        """Update displayed anchor markers without entering anchor-set mode."""
        self._grid_anchor_tl = tl
        self._grid_anchor_br = br
        self._update_display()

    def _widget_to_norm(self, pos) -> Optional[tuple]:
        """Convert widget-space QPoint to normalized image coordinates (0-1)."""
        label_rect = self._label.geometry()
        p = pos - label_rect.topLeft()
        pixmap = self._label.pixmap()
        if pixmap is None or pixmap.isNull():
            return None
        px_w, px_h = pixmap.width(), pixmap.height()
        lbl_w, lbl_h = label_rect.width(), label_rect.height()
        offset_x = (lbl_w - px_w) // 2
        offset_y = (lbl_h - px_h) // 2
        img_x = p.x() - offset_x
        img_y = p.y() - offset_y
        if 0 <= img_x < px_w and 0 <= img_y < px_h:
            return float(img_x) / px_w, float(img_y) / px_h
        return None

    def _widget_to_norm_clamped(self, pos) -> Optional[tuple]:
        """Convert widget-space QPoint to normalized coords, clamped to the image edge."""
        label_rect = self._label.geometry()
        p = pos - label_rect.topLeft()
        pixmap = self._label.pixmap()
        if pixmap is None or pixmap.isNull():
            return None
        px_w, px_h = pixmap.width(), pixmap.height()
        lbl_w, lbl_h = label_rect.width(), label_rect.height()
        offset_x = (lbl_w - px_w) // 2
        offset_y = (lbl_h - px_h) // 2
        if px_w <= 0 or px_h <= 0:
            return None
        img_x = min(max(p.x() - offset_x, 0), px_w - 1)
        img_y = min(max(p.y() - offset_y, 0), px_h - 1)
        den_w = max(px_w - 1, 1)
        den_h = max(px_h - 1, 1)
        return float(img_x) / den_w, float(img_y) / den_h

    def setCursorPos(self, norm_x: float, norm_y: float):
        """Set cursor position from external source (partner sync)."""
        self._cursor_pos = (norm_x, norm_y)
        self._show_zoom = True
        self._update_display()

    def clearCursor(self):
        """Clear cursor when partner cursor leaves."""
        self._cursor_pos = None
        self._show_zoom = False
        self._update_display()

    def mousePressEvent(self, event):
        """Start ROI drawing on left click when in roi_mode or multi draw mode."""
        if event.button() == Qt.LeftButton:
            self.clicked.emit()

        # Legacy single-ROI drag (ROI Nulling)
        if self._roi_mode and event.button() == Qt.LeftButton:
            norm = self._widget_to_norm_clamped(event.pos())
            if norm:
                self._roi_start = norm
                self._roi_current = norm
                self._roi_dragging = True
            event.accept()
            return

        # Multi-ROI modes
        if self._multi_draw_mode != 'idle' and event.button() == Qt.LeftButton:
            norm = self._widget_to_norm_clamped(event.pos())
            if norm is None:
                super().mousePressEvent(event)
                return

            if self._multi_draw_mode == 'single_add':
                # Place ROI centered on click point
                nw, nh = self._add_size_norm
                nx = max(0.0, min(norm[0] - nw / 2, 1.0 - nw))
                ny = max(0.0, min(norm[1] - nh / 2, 1.0 - nh))
                self.multi_roi_drawn.emit(nx, ny, nw, nh)
                # Stay in single_add mode so user can keep placing

            elif self._multi_draw_mode == 'drag':
                self._multi_roi_start = norm
                self._multi_roi_current = norm
                self._multi_roi_dragging = True

            elif self._multi_draw_mode == 'multi_add':
                anchor_type = 'tl' if self._grid_anchor_tl is None else 'br'
                cx, cy = norm
                if anchor_type == 'tl':
                    self._grid_anchor_tl = (cx, cy)
                else:
                    self._grid_anchor_br = (cx, cy)
                self.multi_roi_anchor.emit(anchor_type, cx, cy)
                # Exit after second click (BR) to prevent accidental re-anchoring.
                if anchor_type == 'br':
                    self._multi_draw_mode = 'idle'
                    self.setCursor(Qt.ArrowCursor)

            self._update_display()
            event.accept()
            return

        super().mousePressEvent(event)

    def mouseReleaseEvent(self, event):
        """Finish ROI drawing on left release."""
        # Legacy single-ROI
        if self._roi_mode and self._roi_dragging and event.button() == Qt.LeftButton:
            norm = self._widget_to_norm_clamped(event.pos())
            if norm:
                self._roi_current = norm
            if self._roi_start and self._roi_current:
                x0, y0 = self._roi_start
                x1, y1 = self._roi_current
                nx, ny = min(x0, x1), min(y0, y1)
                nw, nh = abs(x1 - x0), abs(y1 - y0)
                if nw > 0.01 and nh > 0.01:
                    self._active_roi = (nx, ny, nw, nh)
                    self.roi_selected.emit(nx, ny, nw, nh)
            self._roi_dragging = False
            self._roi_mode = False
            self.setCursor(Qt.ArrowCursor)
            self._update_display()
            event.accept()
            return

        # Multi-ROI drag
        if (self._multi_draw_mode == 'drag' and self._multi_roi_dragging
                and event.button() == Qt.LeftButton):
            norm = self._widget_to_norm_clamped(event.pos())
            if norm:
                self._multi_roi_current = norm
            if self._multi_roi_start and self._multi_roi_current:
                x0, y0 = self._multi_roi_start
                x1, y1 = self._multi_roi_current
                nx, ny = min(x0, x1), min(y0, y1)
                nw, nh = abs(x1 - x0), abs(y1 - y0)
                if nw > 0.01 and nh > 0.01:
                    self.multi_roi_drawn.emit(nx, ny, nw, nh)
            self._multi_roi_dragging = False
            self._multi_roi_start = None
            self._multi_roi_current = None
            # Stay in drag mode so user can keep drawing
            self._update_display()
            event.accept()
            return

        super().mouseReleaseEvent(event)

    def mouseMoveEvent(self, event):
        """Track mouse position and emit signal."""
        # Legacy single-ROI drag rubber band
        if self._roi_mode and self._roi_dragging:
            norm = self._widget_to_norm_clamped(event.pos())
            if norm:
                self._roi_current = norm
            self._update_display()
            event.accept()
            return

        # Multi-ROI drag rubber band
        if self._multi_draw_mode == 'drag' and self._multi_roi_dragging:
            norm = self._widget_to_norm_clamped(event.pos())
            if norm:
                self._multi_roi_current = norm
            self._update_display()
            event.accept()
            return

        if self._image is None:
            return

        # Get label geometry
        label_rect = self._label.geometry()
        pos = event.pos() - label_rect.topLeft()

        # Calculate normalized position within the displayed image
        pixmap = self._label.pixmap()
        if pixmap is None or pixmap.isNull():
            return

        # Image is centered in label
        px_w, px_h = pixmap.width(), pixmap.height()
        lbl_w, lbl_h = label_rect.width(), label_rect.height()
        offset_x = (lbl_w - px_w) // 2
        offset_y = (lbl_h - px_h) // 2

        img_x = pos.x() - offset_x
        img_y = pos.y() - offset_y

        clamp_to_edge = self._multi_draw_mode in ('single_add', 'multi_add')
        if clamp_to_edge:
            den_w = max(px_w - 1, 1)
            den_h = max(px_h - 1, 1)
            norm_x = min(max(img_x, 0), px_w - 1) / den_w
            norm_y = min(max(img_y, 0), px_h - 1) / den_h
            self._cursor_pos = (norm_x, norm_y)
            self._show_zoom = True
            self.cursor_moved.emit(norm_x, norm_y)
        elif 0 <= img_x < px_w and 0 <= img_y < px_h:
            norm_x = img_x / px_w
            norm_y = img_y / px_h
            self._cursor_pos = (norm_x, norm_y)
            self._show_zoom = True
            self.cursor_moved.emit(norm_x, norm_y)
        else:
            self._cursor_pos = None
            self._show_zoom = False

        self._update_display()

    def leaveEvent(self, event):
        """Handle mouse leaving widget."""
        self._cursor_pos = None
        self._show_zoom = False
        self.cursor_left.emit()
        self._update_display()

    def _update_display(self):
        """Update the displayed pixmap."""
        if self._image is None:
            self._label.setText(
                f"{self._title}\n\nNo image loaded\n"
                f"Load a folder and select images to begin."
            )
            self._label.setStyleSheet(
                f"border: none; background: {UI_BG_VIEWER}; border-radius: 8px;"
                f" color: {UI_TEXT_MUTED}; font-size: {Typography.FONT_SIZE_SMALL};"
            )
            return

        h, w = self._image.shape[:2]

        # Convert to RGB
        if len(self._image.shape) == 2:
            img_display = self._image
            if img_display.dtype != np.uint8:
                img_display = np.clip(img_display, 0, 255).astype(np.uint8)
            img_rgb = cv2.cvtColor(img_display, cv2.COLOR_GRAY2BGR)
        else:
            img_rgb = self._image.copy()

        # Draw hint box if enabled
        if self._show_hint and self._hint_rect is not None:
            x, y, rw, rh = self._hint_rect
            hint_color = _hex_to_bgr(UI_SUCCESS)
            cv2.rectangle(img_rgb, (x, y), (x + rw, y + rh), hint_color, 2)
            cx, cy = x + rw // 2, y + rh // 2
            mark_len = 10
            cv2.line(img_rgb, (cx - mark_len, cy), (cx + mark_len, cy), hint_color, 1)
            cv2.line(img_rgb, (cx, cy - mark_len), (cx, cy + mark_len), hint_color, 1)
            if self._hint_info:
                font = cv2.FONT_HERSHEY_SIMPLEX
                text_y = max(y - 8, 18)
                cv2.putText(img_rgb, self._hint_info, (x + 1, text_y + 1),
                            font, 0.45, (0, 0, 0), 2, cv2.LINE_AA)
                cv2.putText(img_rgb, self._hint_info, (x, text_y),
                            font, 0.45, hint_color, 1, cv2.LINE_AA)

        # Draw active ROI (persistent cyan box)
        if self._active_roi is not None:
            nx, ny, nw, nh = self._active_roi
            rx, ry = int(nx * w), int(ny * h)
            rw2, rh2 = max(1, int(nw * w)), max(1, int(nh * h))
            cv2.rectangle(img_rgb, (rx, ry), (rx + rw2, ry + rh2), (0, 255, 255), 2)
            cv2.rectangle(img_rgb, (rx + 1, ry + 1), (rx + rw2 - 1, ry + rh2 - 1), (0, 0, 0), 1)
            cv2.putText(img_rgb, "ROI", (rx + 3, ry + 14),
                        cv2.FONT_HERSHEY_SIMPLEX, 0.4, (0, 255, 255), 1, cv2.LINE_AA)

        # Draw legacy rubber band while dragging ROI
        if self._roi_mode and self._roi_start and self._roi_current:
            x0n, y0n = self._roi_start
            x1n, y1n = self._roi_current
            rx = int(min(x0n, x1n) * w)
            ry = int(min(y0n, y1n) * h)
            rw3 = max(1, int(abs(x1n - x0n) * w))
            rh3 = max(1, int(abs(y1n - y0n) * h))
            cv2.rectangle(img_rgb, (rx, ry), (rx + rw3, ry + rh3), (0, 255, 255), 2)

        # Draw multi-ROI overlays -------------------------------------------
        font = cv2.FONT_HERSHEY_SIMPLEX

        # 1. Grid preview (yellow dashed approximated with thin rectangles)
        for rect in self._grid_preview_rects:
            nx, ny, nw, nh = rect
            rx, ry = int(nx * w), int(ny * h)
            rpw, rph = max(1, int(nw * w)), max(1, int(nh * h))
            cv2.rectangle(img_rgb, (rx, ry), (rx + rpw, ry + rph), (0, 215, 255), 1)

        # 2. Confirmed multi-ROIs
        if self._multi_roi_set:
            for roi in self._multi_roi_set.rois:
                nx, ny, nw, nh = roi.norm_rect
                rx, ry = int(nx * w), int(ny * h)
                rpw, rph = max(1, int(nw * w)), max(1, int(nh * h))
                color = roi.color_bgr
                thickness = 2 if roi.roi_type == 'target' else 1
                cv2.rectangle(img_rgb, (rx, ry), (rx + rpw, ry + rph), color, thickness)
                # Inner shadow for readability
                cv2.rectangle(img_rgb, (rx + 1, ry + 1),
                               (rx + rpw - 1, ry + rph - 1), (0, 0, 0), 1)

        # 3. Multi-Add anchor markers
        _ANCHOR_COLOR = (0, 215, 255)  # gold
        for anchor, label in [(self._grid_anchor_tl, 'TL'), (self._grid_anchor_br, 'BR')]:
            if anchor is not None:
                cx, cy = int(anchor[0] * w), int(anchor[1] * h)
                cv2.drawMarker(img_rgb, (cx, cy), _ANCHOR_COLOR,
                               cv2.MARKER_CROSS, 16, 2, cv2.LINE_AA)
                cv2.putText(img_rgb, label, (cx + 5, cy - 5),
                            font, 0.4, _ANCHOR_COLOR, 1, cv2.LINE_AA)

        # 4. Multi-ROI drag rubber band
        if (self._multi_draw_mode == 'drag' and self._multi_roi_dragging
                and self._multi_roi_start and self._multi_roi_current):
            x0n, y0n = self._multi_roi_start
            x1n, y1n = self._multi_roi_current
            rx = int(min(x0n, x1n) * w)
            ry = int(min(y0n, y1n) * h)
            rw3 = max(1, int(abs(x1n - x0n) * w))
            rh3 = max(1, int(abs(y1n - y0n) * h))
            cv2.rectangle(img_rgb, (rx, ry), (rx + rw3, ry + rh3), (0, 255, 0), 2)

        # 5. Single-add / Multi-add cursor preview (ghost rect following mouse)
        if self._multi_draw_mode in ('single_add', 'multi_add') and self._cursor_pos is not None:
            nw_s, nh_s = self._add_size_norm
            cx_n, cy_n = self._cursor_pos
            nx = max(0.0, min(cx_n - nw_s / 2, 1.0 - nw_s))
            ny = max(0.0, min(cy_n - nh_s / 2, 1.0 - nh_s))
            rx, ry = int(nx * w), int(ny * h)
            rpw, rph = max(1, int(nw_s * w)), max(1, int(nh_s * h))
            # single_add → green;  multi_add → gold (matches anchor marker colour)
            ghost_color = (0, 255, 0) if self._multi_draw_mode == 'single_add' else (0, 215, 255)
            cv2.rectangle(img_rgb, (rx, ry), (rx + rpw, ry + rph), ghost_color, 1)

        # Draw zoom overlay if cursor is on image (suppress during ROI draw modes)
        if self._show_zoom and self._cursor_pos is not None and self._multi_draw_mode == 'idle':
            norm_x, norm_y = self._cursor_pos
            src_x, src_y = int(norm_x * w), int(norm_y * h)

            # No cursor marker - just magnifier
            cross_color = (0, 200, 255)  # Cyan in BGR

            # Extract fixed-size source window with edge padding (avoid boundary stretching)
            src_size = max(8, int(round(self.ZOOM_SIZE / self.ZOOM_FACTOR)))
            half_zoom = src_size // 2
            x1_raw, y1_raw = src_x - half_zoom, src_y - half_zoom
            x2_raw, y2_raw = x1_raw + src_size, y1_raw + src_size
            x1, y1 = max(0, x1_raw), max(0, y1_raw)
            x2, y2 = min(w, x2_raw), min(h, y2_raw)

            if x2 > x1 and y2 > y1:
                zoom_region = img_rgb[y1:y2, x1:x2]
                pad_l = max(0, -x1_raw)
                pad_t = max(0, -y1_raw)
                pad_r = max(0, x2_raw - w)
                pad_b = max(0, y2_raw - h)
                if pad_l or pad_t or pad_r or pad_b:
                    zoom_region = cv2.copyMakeBorder(
                        zoom_region,
                        pad_t,
                        pad_b,
                        pad_l,
                        pad_r,
                        borderType=cv2.BORDER_REPLICATE,
                    )

                # Resize to zoom window size
                zoomed = cv2.resize(zoom_region, (self.ZOOM_SIZE, self.ZOOM_SIZE),
                                    interpolation=cv2.INTER_LINEAR)

                # Place magnifier centered on cursor for direct observation.
                zx = src_x - self.ZOOM_SIZE // 2
                zy = src_y - self.ZOOM_SIZE // 2

                # Keep zoom window inside image bounds
                zx = max(0, min(zx, w - self.ZOOM_SIZE))
                zy = max(0, min(zy, h - self.ZOOM_SIZE))

                if zx >= 0 and zy >= 0:
                    # Apply circular mask to zoomed region
                    mask = np.zeros((self.ZOOM_SIZE, self.ZOOM_SIZE), dtype=np.uint8)
                    center = self.ZOOM_SIZE // 2
                    cv2.circle(mask, (center, center), center - 2, 255, -1)

                    # Get background region
                    bg_region = img_rgb[zy:zy + self.ZOOM_SIZE, zx:zx + self.ZOOM_SIZE].copy()

                    # Blend zoomed with background using mask
                    mask_3ch = cv2.merge([mask, mask, mask])
                    result = np.where(mask_3ch == 255, zoomed, bg_region)
                    img_rgb[zy:zy + self.ZOOM_SIZE, zx:zx + self.ZOOM_SIZE] = result

                    # Draw circular border
                    cv2.circle(img_rgb, (zx + center, zy + center), center - 1, cross_color, 2)

        # Convert BGR to RGB for Qt
        img_qt = cv2.cvtColor(img_rgb, cv2.COLOR_BGR2RGB)
        img_qt = np.ascontiguousarray(img_qt)
        qimg = QtGui.QImage(img_qt.data, w, h, w * 3, QtGui.QImage.Format_RGB888).copy()

        pixmap = QtGui.QPixmap.fromImage(qimg)

        # Scale to fit
        scaled = pixmap.scaled(
            self._label.size() - QtCore.QSize(8, 8),
            Qt.KeepAspectRatio,
            Qt.SmoothTransformation
        )
        self._label.setPixmap(scaled)


class SplitViewWidget(QtWidgets.QWidget):
    """Split-view image comparison widget with draggable vertical divider.

    Left side shows the Base image, right side shows the Compare image.
    The vertical divider can be dragged to reveal more/less of each side.
    The slider_blend value (0-100) also controls the divider position.
    """

    clicked = Signal()

    def __init__(self, parent=None):
        super().__init__(parent)
        self._base_image: Optional[np.ndarray] = None
        self._compare_image: Optional[np.ndarray] = None
        self._divider_ratio: float = 0.5  # 0.0 = all compare, 1.0 = all base
        self._dragging: bool = False
        self._base_pix: Optional[QtGui.QPixmap] = None
        self._comp_pix: Optional[QtGui.QPixmap] = None

        self.setMinimumSize(350, 350)
        self.setMouseTracking(True)
        self.setAttribute(QtCore.Qt.WA_OpaquePaintEvent, True)
        self.setStyleSheet(f"""
            QWidget {{
                background-color: {UI_BG_PANEL};
                border: 1px solid {UI_BORDER};
                border-radius: 12px;
            }}
            QLabel {{
                background-color: {UI_BG_VIEWER};
                border: none;
                border-radius: 8px;
                color: #D1D5DB;
            }}
        """)

    # ── Public API ─────────────────────────────────────────────────────────
    def set_images(self, base: Optional[np.ndarray], compare: Optional[np.ndarray]):
        """Set the two images to compare."""
        self._base_image = base
        self._compare_image = compare
        self._base_pix = self._to_pixmap(base)
        self._comp_pix = self._to_pixmap(compare)
        self.update()

    def set_divider(self, ratio: float):
        """Set divider position: 0.0 = all compare at left, 1.0 = all base at right."""
        self._divider_ratio = max(0.0, min(1.0, ratio))
        self.update()

    # ── Helpers ────────────────────────────────────────────────────────────
    @staticmethod
    def _to_pixmap(img: Optional[np.ndarray]) -> Optional[QtGui.QPixmap]:
        if img is None:
            return None
        if img.dtype != np.uint8:
            # Float images in [0, 1] range (from compute pipeline) need scaling to [0, 255]
            if img.max() <= 1.0 + 1e-6:
                img = (img * 255.0).clip(0, 255).astype(np.uint8)
            else:
                img = np.clip(img, 0, 255).astype(np.uint8)
        if img.ndim == 2:
            img_rgb = cv2.cvtColor(img, cv2.COLOR_GRAY2RGB)
        else:
            img_rgb = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
        img_c = np.ascontiguousarray(img_rgb)
        h, w = img_c.shape[:2]
        qimg = QtGui.QImage(img_c.data, w, h, w * 3,
                            QtGui.QImage.Format_RGB888).copy()
        return QtGui.QPixmap.fromImage(qimg)

    def _content_rect(self) -> QtCore.QRect:
        """Return the inner viewer content rect, matching SyncZoomImageWidget margins."""
        return self.rect().adjusted(8, 8, -8, -8)

    def _fitted_target_rect(self) -> QtCore.QRect:
        """Return centered fit-rect using the same KeepAspectRatio logic as standard viewer."""
        content = self._content_rect()
        pix = self._base_pix if self._base_pix is not None else self._comp_pix
        if pix is None or pix.isNull() or content.width() <= 0 or content.height() <= 0:
            return content

        scaled = pix.size().scaled(content.size(), Qt.KeepAspectRatio)
        x = content.x() + (content.width() - scaled.width()) // 2
        y = content.y() + (content.height() - scaled.height()) // 2
        return QtCore.QRect(x, y, scaled.width(), scaled.height())

    # ── Painting ───────────────────────────────────────────────────────────
    def paintEvent(self, event):
        p = QtGui.QPainter(self)
        p.setRenderHint(QtGui.QPainter.SmoothPixmapTransform)

        target_rect = self._fitted_target_rect()
        x0, y0, iw, ih = target_rect.x(), target_rect.y(), target_rect.width(), target_rect.height()
        div_x = x0 + int(self._divider_ratio * iw)

        # Background fill keeps viewer matte black (consistent with standard views)
        p.fillRect(self.rect(), QtGui.QColor(UI_BG_VIEWER))

        if self._base_pix is None and self._comp_pix is None:
            p.setPen(QtGui.QColor(UI_TEXT_MUTED))
            p.drawText(self.rect(), Qt.AlignCenter, "No images loaded")
            p.end()
            return

        # Left half → Base
        left_w = max(0, div_x - x0)
        if self._base_pix and left_w > 0:
            p.save()
            p.setClipRect(x0, y0, left_w, ih)
            p.drawPixmap(target_rect, self._base_pix)
            p.restore()

        # Right half → Compare
        right_w = max(0, iw - left_w)
        if self._comp_pix and right_w > 0:
            p.save()
            p.setClipRect(div_x, y0, right_w, ih)
            p.drawPixmap(target_rect, self._comp_pix)
            p.restore()

        # Divider line
        p.setPen(QtGui.QPen(QtGui.QColor(UI_PRIMARY), 2, Qt.SolidLine))
        p.drawLine(div_x, y0, div_x, y0 + ih)

        # Diamond handle at mid-height
        mid_y = y0 + ih // 2
        s = 9
        diamond = QtGui.QPolygon([
            QtCore.QPoint(div_x, mid_y - s),
            QtCore.QPoint(div_x + s, mid_y),
            QtCore.QPoint(div_x, mid_y + s),
            QtCore.QPoint(div_x - s, mid_y),
        ])
        p.setBrush(QtGui.QBrush(QtGui.QColor(UI_PRIMARY)))
        p.setPen(Qt.NoPen)
        p.drawPolygon(diamond)

        # Side labels
        font = p.font()
        font.setPointSize(8)
        font.setBold(True)
        p.setFont(font)
        p.setPen(QtGui.QColor(UI_TEXT))
        if left_w > 40:
            p.drawText(x0 + 6, y0 + 18, "Base")
        if right_w > 60:
            p.drawText(div_x + 6, y0 + 18, "Compare")

        p.end()

    # ── Mouse interaction ──────────────────────────────────────────────────
    def _near_divider(self, pos_x: int) -> bool:
        target_rect = self._fitted_target_rect()
        x0, iw = target_rect.x(), target_rect.width()
        div_x = x0 + int(self._divider_ratio * iw)
        return abs(pos_x - div_x) <= 14

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self.clicked.emit()
        if event.button() == Qt.LeftButton and self._near_divider(event.pos().x()):
            self._dragging = True
            self.setCursor(Qt.SplitHCursor)
            event.accept()
        else:
            super().mousePressEvent(event)

    def mouseMoveEvent(self, event):
        target_rect = self._fitted_target_rect()
        x0, iw = target_rect.x(), target_rect.width()
        if self._dragging:
            ratio = (event.pos().x() - x0) / max(iw, 1)
            self._divider_ratio = max(0.0, min(1.0, ratio))
            self.update()
            event.accept()
        else:
            self.setCursor(Qt.SplitHCursor if self._near_divider(event.pos().x())
                           else Qt.ArrowCursor)
            super().mouseMoveEvent(event)

    def mouseReleaseEvent(self, event):
        if self._dragging and event.button() == Qt.LeftButton:
            self._dragging = False
            self.setCursor(Qt.ArrowCursor)
            event.accept()
        else:
            super().mouseReleaseEvent(event)

    def leaveEvent(self, event):
        self._dragging = False
        self.setCursor(Qt.ArrowCursor)
        super().leaveEvent(event)

    def resizeEvent(self, event):
        self.update()
        super().resizeEvent(event)


class ImageDisplayWidget(QtWidgets.QLabel):
    """Widget for displaying images with modern styling."""

    def __init__(self, title: str = "", parent=None):
        super().__init__(parent)
        self._title = title
        self._image: Optional[np.ndarray] = None
        self._hint_rect: Optional[tuple] = None  # (x, y, w, h) for hint box
        self._hint_info: str = ""  # Info text for hint
        self._show_hint: bool = False
        self.setMinimumSize(280, 280)  # Square display
        self.setAlignment(Qt.AlignCenter)
        self.setStyleSheet(f"""
            QLabel {{
                background-color: {UI_BG_CARD};
                border: 1px solid {UI_BORDER};
                border-radius: 12px;
            }}
        """)
        self._update_display()

    def setImage(self, image: np.ndarray):
        """Set image to display."""
        self._image = image
        self._update_display()

    def setHint(self, rect: tuple = None, info: str = "", show: bool = True):
        """Set hint rectangle and info.

        Args:
            rect: (x, y, w, h) in image coordinates
            info: Text to display near the hint
            show: Whether to show the hint
        """
        self._hint_rect = rect
        self._hint_info = info
        self._show_hint = show
        self._update_display()

    def setShowHint(self, show: bool):
        """Toggle hint visibility."""
        self._show_hint = show
        self._update_display()

    def _update_display(self):
        """Update the displayed pixmap."""
        if self._image is None:
            self.setText(f"{self._title}\n\nNo image")
            return

        h, w = self._image.shape[:2]

        # Convert to QImage
        if len(self._image.shape) == 2:
            # Grayscale -> Convert to RGB for drawing
            img_display = self._image
            if img_display.dtype != np.uint8:
                img_display = np.clip(img_display, 0, 255).astype(np.uint8)
            img_rgb = cv2.cvtColor(img_display, cv2.COLOR_GRAY2BGR)
        else:
            # Already color
            img_rgb = self._image.copy()

        # Draw hint box if enabled
        if self._show_hint and self._hint_rect is not None:
            x, y, rw, rh = self._hint_rect
            # Draw rectangle (teal accent color - BGR format)
            hint_color = _hex_to_bgr(UI_SUCCESS)
            cv2.rectangle(img_rgb, (x, y), (x + rw, y + rh), hint_color, 2)
            # Draw corner marks instead of full crosshair
            cx, cy = x + rw // 2, y + rh // 2
            mark_len = 10
            cv2.line(img_rgb, (cx - mark_len, cy), (cx + mark_len, cy), hint_color, 1)
            cv2.line(img_rgb, (cx, cy - mark_len), (cx, cy + mark_len), hint_color, 1)
            # Draw info text with shadow for readability
            if self._hint_info:
                font = cv2.FONT_HERSHEY_SIMPLEX
                text_y = max(y - 8, 18)
                # Shadow
                cv2.putText(img_rgb, self._hint_info, (x + 1, text_y + 1),
                            font, 0.45, (0, 0, 0), 2, cv2.LINE_AA)
                # Main text
                cv2.putText(img_rgb, self._hint_info, (x, text_y),
                            font, 0.45, hint_color, 1, cv2.LINE_AA)

        # Convert BGR to RGB for Qt
        img_qt = cv2.cvtColor(img_rgb, cv2.COLOR_BGR2RGB)
        # Make a contiguous copy to ensure data lifetime
        img_qt = np.ascontiguousarray(img_qt)
        qimg = QtGui.QImage(img_qt.data, w, h, w * 3, QtGui.QImage.Format_RGB888).copy()

        pixmap = QtGui.QPixmap.fromImage(qimg)

        # Scale to fit
        scaled = pixmap.scaled(
            self.size() - QtCore.QSize(20, 20),
            Qt.KeepAspectRatio,
            Qt.SmoothTransformation
        )
        self.setPixmap(scaled)

    def resizeEvent(self, event):
        super().resizeEvent(event)
        self._update_display()


class HistogramCanvas(FigureCanvas):
    """Interactive matplotlib canvas for histogram display.

    Click once → set low bound, click again → set high bound and emit range_changed.
    Right-click or clicking a third time resets the selection.
    """

    range_changed = Signal(int, int)  # (lo, hi) gray-level values

    def __init__(self, parent=None):
        self.fig = Figure(figsize=(5, 2.5), dpi=100)
        self.ax = self.fig.add_subplot(111)
        super().__init__(self.fig)
        self.setParent(parent)
        self._style()
        # Range selection state
        self._counts = None
        self._edges = None
        self._lo: Optional[int] = None
        self._hi: Optional[int] = None
        self._click_step: int = 0  # 0=idle, 1=waiting for hi
        self.mpl_connect('button_press_event', self._on_click)

    def _style(self):
        """Apply modern dark theme."""
        self.fig.patch.set_facecolor(UI_BG_WINDOW)
        self.ax.set_facecolor(UI_BG_CARD)
        self.ax.tick_params(colors=UI_TEXT_SECONDARY, labelsize=8, length=0)
        for spine in self.ax.spines.values():
            spine.set_visible(False)
        self.ax.yaxis.grid(True, color=UI_BORDER, linestyle='-', linewidth=0.5, alpha=0.5)
        self.ax.set_axisbelow(True)

    def clear_range(self):
        """Reset range selection without re-plotting data."""
        self._lo = None
        self._hi = None
        self._click_step = 0
        if self._counts is not None:
            self._draw(self._counts, self._edges)

    def plot_histogram(self, counts: np.ndarray, edges: np.ndarray):
        """Plot histogram data, preserving any active range selection."""
        self._counts = counts
        self._edges = edges
        self._draw(counts, edges)

    def _draw(self, counts, edges):
        """Internal render with range overlay."""
        self.ax.clear()
        self._style()

        if counts is None or len(counts) == 0:
            self.ax.text(0.5, 0.5, "No data", ha='center', va='center',
                         color=UI_TEXT_SECONDARY, transform=self.ax.transAxes, fontsize=11)
            self.draw()
            return

        # Base bars
        lo_set = self._lo is not None
        hi_set = self._hi is not None and hi_set if False else (self._hi is not None)
        lo_v = self._lo if lo_set else -1
        hi_v = self._hi if hi_set else 256

        bar_colors = []
        for edge in edges[:-1]:
            if lo_set and hi_set and lo_v <= edge <= hi_v:
                bar_colors.append(UI_SUCCESS)  # success highlight
            elif lo_set and not hi_set and edge >= lo_v:
                bar_colors.append(UI_SUCCESS)  # pending second click
            else:
                bar_colors.append(UI_PRIMARY)

        self.ax.bar(edges[:-1], counts, width=1, color=bar_colors, alpha=0.85,
                    edgecolor='none', linewidth=0)

        # Range shading
        if lo_set and hi_set:
            self.ax.axvspan(lo_v, hi_v, alpha=0.12, color=UI_SUCCESS, zorder=0)

        # Vertical markers
        if lo_set:
            self.ax.axvline(lo_v, color=UI_SUCCESS, linewidth=1.5, linestyle='--')
        if hi_set:
            self.ax.axvline(hi_v, color=UI_SUCCESS, linewidth=1.5, linestyle='--')

        # Annotation
        if lo_set and not hi_set:
            self.ax.text(0.02, 0.96, f"lo={lo_v}  ← click to set hi",
                         transform=self.ax.transAxes, fontsize=7.5,
                         color=UI_SUCCESS, va='top')
        elif lo_set and hi_set:
            pct = np.sum((edges[:-1] >= lo_v) & (edges[:-1] <= hi_v) * counts)
            total = counts.sum() or 1
            self.ax.text(0.02, 0.96, f"GL {lo_v}–{hi_v}  ({pct / total * 100:.1f}%)",
                         transform=self.ax.transAxes, fontsize=7.5,
                         color=UI_SUCCESS, va='top')

        # Y-axis formatter
        from matplotlib.ticker import FuncFormatter
        self.ax.yaxis.set_major_formatter(FuncFormatter(
            lambda x, _: f'{x / 1000:.0f}K' if x >= 1000 else f'{x:.0f}'
        ))
        self.ax.set_xlabel("Gray Level", color=UI_TEXT_SECONDARY, fontsize=9)
        self.ax.set_ylabel("Count", color=UI_TEXT_SECONDARY, fontsize=9)
        self.ax.set_xlim(0, 255)

        try:
            self.fig.tight_layout(pad=1.2)
        except Exception:
            pass
        self.draw()

    def _on_click(self, event):
        """Handle mouse click: set lo then hi, right-click resets."""
        if event.inaxes != self.ax or event.xdata is None:
            return
        x = int(round(max(0, min(255, event.xdata))))

        if event.button == 3:  # Right-click → reset
            self._lo = None
            self._hi = None
            self._click_step = 0
            self._draw(self._counts, self._edges)
            return

        if self._click_step == 0:
            # First click: set lo
            self._lo = x
            self._hi = None
            self._click_step = 1
        elif self._click_step == 1:
            # Second click: set hi, emit signal
            lo, hi = sorted([self._lo, x])
            self._lo, self._hi = lo, hi
            self._click_step = 0
            self._draw(self._counts, self._edges)
            self.range_changed.emit(lo, hi)
            return

        self._draw(self._counts, self._edges)


class _Worker(QtCore.QObject):
    finished = Signal(object)
    error = Signal(str)
    progress = Signal(int, str)   # (completed_count, current_label)

    def __init__(self, fn):
        super().__init__()
        self._fn = fn
        self.abort_requested: bool = False

    def run(self):
        try:
            result = self._fn(self)
        except Exception as exc:
            self.error.emit(str(exc))
            return
        self.finished.emit(result)


class AlignmentScoreWidget(QtWidgets.QFrame):
    """Alignment summary card styled to match the design mockup."""

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setObjectName("AlignCard")
        self.setMinimumWidth(248)
        self.setSizePolicy(QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Expanding)
        self.setStyleSheet(f"""
            QFrame#AlignCard {{
                background-color: {UI_BG_CARD};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.MD};
            }}
            QFrame#AlignCard QLabel {{
                border: none;
                background: transparent;
            }}
            QFrame#AlignCard QLabel[cardTitle="true"] {{
                color: {UI_TEXT};
                font-size: {Typography.FONT_SIZE_BODY};
                font-weight: {Typography.FONT_WEIGHT_BOLD};
            }}
            QFrame#AlignCard QLabel[cardSubTitle="true"] {{
                color: {UI_TEXT_MUTED};
                font-size: {Typography.FONT_SIZE_CAPTION};
            }}
            QFrame#AlignCard QLabel[statLabel="true"] {{
                color: {UI_TEXT_MUTED};
                font-size: {Typography.FONT_SIZE_SMALL};
            }}
            QFrame#AlignCard QLabel[statValue="true"] {{
                color: {UI_TEXT};
                font-size: {Typography.FONT_SIZE_BODY};
                font-weight: {Typography.FONT_WEIGHT_BOLD};
                font-family: {Typography.FONT_FAMILY_MONO};
            }}
            QFrame#AlignCard QLabel[statValueLong="true"] {{
                color: {UI_TEXT_SECONDARY};
                font-size: {Typography.FONT_SIZE_SMALL};
                font-weight: {Typography.FONT_WEIGHT_SEMIBOLD};
                font-family: {Typography.FONT_FAMILY_MONO};
            }}
            QFrame#AlignCard QFrame[divider="true"] {{
                background-color: {UI_BORDER};
                min-height: 1px;
                max-height: 1px;
                border: none;
            }}
        """)

        layout = QtWidgets.QVBoxLayout(self)
        layout.setContentsMargins(16, 14, 16, 14)
        layout.setSpacing(8)

        header = QtWidgets.QVBoxLayout()
        header.setSpacing(4)
        top_row = QtWidgets.QHBoxLayout()
        top_row.setSpacing(8)
        dot = QtWidgets.QLabel()
        dot.setFixedSize(8, 8)
        dot.setStyleSheet(f"background-color: {UI_PRIMARY}; border-radius: 4px;")
        title = QtWidgets.QLabel("Alignment Quality")
        title.setProperty("cardTitle", True)
        top_row.addWidget(dot, 0, Qt.AlignVCenter)
        top_row.addWidget(title)
        top_row.addStretch()
        header.addLayout(top_row)
        sub = QtWidgets.QLabel("Single-pair alignment confidence")
        sub.setProperty("cardSubTitle", True)
        header.addWidget(sub)
        layout.addLayout(header)

        divider = QtWidgets.QFrame()
        divider.setProperty("divider", True)
        layout.addWidget(divider)

        self.rows: Dict[str, QtWidgets.QLabel] = {}
        for key, label in [
            ("phase", "Phase Shift"),
            ("ncc", "NCC"),
            ("residual", "Residual"),
            ("final", "Final Score"),
            ("shift", "Shift"),
        ]:
            row = QtWidgets.QHBoxLayout()
            row.setSpacing(8)
            lbl_name = QtWidgets.QLabel(label)
            lbl_name.setProperty("statLabel", True)
            lbl_value = QtWidgets.QLabel("--")
            lbl_value.setProperty("statValue", True)
            lbl_value.setAlignment(Qt.AlignRight | Qt.AlignVCenter)
            if key == "shift":
                lbl_value.setProperty("statValueLong", True)
            row.addWidget(lbl_name)
            row.addStretch()
            row.addWidget(lbl_value)
            layout.addLayout(row)
            self.rows[key] = lbl_value

        divider2 = QtWidgets.QFrame()
        divider2.setProperty("divider", True)
        layout.addWidget(divider2)

        # Backward-compatible aliases used by the existing update logic.
        self.lbl_phase = self.rows["phase"]
        self.lbl_ncc = self.rows["ncc"]
        self.lbl_residual = self.rows["residual"]
        self.lbl_final = self.rows["final"]
        self.lbl_shift = self.rows["shift"]

        status_row = QtWidgets.QHBoxLayout()
        self.lbl_status_text = QtWidgets.QLabel("Status")
        self.lbl_status_text.setProperty("statLabel", True)
        self.lbl_status = QtWidgets.QLabel("--")
        self.lbl_status.setProperty("statValue", True)
        self.lbl_status.setAlignment(Qt.AlignRight | Qt.AlignVCenter)
        status_row.addWidget(self.lbl_status_text)
        status_row.addStretch()
        status_row.addWidget(self.lbl_status)
        layout.addLayout(status_row)
        layout.addStretch(1)

    def update_scores(self, result: CombineResult):
        if not result or not result.alignments:
            self._set_empty()
            return

        n = len(result.alignments)
        avg_phase = sum(a.score_phase for a in result.alignments) / n
        avg_ncc = sum(a.score_ncc for a in result.alignments) / n
        avg_residual = sum(a.score_residual for a in result.alignments) / n
        avg_final = sum(a.final_score for a in result.alignments) / n
        avg_dx = sum(a.dx for a in result.alignments) / n
        avg_dy = sum(a.dy for a in result.alignments) / n

        self.rows["phase"].setText(f"{avg_phase:.3f}")
        self.rows["ncc"].setText(f"{avg_ncc:.3f}")
        self.rows["residual"].setText(f"{avg_residual:.3f}")
        self.rows["final"].setText(f"{avg_final:.1f}")
        self.rows["shift"].setText(f"({avg_dx:+.1f}, {avg_dy:+.1f})")

        worst = result.worst_alignment_score
        if worst >= 75:
            status_text = "OK"
            status_color = UI_SUCCESS
        elif worst >= 55:
            status_text = "WARN"
            status_color = UI_PRIMARY_HOVER
        else:
            status_text = "FAIL"
            status_color = UI_WARNING
        self.lbl_status.setText(status_text)
        self.lbl_status.setStyleSheet(
            f"color: {status_color}; font-size: {Typography.FONT_SIZE_BODY}; font-weight: {Typography.FONT_WEIGHT_BOLD}; border: none;"
        )

    def _set_empty(self):
        for key, lbl in self.rows.items():
            lbl.setText("--")
            prop = "true" if key == "shift" else None
            if prop:
                lbl.setProperty("statValueLong", True)
                lbl.style().unpolish(lbl)
                lbl.style().polish(lbl)
        self.lbl_status.setText("--")
        self.lbl_status.setStyleSheet(
            f"color: {UI_TEXT}; font-size: {Typography.FONT_SIZE_BODY}; font-weight: {Typography.FONT_WEIGHT_BOLD}; border: none;"
        )


class StatisticsWidget(QtWidgets.QFrame):
    """Merged Analysis card — shows Alignment and Difference metrics in one panel."""

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setObjectName("StatsCard")
        self.setMinimumWidth(248)
        self.setSizePolicy(QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Expanding)
        self.setStyleSheet(f"""
            QFrame#StatsCard {{
                background-color: {UI_BG_CARD};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.MD};
            }}
            QFrame#StatsCard QLabel {{
                border: none;
                background: transparent;
            }}
            QFrame#StatsCard QLabel[cardTitle="true"] {{
                color: {UI_TEXT};
                font-size: {Typography.FONT_SIZE_BODY};
                font-weight: {Typography.FONT_WEIGHT_BOLD};
            }}
            QFrame#StatsCard QLabel[sectionHeader="true"] {{
                color: {UI_TEXT_SECONDARY};
                font-size: {Typography.FONT_SIZE_CAPTION};
                font-weight: {Typography.FONT_WEIGHT_BOLD};
                text-transform: uppercase;
                letter-spacing: 0.5px;
            }}
            QFrame#StatsCard QLabel[statLabel="true"] {{
                color: {UI_TEXT_MUTED};
                font-size: {Typography.FONT_SIZE_SMALL};
            }}
            QFrame#StatsCard QLabel[statValue="true"] {{
                color: {UI_TEXT};
                font-size: {Typography.FONT_SIZE_BODY};
                font-weight: {Typography.FONT_WEIGHT_BOLD};
                font-family: {Typography.FONT_FAMILY_MONO};
            }}
            QFrame#StatsCard QLabel[statValueLong="true"] {{
                color: {UI_TEXT_SECONDARY};
                font-size: {Typography.FONT_SIZE_SMALL};
                font-weight: {Typography.FONT_WEIGHT_SEMIBOLD};
                font-family: {Typography.FONT_FAMILY_MONO};
            }}
            QFrame#StatsCard QFrame[divider="true"] {{
                background-color: {UI_BORDER};
                min-height: 1px;
                max-height: 1px;
                border: none;
            }}
        """)

        layout = QtWidgets.QVBoxLayout(self)
        layout.setContentsMargins(16, 14, 16, 14)
        layout.setSpacing(6)

        # Card title
        top_row = QtWidgets.QHBoxLayout()
        top_row.setSpacing(8)
        dot = QtWidgets.QLabel()
        dot.setFixedSize(8, 8)
        dot.setStyleSheet(f"background-color: {UI_PRIMARY}; border-radius: 4px;")
        title = QtWidgets.QLabel("Analysis")
        title.setProperty("cardTitle", True)
        top_row.addWidget(dot, 0, Qt.AlignVCenter)
        top_row.addWidget(title)
        top_row.addStretch()
        layout.addLayout(top_row)

        divider0 = QtWidgets.QFrame()
        divider0.setProperty("divider", True)
        layout.addWidget(divider0)

        # ── ALIGNMENT section ─────────────────────────────────────────────────
        lbl_align_hdr = QtWidgets.QLabel("Alignment")
        lbl_align_hdr.setProperty("sectionHeader", True)
        layout.addWidget(lbl_align_hdr)

        self.align_labels: Dict[str, QtWidgets.QLabel] = {}
        for key, label in [
            ("phase", "Phase Shift"),
            ("ncc", "NCC"),
            ("residual", "Residual"),
            ("final", "Final Score"),
            ("shift", "Shift"),
        ]:
            row = QtWidgets.QHBoxLayout()
            row.setSpacing(8)
            lbl_name = QtWidgets.QLabel(label)
            lbl_name.setProperty("statLabel", True)
            lbl_value = QtWidgets.QLabel("--")
            lbl_value.setAlignment(Qt.AlignRight | Qt.AlignVCenter)
            if key == "shift":
                lbl_value.setProperty("statValueLong", True)
            else:
                lbl_value.setProperty("statValue", True)
            row.addWidget(lbl_name)
            row.addStretch()
            row.addWidget(lbl_value)
            layout.addLayout(row)
            self.align_labels[key] = lbl_value

        # Status row (align)
        status_row = QtWidgets.QHBoxLayout()
        lbl_status_name = QtWidgets.QLabel("Status")
        lbl_status_name.setProperty("statLabel", True)
        self.lbl_align_status = QtWidgets.QLabel("--")
        self.lbl_align_status.setProperty("statValue", True)
        self.lbl_align_status.setAlignment(Qt.AlignRight | Qt.AlignVCenter)
        status_row.addWidget(lbl_status_name)
        status_row.addStretch()
        status_row.addWidget(self.lbl_align_status)
        layout.addLayout(status_row)

        divider1 = QtWidgets.QFrame()
        divider1.setProperty("divider", True)
        layout.addWidget(divider1)

        # ── DIFFERENCE section ────────────────────────────────────────────────
        lbl_diff_hdr = QtWidgets.QLabel("Difference")
        lbl_diff_hdr.setProperty("sectionHeader", True)
        layout.addWidget(lbl_diff_hdr)

        self.stats_labels: Dict[str, QtWidgets.QLabel] = {}
        for key, label in [
            ("diff_mean", "Diff Mean"),
            ("diff_std", "Diff Std"),
            ("hot_pixels", "Hot Pixels"),
            ("norm_coeff", "Normalize"),
            ("subpixel_shift", "Sub-px Shift"),
        ]:
            row = QtWidgets.QHBoxLayout()
            row.setSpacing(8)
            lbl_name = QtWidgets.QLabel(label)
            lbl_name.setProperty("statLabel", True)
            lbl_value = QtWidgets.QLabel("--")
            lbl_value.setAlignment(Qt.AlignRight | Qt.AlignVCenter)
            if key in {"norm_coeff", "subpixel_shift"}:
                lbl_value.setProperty("statValueLong", True)
            else:
                lbl_value.setProperty("statValue", True)
            row.addWidget(lbl_name)
            row.addStretch()
            row.addWidget(lbl_value)
            layout.addLayout(row)
            self.stats_labels[key] = lbl_value

        layout.addStretch(1)
        self.reset()

    def update_stats(self, result: CombineResult):
        if not result or not result.stats:
            for lbl in self.stats_labels.values():
                lbl.setText("--")
            return

        s = result.stats
        self.stats_labels["diff_mean"].setText(f"{s.get('diff_mean', 0):.4f}")
        self.stats_labels["diff_std"].setText(f"{s.get('diff_std', 0):.4f}")
        self.stats_labels["hot_pixels"].setText(f"{s.get('hot_pixels', 0)}")
        self.stats_labels["norm_coeff"].setText(f"a={s.get('norm_a', 0):.4f}, b={s.get('norm_b', 0):.2f}")
        self.stats_labels["subpixel_shift"].setText(
            f"({s.get('subpixel_dx', 0):+.2f}, {s.get('subpixel_dy', 0):+.2f})"
        )

    def update_alignment(self, phase: str, ncc: str, residual: str,
                         final: str, shift: str, status: str, status_color: str):
        """Update the alignment section of the merged analysis card."""
        self.align_labels["phase"].setText(phase)
        self.align_labels["ncc"].setText(ncc)
        self.align_labels["residual"].setText(residual)
        self.align_labels["final"].setText(final)
        self.align_labels["shift"].setText(shift)
        self.lbl_align_status.setText(f" {status} ")
        self.lbl_align_status.setStyleSheet(
            f"color: #FFFFFF; background-color: {status_color}; border: 1px solid {status_color};"
            f" border-radius: {BorderRadius.SM}; padding: 1px 8px;"
            f" font-size: {Typography.FONT_SIZE_SMALL}; font-weight: {Typography.FONT_WEIGHT_BOLD};"
        )

    def reset(self):
        """Reset Analysis card to default empty-state values."""
        for lbl in self.align_labels.values():
            lbl.setText("--")
        for lbl in self.stats_labels.values():
            lbl.setText("--")
        self.lbl_align_status.setText(" -- ")
        self.lbl_align_status.setStyleSheet(
            f"color: {UI_TEXT_MUTED}; background-color: {UI_BG_SUBTLE}; border: 1px solid {UI_BORDER};"
            f" border-radius: {BorderRadius.SM}; padding: 1px 8px;"
            f" font-size: {Typography.FONT_SIZE_SMALL}; font-weight: {Typography.FONT_WEIGHT_BOLD};"
        )


class AlignmentPanelWidget(QtWidgets.QFrame):
    """Compact alignment-only card for STATE B bottom panel (LEFT slot).

    Shows only alignment-quality metrics using a QFormLayout.
    Difference / statistics fields have been moved to DiffROIAnalysisPanelWidget.
    """

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setObjectName("AlignCard")
        self.setSizePolicy(QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Expanding)
        self.setStyleSheet(f"""
            QFrame#AlignCard {{
                background-color: {UI_BG_CARD};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.MD};
            }}
            QFrame#AlignCard QLabel {{
                border: none;
                background: transparent;
            }}
        """)

        layout = QtWidgets.QVBoxLayout(self)
        layout.setContentsMargins(14, 12, 14, 12)
        layout.setSpacing(6)

        # Card title row
        top_row = QtWidgets.QHBoxLayout()
        top_row.setSpacing(8)
        dot = QtWidgets.QLabel()
        dot.setFixedSize(8, 8)
        dot.setStyleSheet(f"background-color: {UI_PRIMARY}; border-radius: 4px;")
        title = QtWidgets.QLabel("Alignment")
        title.setStyleSheet(
            f"color: {UI_TEXT}; font-size: {Typography.FONT_SIZE_BODY};"
            f" font-weight: {Typography.FONT_WEIGHT_BOLD};"
        )
        top_row.addWidget(dot, 0, Qt.AlignVCenter)
        top_row.addWidget(title)
        top_row.addStretch()
        layout.addLayout(top_row)

        divider = QtWidgets.QFrame()
        divider.setStyleSheet(f"background-color: {UI_BORDER}; min-height:1px; max-height:1px; border:none;")
        layout.addWidget(divider)

        # Form layout for compact two-column display
        form = QtWidgets.QFormLayout()
        form.setSpacing(5)
        form.setContentsMargins(0, 4, 0, 0)
        form.setLabelAlignment(Qt.AlignLeft | Qt.AlignVCenter)
        form.setFormAlignment(Qt.AlignLeft | Qt.AlignTop)

        _label_style = (
            f"color: {UI_TEXT_MUTED}; font-size: {Typography.FONT_SIZE_SMALL};"
        )
        _value_style = (
            f"color: {UI_TEXT}; font-size: {Typography.FONT_SIZE_BODY};"
            f" font-weight: {Typography.FONT_WEIGHT_BOLD};"
            f" font-family: {Typography.FONT_FAMILY_MONO};"
        )

        self._fields: Dict[str, QtWidgets.QLabel] = {}
        for key, label_text in [
            ("phase", "Phase Shift"),
            ("ncc", "NCC"),
            ("residual", "Residual"),
            ("final", "Final Score"),
            ("shift", "Shift (dx, dy)"),
        ]:
            lbl_name = QtWidgets.QLabel(label_text)
            lbl_name.setStyleSheet(_label_style)
            lbl_val = QtWidgets.QLabel("--")
            lbl_val.setStyleSheet(_value_style)
            form.addRow(lbl_name, lbl_val)
            self._fields[key] = lbl_val

        layout.addLayout(form)

        # Status row with colored badge
        status_row = QtWidgets.QHBoxLayout()
        lbl_status_name = QtWidgets.QLabel("Status")
        lbl_status_name.setStyleSheet(_label_style)
        self._lbl_status = QtWidgets.QLabel(" -- ")
        status_row.addWidget(lbl_status_name)
        status_row.addStretch()
        status_row.addWidget(self._lbl_status)
        layout.addLayout(status_row)

        layout.addStretch(1)
        self.reset()

    def update_alignment(self, phase: str, ncc: str, residual: str,
                         final: str, shift: str, status: str, status_color: str):
        self._fields["phase"].setText(phase)
        self._fields["ncc"].setText(ncc)
        self._fields["residual"].setText(residual)
        self._fields["final"].setText(final)
        self._fields["shift"].setText(shift)
        self._lbl_status.setText(f" {status} ")
        self._lbl_status.setStyleSheet(
            f"color: #FFFFFF; background-color: {status_color}; border: 1px solid {status_color};"
            f" border-radius: {BorderRadius.SM}; padding: 1px 8px;"
            f" font-size: {Typography.FONT_SIZE_SMALL}; font-weight: {Typography.FONT_WEIGHT_BOLD};"
        )

    def reset(self):
        for lbl in self._fields.values():
            lbl.setText("--")
        self._lbl_status.setText(" -- ")
        self._lbl_status.setStyleSheet(
            f"color: {UI_TEXT_MUTED}; background-color: {UI_BG_SUBTLE};"
            f" border: 1px solid {UI_BORDER}; border-radius: {BorderRadius.SM};"
            f" padding: 1px 8px; font-size: {Typography.FONT_SIZE_SMALL};"
            f" font-weight: {Typography.FONT_WEIGHT_BOLD};"
        )


class DiffROIAnalysisPanelWidget(QtWidgets.QFrame):
    """Center panel for STATE B bottom row — Diff / ROI Analysis.

    Two display states:
    - ROI defined  : shows compact ROI quantification summary for the current result.
    - No ROI       : shows a fallback message with button to open ROI Manager.

    The [ROI Details...] button opens the full ROIIntensityProfileDialog.
    The [Open ROI Manager] button opens the ROI Manager configuration dialog.
    """

    # Signals emitted by the embedded action buttons
    open_roi_manager_requested = Signal()
    roi_details_requested = Signal()

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setObjectName("DiffROICard")
        self.setSizePolicy(QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Expanding)
        self.setStyleSheet(f"""
            QFrame#DiffROICard {{
                background-color: {UI_BG_CARD};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.MD};
            }}
            QFrame#DiffROICard QLabel {{
                border: none;
                background: transparent;
            }}
        """)

        outer = QtWidgets.QVBoxLayout(self)
        outer.setContentsMargins(14, 12, 14, 12)
        outer.setSpacing(6)

        # Card title row
        top_row = QtWidgets.QHBoxLayout()
        top_row.setSpacing(8)
        dot = QtWidgets.QLabel()
        dot.setFixedSize(8, 8)
        dot.setStyleSheet(f"background-color: {UI_INFO}; border-radius: 4px;")
        title = QtWidgets.QLabel("Diff / ROI Analysis")
        title.setStyleSheet(
            f"color: {UI_TEXT}; font-size: {Typography.FONT_SIZE_BODY};"
            f" font-weight: {Typography.FONT_WEIGHT_BOLD};"
        )
        top_row.addWidget(dot, 0, Qt.AlignVCenter)
        top_row.addWidget(title)
        top_row.addStretch()
        outer.addLayout(top_row)

        divider = QtWidgets.QFrame()
        divider.setStyleSheet(f"background-color: {UI_BORDER}; min-height:1px; max-height:1px; border:none;")
        outer.addWidget(divider)

        # Stacked widget: page 0 = ROI summary, page 1 = no-ROI fallback
        self._stack = QtWidgets.QStackedWidget()
        outer.addWidget(self._stack, 1)

        # ── Page 0: ROI summary ────────────────────────────────────────────
        page_roi = QtWidgets.QWidget()
        roi_layout = QtWidgets.QVBoxLayout(page_roi)
        roi_layout.setContentsMargins(0, 2, 0, 0)
        roi_layout.setSpacing(4)

        _sec_style = (
            f"color: {UI_TEXT_SECONDARY}; font-size: {Typography.FONT_SIZE_CAPTION};"
            f" font-weight: {Typography.FONT_WEIGHT_BOLD}; text-transform: uppercase;"
        )
        _label_style = (
            f"color: {UI_TEXT_MUTED}; font-size: {Typography.FONT_SIZE_SMALL};"
        )
        _value_style = (
            f"color: {UI_TEXT}; font-size: {Typography.FONT_SIZE_BODY};"
            f" font-weight: {Typography.FONT_WEIGHT_BOLD};"
            f" font-family: {Typography.FONT_FAMILY_MONO};"
        )
        _value_accent = (
            f"color: {UI_PRIMARY}; font-size: {Typography.FONT_SIZE_BODY};"
            f" font-weight: {Typography.FONT_WEIGHT_BOLD};"
            f" font-family: {Typography.FONT_FAMILY_MONO};"
        )

        def _make_form_row(label_text: str, value_style: str = _value_style):
            row = QtWidgets.QHBoxLayout()
            row.setSpacing(6)
            lbl = QtWidgets.QLabel(label_text)
            lbl.setStyleSheet(_label_style)
            val = QtWidgets.QLabel("--")
            val.setStyleSheet(value_style)
            val.setAlignment(Qt.AlignRight | Qt.AlignVCenter)
            row.addWidget(lbl)
            row.addStretch()
            row.addWidget(val)
            return row, val

        # Section 1: ROI Summary
        hdr1 = QtWidgets.QLabel("ROI Summary")
        hdr1.setStyleSheet(_sec_style)
        roi_layout.addWidget(hdr1)

        row, self._lbl_target_count = _make_form_row("Target ROI")
        roi_layout.addLayout(row)
        row, self._lbl_ref_count = _make_form_row("Reference ROIs")
        roi_layout.addLayout(row)
        row, self._lbl_roi_mode = _make_form_row("Calibration Mode")
        roi_layout.addLayout(row)

        div1 = QtWidgets.QFrame()
        div1.setStyleSheet(f"background-color: {UI_BORDER}; min-height:1px; max-height:1px; border:none;")
        roi_layout.addWidget(div1)

        # Section 2: Calibration
        hdr2 = QtWidgets.QLabel("Calibration")
        hdr2.setStyleSheet(_sec_style)
        roi_layout.addWidget(hdr2)

        row, self._lbl_alpha = _make_form_row("ROI-match α", _value_accent)
        roi_layout.addLayout(row)

        div2 = QtWidgets.QFrame()
        div2.setStyleSheet(f"background-color: {UI_BORDER}; min-height:1px; max-height:1px; border:none;")
        roi_layout.addWidget(div2)

        # Section 3: Diff Quantification
        hdr3 = QtWidgets.QLabel("Diff Quantification")
        hdr3.setStyleSheet(_sec_style)
        roi_layout.addWidget(hdr3)

        row, self._lbl_target_mean = _make_form_row("Target Mean Diff")
        roi_layout.addLayout(row)
        row, self._lbl_ref_mean = _make_form_row("Ref Mean Diff")
        roi_layout.addLayout(row)
        row, self._lbl_ref_std = _make_form_row("Ref Std Diff")
        roi_layout.addLayout(row)
        row, self._lbl_delta = _make_form_row("Δ (Target − Ref)")
        roi_layout.addLayout(row)
        row, self._lbl_snr = _make_form_row("SNR (Δ / σ_ref)", _value_accent)
        roi_layout.addLayout(row)

        roi_layout.addStretch(1)

        btn_details = QtWidgets.QPushButton("ROI Details…")
        btn_details.setFixedHeight(26)
        btn_details.setStyleSheet(
            f"QPushButton {{ background-color: {UI_PRIMARY}; color: #FFFFFF;"
            f" border: none; border-radius: {BorderRadius.SM};"
            f" font-size: {Typography.FONT_SIZE_SMALL}; font-weight: {Typography.FONT_WEIGHT_SEMIBOLD}; }}"
            f"QPushButton:hover {{ background-color: {UI_PRIMARY_HOVER}; }}"
        )
        btn_details.clicked.connect(self.roi_details_requested)
        roi_layout.addWidget(btn_details)

        self._stack.addWidget(page_roi)

        # ── Page 1: No-ROI fallback ────────────────────────────────────────
        page_no_roi = QtWidgets.QWidget()
        no_roi_layout = QtWidgets.QVBoxLayout(page_no_roi)
        no_roi_layout.setContentsMargins(0, 8, 0, 0)
        no_roi_layout.setSpacing(8)
        no_roi_layout.setAlignment(Qt.AlignHCenter | Qt.AlignVCenter)

        icon_lbl = QtWidgets.QLabel("◻")
        icon_lbl.setAlignment(Qt.AlignCenter)
        icon_lbl.setStyleSheet(
            f"color: {UI_TEXT_MUTED}; font-size: 28px; border: none; background: transparent;"
        )
        no_roi_layout.addWidget(icon_lbl)

        msg1 = QtWidgets.QLabel("No ROI defined.")
        msg1.setAlignment(Qt.AlignCenter)
        msg1.setStyleSheet(
            f"color: {UI_TEXT_SECONDARY}; font-size: {Typography.FONT_SIZE_BODY};"
            f" font-weight: {Typography.FONT_WEIGHT_SEMIBOLD}; border: none; background: transparent;"
        )
        no_roi_layout.addWidget(msg1)

        msg2 = QtWidgets.QLabel("Define ROI in ROI Manager\nto enable defect quantification.")
        msg2.setAlignment(Qt.AlignCenter)
        msg2.setStyleSheet(
            f"color: {UI_TEXT_MUTED}; font-size: {Typography.FONT_SIZE_SMALL};"
            f" border: none; background: transparent;"
        )
        no_roi_layout.addWidget(msg2)

        no_roi_layout.addStretch(1)

        btn_open_roi = QtWidgets.QPushButton("Open ROI Manager")
        btn_open_roi.setFixedHeight(26)
        btn_open_roi.setStyleSheet(
            f"QPushButton {{ background-color: {UI_BG_SUBTLE}; color: {UI_TEXT};"
            f" border: 1px solid {UI_BORDER}; border-radius: {BorderRadius.SM};"
            f" font-size: {Typography.FONT_SIZE_SMALL}; }}"
            f"QPushButton:hover {{ border-color: {UI_PRIMARY}; color: {UI_PRIMARY}; }}"
        )
        btn_open_roi.clicked.connect(self.open_roi_manager_requested)
        no_roi_layout.addWidget(btn_open_roi)

        self._stack.addWidget(page_no_roi)

        # Default: show no-ROI fallback
        self._stack.setCurrentIndex(1)

    # ------------------------------------------------------------------
    # Public update API
    # ------------------------------------------------------------------

    def show_no_roi(self):
        """Switch to the no-ROI fallback page."""
        self._stack.setCurrentIndex(1)

    def update_result(self, result: "SinglePairResult", roi_full: "Optional[ROIFullResult]",
                      n_target: int, n_ref: int):
        """Update the center panel for the currently displayed result.

        Parameters
        ----------
        result:    The SinglePairResult for the displayed pair.
        roi_full:  The ROIFullResult from the last ROI analysis run, or None.
        n_target:  Number of target ROIs currently defined.
        n_ref:     Number of reference ROIs currently defined.
        """
        if roi_full is None or (n_target == 0 and n_ref == 0):
            self.show_no_roi()
            return

        self._stack.setCurrentIndex(0)

        # Section 1: ROI summary
        self._lbl_target_count.setText(str(n_target))
        self._lbl_ref_count.setText(str(n_ref))
        self._lbl_roi_mode.setText("Reference-only" if n_ref > 0 else "—")

        # Section 2: Calibration alpha
        alpha = result.roi_match_alpha
        self._lbl_alpha.setText(f"{alpha:.4f}" if alpha is not None else "N/A (no ROI-match)")

        # Section 3: Diff quantification — look up this pair's SNR entry
        compare_label = result.compare_label
        snr_entry = roi_full.snr_per_diff.get(compare_label)
        if snr_entry is not None:
            # Stored values are in [0, 1] normalized range; scale to GLV (0-255)
            # for display so they match the visual diff image brightness.
            mu_t = snr_entry.mu_target * 255.0
            mu_r = snr_entry.mu_ref * 255.0
            sigma_r = snr_entry.sigma_ref * 255.0
            delta = mu_t - mu_r
            snr_val = snr_entry.snr

            self._lbl_target_mean.setText(f"{mu_t:.2f}")
            self._lbl_ref_mean.setText(f"{mu_r:.2f}")
            self._lbl_ref_std.setText(f"{sigma_r:.2f}")
            self._lbl_delta.setText(f"{delta:+.2f}")
            self._lbl_snr.setText(f"{snr_val:.3f}")
        else:
            for lbl in (self._lbl_target_mean, self._lbl_ref_mean,
                        self._lbl_ref_std, self._lbl_delta, self._lbl_snr):
                lbl.setText("--")

    def reset(self):
        """Reset to no-ROI fallback state."""
        self.show_no_roi()


class GLVMaskPreviewDialog(QtWidgets.QDialog):
    """Interactive dialog that visualises the GLV-Mask on the Base image.

    The user can adjust GLV Low / High directly inside the dialog and see the
    mask update in real-time.  Clicking 'Apply to Main' syncs the values back
    to the parent dialog's spinboxes.

    Visualisation:
      • In-mask pixels  : green-tinted highlight (R×0.65, G×1.0, B×0.65)
      • Out-of-mask pixels: dimmed to 8 % so the selected band stands out sharply
    """

    def __init__(
            self,
            base_img: np.ndarray,
            glv_low: int,
            glv_high: int,
            parent=None,
            apply_callback=None,  # callable(low: int, high: int) → None
    ):
        super().__init__(parent)
        self._base_img = base_img
        self._apply_callback = apply_callback
        self.setWindowTitle("GLV Mask Preview")
        self.setMinimumSize(540, 600)
        self.setStyleSheet(DIALOG_STYLE)
        self._setup_ui(glv_low, glv_high)
        self._update_preview()

    # ── UI construction ───────────────────────────────────────────────────────

    def _setup_ui(self, glv_low: int, glv_high: int):
        layout = QtWidgets.QVBoxLayout(self)
        layout.setSpacing(8)
        layout.setContentsMargins(12, 12, 12, 12)

        # ── GLV range controls ────────────────────────────────────────────────
        ctrl_row = QtWidgets.QHBoxLayout()
        ctrl_row.addWidget(QtWidgets.QLabel("GLV range:"))

        self.spn_low = QtWidgets.QSpinBox()
        self.spn_low.setRange(0, 254)
        self.spn_low.setValue(glv_low)
        self.spn_low.setFixedWidth(65)
        self.spn_low.setToolTip("Lower bound (inclusive, 0–255)")
        ctrl_row.addWidget(self.spn_low)

        ctrl_row.addWidget(QtWidgets.QLabel("–"))

        self.spn_high = QtWidgets.QSpinBox()
        self.spn_high.setRange(1, 255)
        self.spn_high.setValue(glv_high)
        self.spn_high.setFixedWidth(65)
        self.spn_high.setToolTip("Upper bound (inclusive, 0–255)")
        ctrl_row.addWidget(self.spn_high)

        ctrl_row.addStretch()
        layout.addLayout(ctrl_row)

        # ── Image display ─────────────────────────────────────────────────────
        self._lbl_image = QtWidgets.QLabel()
        self._lbl_image.setAlignment(Qt.AlignCenter)
        self._lbl_image.setMinimumSize(480, 420)
        layout.addWidget(self._lbl_image, stretch=1)

        # ── Info bar (pixel count / %) ────────────────────────────────────────
        self._lbl_info = QtWidgets.QLabel()
        self._lbl_info.setAlignment(Qt.AlignCenter)
        layout.addWidget(self._lbl_info)

        # ── Buttons ───────────────────────────────────────────────────────────
        btn_row = QtWidgets.QHBoxLayout()
        btn_row.addStretch()
        if self._apply_callback is not None:
            self._btn_apply = QtWidgets.QPushButton("Apply to Main")
            self._btn_apply.setToolTip("Copy the current GLV range back to the main dialog")
            self._btn_apply.clicked.connect(self._on_apply)
            btn_row.addWidget(self._btn_apply)
        btn_close = QtWidgets.QPushButton("Close")
        btn_close.clicked.connect(self.accept)
        btn_row.addWidget(btn_close)
        layout.addLayout(btn_row)

        # Live update on spinbox change
        self.spn_low.valueChanged.connect(self._update_preview)
        self.spn_high.valueChanged.connect(self._update_preview)

    # ── Preview rendering ─────────────────────────────────────────────────────

    def _update_preview(self):
        """Recompute mask and refresh the image label."""
        glv_low = self.spn_low.value()
        glv_high = self.spn_high.value()

        img_f = self._base_img.astype(np.float32)
        mask = (img_f >= glv_low) & (img_f <= glv_high)

        # Stats
        n_in = int(mask.sum())
        pct = 100.0 * n_in / max(mask.size, 1)
        status = "too few pixels — will fall back to full-image Percentile" if n_in < 50 else f"{n_in:,} px in mask  ({pct:.1f} %)"
        self._lbl_info.setText(f"GLV {glv_low} – {glv_high}   |   {status}")

        # Build RGB visualisation
        rgb = self._build_rgb_preview(self._base_img, mask)

        h, w = rgb.shape[:2]
        bytes_per_line = 3 * w
        qimg = QtGui.QImage(rgb.tobytes(), w, h, bytes_per_line, QtGui.QImage.Format_RGB888)
        pixmap = QtGui.QPixmap.fromImage(qimg)

        avail = self._lbl_image.size()
        scaled = pixmap.scaled(
            avail.width() - 4, avail.height() - 4,
            Qt.KeepAspectRatio, Qt.SmoothTransformation
        )
        self._lbl_image.setPixmap(scaled)

        self.setWindowTitle(f"GLV Mask Preview  [{glv_low} – {glv_high}]")

    def resizeEvent(self, event):
        super().resizeEvent(event)
        self._update_preview()

    @staticmethod
    def _build_rgb_preview(img: np.ndarray, mask: np.ndarray) -> np.ndarray:
        """Return uint8 RGB array.

        In-mask  : green highlight (R×0.65, G×1.0, B×0.65)
        Out-mask : dimmed to 8 % as neutral grey
        """
        img_f = img.astype(np.float32)
        h, w = img_f.shape[:2]
        rgb = np.zeros((h, w, 3), dtype=np.uint8)

        # Out-of-mask: dim to 8 %
        dim = np.clip(img_f * 0.08, 0, 255).astype(np.uint8)
        rgb[..., 0] = dim
        rgb[..., 1] = dim
        rgb[..., 2] = dim

        # In-mask: green-tinted highlight
        v = img_f[mask]
        rgb[mask, 0] = np.clip(v * 0.65, 0, 255).astype(np.uint8)
        rgb[mask, 1] = np.clip(v, 0, 255).astype(np.uint8)
        rgb[mask, 2] = np.clip(v * 0.65, 0, 255).astype(np.uint8)

        return rgb

    # ── Actions ───────────────────────────────────────────────────────────────

    def _on_apply(self):
        """Sync current spinbox values back to the parent dialog."""
        self._apply_callback(self.spn_low.value(), self.spn_high.value())


_NORMALIZE_METHOD_LABELS = {
    'percentile': 'Percentile (P2–P98)',
    'glv_mask': 'GLV-Mask',
    'skip': 'Skip (raw ÷ 255)',
    'roi_match': 'ROI-Match (EPI Nulling)',
}


class NormalizedCompareDialog(QtWidgets.QDialog):
    """Dialog to preview normalization effect: shows raw compare, normalized compare, and base."""

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Normalize Preview")
        self.resize(900, 800)
        self.setStyleSheet(DIALOG_STYLE)

        layout = QtWidgets.QVBoxLayout(self)
        layout.setContentsMargins(12, 12, 12, 12)
        layout.setSpacing(8)

        self.lbl_info = QtWidgets.QLabel("")
        self.lbl_info.setStyleSheet(f"color: {UI_TEXT_SECONDARY}; font-size: {Typography.FONT_SIZE_SMALL};")
        layout.addWidget(self.lbl_info)

        label_row = QtWidgets.QHBoxLayout()
        self.lbl_compare_title = QtWidgets.QLabel("Aligned Compare (Before)")
        self.lbl_norm_title = QtWidgets.QLabel("After Normalize")
        self.lbl_base_title = QtWidgets.QLabel("Base (Target)")
        for _lbl in (self.lbl_compare_title, self.lbl_norm_title, self.lbl_base_title):
            _lbl.setAlignment(Qt.AlignCenter)
            _lbl.setStyleSheet(
                f"font-size: {Typography.FONT_SIZE_SMALL}; font-weight: {Typography.FONT_WEIGHT_BOLD};"
                f" color: {UI_TEXT};"
            )
            label_row.addWidget(_lbl, 1)
        layout.addLayout(label_row)

        image_row = QtWidgets.QHBoxLayout()
        self.img_compare = SyncZoomImageWidget("Aligned Compare (Before)")
        self.img_norm_compare = SyncZoomImageWidget("After Normalize")
        self.img_base = SyncZoomImageWidget("Base (Target)")
        image_row.addWidget(self.img_compare, 1)
        image_row.addWidget(self.img_norm_compare, 1)
        image_row.addWidget(self.img_base, 1)
        layout.addLayout(image_row, 1)

        hist_row = QtWidgets.QHBoxLayout()
        self.hist_compare = HistogramCanvas()
        self.hist_compare.setFixedHeight(140)
        self.hist_norm_compare = HistogramCanvas()
        self.hist_norm_compare.setFixedHeight(140)
        self.hist_base = HistogramCanvas()
        self.hist_base.setFixedHeight(140)
        hist_row.addWidget(self.hist_compare, 1)
        hist_row.addWidget(self.hist_norm_compare, 1)
        hist_row.addWidget(self.hist_base, 1)
        layout.addLayout(hist_row)

        button_row = QtWidgets.QHBoxLayout()
        button_row.addStretch()
        btn_close = QtWidgets.QPushButton("Close")
        btn_close.clicked.connect(self.close)
        button_row.addWidget(btn_close)
        layout.addLayout(button_row)

    def set_images(
            self,
            compare_image: np.ndarray,
            normalized_compare: np.ndarray,
            base_image: np.ndarray,
            base_label: str,
            compare_label: str,
            norm_a: float,
            norm_b: float,
            method_name: str = 'percentile',
    ):
        method_display = _NORMALIZE_METHOD_LABELS.get(method_name, method_name)
        is_linear = method_name in ('percentile', 'glv_mask')
        if is_linear:
            coeff_str = f"  |  a={norm_a:.4f}, b={norm_b:.1f}"
        else:
            coeff_str = "  |  (non-linear)"
        self.lbl_info.setText(
            f"Method: {method_display}  |  Compare: {compare_label}  →  Base: {base_label}{coeff_str}"
        )
        self.lbl_compare_title.setText(f"Aligned Compare — Before\n({compare_label})")
        self.lbl_norm_title.setText(f"After Normalize\n[{method_display}]")
        self.lbl_base_title.setText(f"Base — Target\n({base_label})")
        self.setWindowTitle(f"Normalize Preview  [{method_display}]  —  {compare_label}")
        if compare_image is not None:
            self.img_compare.setImage(compare_image)
        if normalized_compare is not None:
            self.img_norm_compare.setImage(normalized_compare)
        if base_image is not None:
            self.img_base.setImage(base_image)

        self._update_histogram(self.hist_compare, compare_image)
        self._update_histogram(self.hist_norm_compare, normalized_compare)
        self._update_histogram(self.hist_base, base_image)

    @staticmethod
    def _update_histogram(canvas: HistogramCanvas, image: Optional[np.ndarray]):
        if image is None or image.size == 0:
            canvas.plot_histogram(np.zeros(256), np.arange(257))
            return
        img_uint8 = image.astype(np.uint8)
        counts, edges = np.histogram(img_uint8.flatten(), bins=256, range=(0, 256))
        canvas.plot_histogram(counts, edges)


class _NumericSortItem(QtWidgets.QTableWidgetItem):
    """QTableWidgetItem that compares numerically when a sort_key is supplied."""

    def __init__(self, text: str, sort_key=None):
        super().__init__(text)
        self._sort_key = sort_key  # float or None

    def __lt__(self, other: QtWidgets.QTableWidgetItem) -> bool:
        if isinstance(other, _NumericSortItem):
            a, b = self._sort_key, other._sort_key
            if a is not None and b is not None:
                return float(a) < float(b)
            if a is None and b is not None:
                return True   # None sorts before numbers
            if a is not None and b is None:
                return False
        return super().__lt__(other)


class ROIIntensityProfileDialog(QtWidgets.QDialog):
    """Detailed ROI analysis dialog — opened via [ROI Details…] button.

    Supports multi-base results (auto-pair mode) by grouping per base_label.

    Tabs
    ----
    1. LE Summary     – engineering table (one row per pair) + CSV export.
    2. SNR Chart      – bar chart with Δ subplot and σ_ref error bars.
    3. Per-ROI Mean   – selectable ROI; base / compare / diff mean lines.
    4. Raw Table      – all layer × ROI stats (original detailed view).
    """

    # Shared axis / spine style helpers
    _BG_FIG  = '#1F2937'
    _BG_AX   = '#111827'
    _COL_TXT = '#D1D5DB'
    _COL_MUT = '#9CA3AF'
    _COL_SPL = '#4B5563'
    _COL_GRD = '#374151'

    # SNR quality thresholds for row color-coding
    _SNR_GOOD  = 2.0   # ≥ this → green
    _SNR_OK    = 1.0   # ≥ this → yellow; below → red

    # Align status colors
    _STATUS_COLOR = {
        'ok':   '#4ADE80',
        'warn': '#FCD34D',
        'fail': '#F87171',
    }

    def __init__(
        self,
        roi_results: Dict[str, ROIFullResult],
        all_results: List[SinglePairResult],
        parent=None,
    ):
        super().__init__(parent)
        # roi_results: base_label → ROIFullResult
        # all_results: full list of SinglePairResult for alpha / align_score lookup
        self._roi_results = roi_results
        self._all_results = all_results
        self._base_labels = list(roi_results.keys())

        self.setWindowTitle("ROI Intensity Analysis  —  ROI Details")
        self.setWindowFlags(self.windowFlags() | Qt.Window)
        self.resize(1040, 680)
        self._build_ui()

    # ------------------------------------------------------------------
    # Top-level layout
    # ------------------------------------------------------------------

    @property
    def _is_auto_pair(self) -> bool:
        """True when results come from auto-pair mode (multiple distinct base images)."""
        return len(set(r.base_label for r in self._all_results)) > 1

    # ------------------------------------------------------------------
    # Light-theme stylesheet (matches main UI DIALOG_STYLE palette)
    # ------------------------------------------------------------------
    _DIALOG_QSS = """
        QDialog, QWidget {
            background-color: #F3F4F6;
            color: #111827;
            font-family: 'Liberation Sans', Arial, 'Helvetica Neue', 'Segoe UI';
            font-size: 11px;
        }
        QTabWidget::pane {
            border: 1px solid #E5E7EB;
            background-color: #FFFFFF;
            border-radius: 0px 4px 4px 4px;
        }
        QTabBar::tab {
            background-color: #F9FAFB;
            color: #6B7280;
            padding: 6px 16px;
            border: 1px solid #E5E7EB;
            border-bottom: none;
            border-top-left-radius: 5px;
            border-top-right-radius: 5px;
            margin-right: 2px;
            font-size: 11px;
        }
        QTabBar::tab:selected {
            background-color: #FFFFFF;
            color: #111827;
            border-bottom: 2px solid #F59E0B;
            font-weight: bold;
        }
        QTabBar::tab:hover:!selected {
            background-color: #FFF8ED;
            color: #374151;
        }
        QTableWidget {
            background-color: #FFFFFF;
            alternate-background-color: #F9FAFB;
            color: #111827;
            gridline-color: #E5E7EB;
            border: 1px solid #E5E7EB;
            border-radius: 4px;
            selection-background-color: #FEF3C7;
            selection-color: #111827;
            font-size: 11px;
        }
        QHeaderView::section {
            background-color: #F3F4F6;
            color: #6B7280;
            padding: 5px 8px;
            border: none;
            border-right: 1px solid #E5E7EB;
            border-bottom: 1px solid #D1D5DB;
            font-weight: bold;
            font-size: 11px;
        }
        QComboBox {
            background-color: #FFFFFF;
            color: #111827;
            border: 1px solid #D1D5DB;
            border-radius: 4px;
            padding: 3px 10px;
            min-height: 24px;
            font-size: 11px;
        }
        QComboBox::drop-down {
            border: none;
            width: 20px;
        }
        QComboBox QAbstractItemView {
            background-color: #FFFFFF;
            color: #111827;
            selection-background-color: #FEF3C7;
            border: 1px solid #D1D5DB;
        }
        QPushButton {
            background-color: #FFFFFF;
            color: #374151;
            border: 1px solid #D1D5DB;
            border-radius: 5px;
            padding: 5px 16px;
            min-height: 26px;
            font-size: 11px;
        }
        QPushButton:hover {
            background-color: #FFF8ED;
            color: #111827;
            border-color: #F59E0B;
        }
        QPushButton:pressed {
            background-color: #FDE68A;
            border-color: #D97706;
        }
        QLabel {
            color: #6B7280;
            font-size: 11px;
        }
        QScrollBar:vertical {
            background-color: #F3F4F6;
            width: 10px;
            margin: 0;
            border-radius: 5px;
        }
        QScrollBar::handle:vertical {
            background-color: #D1D5DB;
            border-radius: 5px;
            min-height: 24px;
        }
        QScrollBar::handle:vertical:hover { background-color: #9CA3AF; }
        QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical { height: 0; }
        QScrollBar:horizontal {
            background-color: #F3F4F6;
            height: 10px;
            margin: 0;
            border-radius: 5px;
        }
        QScrollBar::handle:horizontal {
            background-color: #D1D5DB;
            border-radius: 5px;
            min-width: 24px;
        }
        QScrollBar::handle:horizontal:hover { background-color: #9CA3AF; }
        QScrollBar::add-line:horizontal, QScrollBar::sub-line:horizontal { width: 0; }
    """

    def _build_ui(self) -> None:
        self.setStyleSheet(self._DIALOG_QSS)

        root = QtWidgets.QVBoxLayout(self)
        root.setSpacing(0)
        root.setContentsMargins(0, 0, 0, 0)

        # ── Header banner ─────────────────────────────────────────────
        header = QtWidgets.QWidget()
        header.setFixedHeight(44)
        header.setStyleSheet(
            "background: qlineargradient(x1:0,y1:0,x2:1,y2:0,"
            "stop:0 #FEF3C7, stop:1 #FFF8ED);"
            "border-bottom: 2px solid #F59E0B;"
        )
        h_lay = QtWidgets.QHBoxLayout(header)
        h_lay.setContentsMargins(16, 0, 16, 0)
        lbl_title = QtWidgets.QLabel("ROI Intensity Analysis")
        lbl_title.setStyleSheet(
            "color: #111827; font-size: 15px; font-weight: bold; background: transparent;"
        )
        n_bases = len(self._base_labels)
        n_pairs = len(self._all_results)
        lbl_info = QtWidgets.QLabel(
            f"  {n_bases} base image{'s' if n_bases != 1 else ''}  ·  {n_pairs} pair{'s' if n_pairs != 1 else ''}"
        )
        lbl_info.setStyleSheet(
            "color: #92400E; font-size: 11px; background: transparent;"
        )
        h_lay.addWidget(lbl_title)
        h_lay.addWidget(lbl_info)
        h_lay.addStretch()
        root.addWidget(header)

        # ── Tab area ──────────────────────────────────────────────────
        inner = QtWidgets.QWidget()
        inner_lay = QtWidgets.QVBoxLayout(inner)
        inner_lay.setSpacing(8)
        inner_lay.setContentsMargins(10, 8, 10, 8)

        tabs = QtWidgets.QTabWidget()
        tabs.addTab(self._build_summary_tab(),    "📋  Pair Summary")
        if self._is_auto_pair:
            tabs.addTab(self._build_matrix_tab(),      "🔢  SNR Pair Matrix")
            tabs.addTab(self._build_diff_matrix_tab(), "🖼  Diff Image Matrix")
        tabs.addTab(self._build_snr_chart_tab(), "📊  SNR Bar Chart")
        tabs.addTab(self._build_mean_tab(),      "📈  Intensity Profile")
        tabs.addTab(self._build_table_tab(),     "🗂  Raw Stats")
        inner_lay.addWidget(tabs, stretch=1)

        bottom = QtWidgets.QHBoxLayout()
        btn_close = QtWidgets.QPushButton("Close")
        btn_close.setFixedWidth(100)
        btn_close.setStyleSheet(
            "QPushButton { background:#FFFFFF; color:#374151; border:1px solid #D1D5DB;"
            " border-radius:5px; padding:5px 16px; font-weight:bold; }"
            "QPushButton:hover { background:#FFF8ED; border-color:#F59E0B; color:#92400E; }"
            "QPushButton:pressed { background:#FDE68A; }"
        )
        btn_close.clicked.connect(self.close)
        bottom.addStretch()
        bottom.addWidget(btn_close)
        inner_lay.addLayout(bottom)

        root.addWidget(inner, stretch=1)

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _make_base_selector(self, label: str = "Base:") -> Tuple[
            QtWidgets.QHBoxLayout, Optional[QtWidgets.QComboBox]]:
        """Return a (layout, combo) pair.  combo is None when only one base exists."""
        row = QtWidgets.QHBoxLayout()
        if len(self._base_labels) <= 1:
            return row, None
        row.addWidget(QtWidgets.QLabel(label))
        cmb = QtWidgets.QComboBox()
        cmb.addItems(self._base_labels)
        row.addWidget(cmb)
        row.addStretch()
        return row, cmb

    def _current_roi_result(self, cmb: Optional[QtWidgets.QComboBox]) -> Optional[ROIFullResult]:
        key = cmb.currentText() if cmb is not None else (self._base_labels[0] if self._base_labels else None)
        return self._roi_results.get(key) if key else None

    def _style_ax(self, ax) -> None:
        ax.set_facecolor(self._BG_AX)
        ax.tick_params(colors=self._COL_TXT)
        for spine in ('bottom', 'left'):
            ax.spines[spine].set_color(self._COL_SPL)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(True, color=self._COL_GRD, linewidth=0.5)

    def _results_for_base(self, base_label: str) -> List[SinglePairResult]:
        return [r for r in self._all_results if r.base_label == base_label]

    def _compute_base_snr(self, base_label: str) -> Optional[float]:
        """Return Base SNR for *base_label* (same formula as SNR Pair Matrix on raw base image).

        SNR = max(0, (μ_target − μ_ref) / σ_ref)
        where μ_ref = mean of reference-ROI means, σ_ref = std of those means (≥2 refs)
        or the single reference ROI's pixel std (1 ref).
        """
        import numpy as _np
        roi_full = self._roi_results.get(base_label)
        if roi_full is None:
            return None
        base_layer = roi_full.get_base_layer()
        target_roi = roi_full.roi_set.get_target()
        ref_rois   = roi_full.roi_set.get_references()
        if not (target_roi and ref_rois and base_layer):
            return None
        t_stat = base_layer.roi_stats.get(target_roi.id)
        ref_stats = [base_layer.roi_stats[r.id]
                     for r in ref_rois if r.id in base_layer.roi_stats]
        if not (t_stat and ref_stats):
            return None
        ref_means = _np.array([rs.mean for rs in ref_stats], dtype=_np.float32)
        mu_t  = float(t_stat.mean)
        mu_r  = float(_np.mean(ref_means))
        sigma = float(_np.std(ref_means)) if len(ref_means) >= 2 else float(ref_stats[0].std)
        if sigma <= 1e-7:
            return 0.0
        return max(0.0, (mu_t - mu_r) / sigma)

    def _save_figure(self, fig, default_stem: str = "roi_chart") -> None:
        """Open a save dialog and write *fig* to PNG / PDF / SVG."""
        if fig is None:
            return
        path, _ = QtWidgets.QFileDialog.getSaveFileName(
            self, "Save Chart", default_stem,
            "PNG Image (*.png);;PDF Document (*.pdf);;SVG Vector (*.svg)"
        )
        if not path:
            return
        try:
            fig.savefig(path, dpi=150, bbox_inches='tight',
                        facecolor=self._BG_FIG, edgecolor='none')
            QtWidgets.QMessageBox.information(self, "Save Chart", f"Saved to:\n{path}")
        except Exception as exc:
            QtWidgets.QMessageBox.warning(self, "Save Chart", f"Failed:\n{exc}")

    def _make_save_btn(self, get_fig_fn, stem: str = "chart") -> QtWidgets.QPushButton:
        btn = QtWidgets.QPushButton("Save Chart…")
        btn.setFixedHeight(24)
        btn.setMaximumWidth(120)
        btn.clicked.connect(lambda: self._save_figure(get_fig_fn(), stem))
        return btn

    # ------------------------------------------------------------------
    # Tab 1 — LE Summary (engineering table + CSV export)
    # ------------------------------------------------------------------

    def _build_summary_tab(self) -> QtWidgets.QWidget:
        w = QtWidgets.QWidget()
        lay = QtWidgets.QVBoxLayout(w)
        lay.setSpacing(6)

        ctrl_row, self._cmb_summary_base = self._make_base_selector("Filter by Base:")
        if self._cmb_summary_base is not None:
            self._cmb_summary_base.currentTextChanged.connect(self._refresh_summary_table)
        export_btn = QtWidgets.QPushButton("Export CSV…")
        export_btn.setFixedHeight(26)
        export_btn.clicked.connect(self._on_export_csv)
        ctrl_row.addWidget(export_btn)
        lay.addLayout(ctrl_row)

        # Column indices (keep in sync with _SUMMARY_HEADERS below)
        self._SUMMARY_HEADERS = [
            'Base', 'Compare (LE)', 'Align Status', 'ROI-match α', 'Align Score',
            'T Mean Diff', 'R Mean Diff', 'R Std Diff', 'Δ (T−R)', 'Pair SNR', 'Base SNR',
        ]
        self._COL_SNR      = self._SUMMARY_HEADERS.index('Pair SNR')
        self._COL_BASE_SNR = self._SUMMARY_HEADERS.index('Base SNR')
        self._COL_STATUS   = self._SUMMARY_HEADERS.index('Align Status')

        self._summary_table = QtWidgets.QTableWidget(0, len(self._SUMMARY_HEADERS))
        self._summary_table.setHorizontalHeaderLabels(self._SUMMARY_HEADERS)
        self._summary_table.horizontalHeader().setSectionResizeMode(
            QtWidgets.QHeaderView.ResizeToContents
        )
        self._summary_table.horizontalHeader().setStretchLastSection(True)
        self._summary_table.setEditTriggers(QtWidgets.QAbstractItemView.NoEditTriggers)
        self._summary_table.setAlternatingRowColors(True)
        self._summary_table.setSelectionBehavior(QtWidgets.QAbstractItemView.SelectRows)
        self._summary_table.setSortingEnabled(True)
        self._summary_table.horizontalHeader().setSortIndicator(
            self._COL_SNR, Qt.DescendingOrder
        )
        lay.addWidget(self._summary_table, stretch=1)

        self._refresh_summary_table()
        return w

    def _get_summary_data(self, base_filter: Optional[str] = None) -> List[dict]:
        """Return structured data dicts for each pair (used by table and CSV)."""
        # Pre-compute Base SNR once per base label to avoid repeated work.
        _base_snr_cache: Dict[str, Optional[float]] = {}

        data = []
        for r in self._all_results:
            if base_filter is not None and r.base_label != base_filter:
                continue
            roi_full  = self._roi_results.get(r.base_label)
            snr_entry = roi_full.snr_per_diff.get(r.compare_label) if roi_full else None
            align_status  = (r.alignment.status if r.alignment else None) or '—'
            align_score_v = r.alignment.final_score if r.alignment else None
            alpha_str     = f"{r.roi_match_alpha:.4f}" if r.roi_match_alpha is not None else "—"
            align_str     = f"{align_score_v:.1f}" if align_score_v is not None else "—"

            # Base SNR (same for all pairs sharing the same base)
            if r.base_label not in _base_snr_cache:
                _base_snr_cache[r.base_label] = self._compute_base_snr(r.base_label)
            base_snr_v = _base_snr_cache[r.base_label]
            base_snr_s = f"{base_snr_v:.4f}" if base_snr_v is not None else "—"

            if snr_entry is not None:
                # Scale from [0,1] normalized range to GLV (0-255) for display.
                mu_t_glv    = snr_entry.mu_target * 255.0
                mu_r_glv    = snr_entry.mu_ref    * 255.0
                sigma_r_glv = snr_entry.sigma_ref * 255.0
                delta_v     = mu_t_glv - mu_r_glv
                data.append({
                    'base':         r.base_label,
                    'compare':      r.compare_label,
                    'align_status': align_status,
                    'alpha':        alpha_str,
                    'align_score':  align_str,
                    'align_score_v': align_score_v,
                    'mu_t':         f"{mu_t_glv:.2f}",
                    'mu_r':         f"{mu_r_glv:.2f}",
                    'sigma_r':      f"{sigma_r_glv:.2f}",
                    'delta':        f"{delta_v:+.2f}",
                    'delta_v':      delta_v,
                    'snr':          f"{snr_entry.snr:.4f}",
                    'snr_v':        snr_entry.snr,
                    'base_snr':     base_snr_s,
                    'base_snr_v':   base_snr_v,
                })
            else:
                data.append({
                    'base': r.base_label, 'compare': r.compare_label,
                    'align_status': align_status,
                    'alpha': alpha_str, 'align_score': align_str, 'align_score_v': align_score_v,
                    'mu_t': '—', 'mu_r': '—', 'sigma_r': '—',
                    'delta': '—', 'delta_v': None,
                    'snr': '—',  'snr_v': None,
                    'base_snr': base_snr_s, 'base_snr_v': base_snr_v,
                })
        return data

    def _get_summary_rows(self, base_filter: Optional[str] = None) -> List[List[str]]:
        """Flat string rows in SUMMARY_HEADERS column order (for CSV export)."""
        return [
            [d['base'], d['compare'], d['align_status'], d['alpha'], d['align_score'],
             d['mu_t'], d['mu_r'], d['sigma_r'], d['delta'], d['snr'], d['base_snr']]
            for d in self._get_summary_data(base_filter)
        ]

    @staticmethod
    def _snr_bg(snr_v: Optional[float]) -> Optional[QtGui.QColor]:
        if snr_v is None:
            return None
        if snr_v >= ROIIntensityProfileDialog._SNR_GOOD:
            return QtGui.QColor('#166534')   # dark green bg
        if snr_v >= ROIIntensityProfileDialog._SNR_OK:
            return QtGui.QColor('#854D0E')   # dark yellow/amber bg
        return QtGui.QColor('#7F1D1D')       # dark red bg

    @staticmethod
    def _status_bg(status: str) -> Optional[QtGui.QColor]:
        color_hex = ROIIntensityProfileDialog._STATUS_COLOR.get(status.lower())
        if color_hex is None:
            return None
        c = QtGui.QColor(color_hex)
        # Darken for use as background (text stays light)
        return QtGui.QColor(c.red() // 2, c.green() // 2, c.blue() // 2)

    def _make_sort_item(self, text: str, sort_key=None) -> QtWidgets.QTableWidgetItem:
        """QTableWidgetItem that sorts numerically when sort_key is a float."""
        item = _NumericSortItem(text, sort_key)
        item.setTextAlignment(Qt.AlignCenter)
        return item

    def _refresh_summary_table(self) -> None:
        base_filter: Optional[str] = None
        if self._cmb_summary_base is not None:
            base_filter = self._cmb_summary_base.currentText() or None

        # Disable sorting while filling to avoid mid-insert reordering
        self._summary_table.setSortingEnabled(False)
        self._summary_table.setRowCount(0)

        for d in self._get_summary_data(base_filter):
            row_idx = self._summary_table.rowCount()
            self._summary_table.insertRow(row_idx)

            # Build items in column order matching _SUMMARY_HEADERS
            items = [
                self._make_sort_item(d['base']),
                self._make_sort_item(d['compare']),
                self._make_sort_item(d['align_status']),
                self._make_sort_item(d['alpha']),
                self._make_sort_item(d['align_score'], d['align_score_v']),
                self._make_sort_item(d['mu_t']),
                self._make_sort_item(d['mu_r']),
                self._make_sort_item(d['sigma_r']),
                self._make_sort_item(d['delta'], d['delta_v']),
                self._make_sort_item(d['snr'], d['snr_v']),
                self._make_sort_item(d['base_snr'], d['base_snr_v']),
            ]

            # Color: Align Status cell
            status_bg = self._status_bg(d['align_status'])
            if status_bg is not None:
                items[self._COL_STATUS].setBackground(status_bg)
                items[self._COL_STATUS].setForeground(
                    QtGui.QColor(self._STATUS_COLOR.get(d['align_status'].lower(), '#D1D5DB'))
                )

            # Color: SNR cell
            snr_bg = self._snr_bg(d['snr_v'])
            if snr_bg is not None:
                items[self._COL_SNR].setBackground(snr_bg)
                items[self._COL_SNR].setForeground(QtGui.QColor('#D1D5DB'))

            # Color: Base SNR cell
            base_snr_bg = self._snr_bg(d['base_snr_v'])
            if base_snr_bg is not None:
                items[self._COL_BASE_SNR].setBackground(base_snr_bg)
                items[self._COL_BASE_SNR].setForeground(QtGui.QColor('#D1D5DB'))

            for col, item in enumerate(items):
                self._summary_table.setItem(row_idx, col, item)

        self._summary_table.setSortingEnabled(True)

    def _on_export_csv(self) -> None:
        """Export the LE Summary table to a CSV file chosen by the user."""
        path, _ = QtWidgets.QFileDialog.getSaveFileName(
            self, "Export LE Summary", "", "CSV Files (*.csv)"
        )
        if not path:
            return
        headers = [
            'base', 'compare_le', 'align_status', 'roi_match_alpha', 'align_score',
            'target_mean_diff_glv', 'ref_mean_diff_glv', 'ref_std_diff_glv', 'delta_glv',
            'snr', 'base_snr',
        ]
        try:
            with open(path, 'w', newline='', encoding='utf-8') as f:
                import csv
                writer = csv.writer(f)
                writer.writerow(headers)
                writer.writerows(self._get_summary_rows())  # all bases
            QtWidgets.QMessageBox.information(self, "Export CSV", f"Saved to:\n{path}")
        except Exception as exc:
            QtWidgets.QMessageBox.warning(self, "Export CSV", f"Failed to save:\n{exc}")

    # ------------------------------------------------------------------
    # Tab 2 — SNR Chart (bar + Δ subplot; "All Bases" overlay mode)
    # ------------------------------------------------------------------

    # Colors used for multi-base line overlay (cycles if > 8 bases)
    _MULTI_COLORS = [
        '#F59E0B', '#60A5FA', '#34D399', '#F87171',
        '#A78BFA', '#FB923C', '#2DD4BF', '#E879F9',
    ]
    _ALL_BASES_KEY = "All Bases"

    def _build_snr_chart_tab(self) -> QtWidgets.QWidget:
        w = QtWidgets.QWidget()
        lay = QtWidgets.QVBoxLayout(w)
        lay.setSpacing(4)

        ctrl_row = QtWidgets.QHBoxLayout()
        self._cmb_chart_base: Optional[QtWidgets.QComboBox] = None
        if len(self._base_labels) > 1:
            ctrl_row.addWidget(QtWidgets.QLabel("Base:"))
            self._cmb_chart_base = QtWidgets.QComboBox()
            self._cmb_chart_base.addItem(self._ALL_BASES_KEY)
            self._cmb_chart_base.addItems(self._base_labels)
            self._cmb_chart_base.currentTextChanged.connect(self._refresh_snr_chart)
            ctrl_row.addWidget(self._cmb_chart_base)
        ctrl_row.addStretch()
        ctrl_row.addWidget(self._make_save_btn(lambda: self._chart_canvas.figure, "roi_snr_chart"))
        lay.addLayout(ctrl_row)

        fig = Figure(figsize=(6, 4), tight_layout=True)
        fig.patch.set_facecolor(self._BG_FIG)
        self._chart_ax_snr, self._chart_ax_delta = fig.subplots(2, 1, sharex=True)
        self._chart_canvas = FigureCanvas(fig)
        lay.addWidget(self._chart_canvas, stretch=1)

        self._refresh_snr_chart()
        return w

    def _refresh_snr_chart(self) -> None:
        sel = self._cmb_chart_base.currentText() if self._cmb_chart_base else ""
        if sel == self._ALL_BASES_KEY:
            self._draw_snr_all_bases()
        else:
            self._draw_snr_single_base(sel or (self._base_labels[0] if self._base_labels else ""))

    def _draw_snr_single_base(self, base_label: str) -> None:
        """Bar chart mode — one base, bars per compare LE."""
        roi_result = self._roi_results.get(base_label)
        ax_snr, ax_delta = self._chart_ax_snr, self._chart_ax_delta
        for ax in (ax_snr, ax_delta):
            ax.cla()
            self._style_ax(ax)

        if roi_result is None or not roi_result.snr_per_diff:
            ax_snr.text(0.5, 0.5, "No ROI data available",
                        ha='center', va='center', color=self._COL_MUT,
                        transform=ax_snr.transAxes)
            self._chart_canvas.draw()
            return

        snr_data   = roi_result.snr_per_diff
        labels     = list(snr_data.keys())
        xs         = list(range(len(labels)))
        snr_vals   = [snr_data[k].snr for k in labels]
        # Scale to GLV (0-255) so chart y-axis matches displayed table values.
        delta_vals = [(snr_data[k].mu_target - snr_data[k].mu_ref) * 255.0 for k in labels]
        sigma_refs = [snr_data[k].sigma_ref * 255.0 for k in labels]
        bar_w = 0.55

        bars = ax_snr.bar(xs, snr_vals, width=bar_w, color='#F59E0B', alpha=0.85, zorder=3)
        ax_snr.axhline(0, color=self._COL_SPL, linewidth=0.8, linestyle='--')
        ax_snr.set_ylabel("SNR", color=self._COL_TXT, fontsize=9)
        ax_snr.set_title(f"SNR = Δ / σ_Ref   [Base: {base_label}]",
                         color=self._COL_MUT, fontsize=9)
        spread = max(snr_vals) - min(snr_vals) if snr_vals else 1
        for i, (_, val) in enumerate(zip(bars, snr_vals)):
            ax_snr.text(i, val + spread * 0.03, f"{val:.2f}",
                        ha='center', va='bottom', fontsize=8, color=self._COL_TXT)

        bar_colors = ['#34D399' if d >= 0 else '#F87171' for d in delta_vals]
        ax_delta.bar(xs, delta_vals, width=bar_w, color=bar_colors, alpha=0.8, zorder=3, label='Δ (T−R)')
        ax_delta.errorbar(xs, delta_vals, yerr=sigma_refs,
                          fmt='none', color='#D1D5DB', capsize=5, linewidth=1.5, zorder=4, label='±σ_Ref')
        ax_delta.axhline(0, color=self._COL_SPL, linewidth=0.8, linestyle='--')
        ax_delta.set_ylabel("Δ = T−R", color=self._COL_TXT, fontsize=9)
        ax_delta.set_xticks(xs)
        ax_delta.set_xticklabels(labels, rotation=20, ha='right', color=self._COL_TXT, fontsize=9)
        ax_delta.legend(facecolor=self._BG_FIG, labelcolor=self._COL_TXT, fontsize=8, loc='upper right')
        self._chart_canvas.draw()

    def _draw_snr_all_bases(self) -> None:
        """Overlay mode — one colored line per base, shared compare-LE x-axis."""
        ax_snr, ax_delta = self._chart_ax_snr, self._chart_ax_delta
        for ax in (ax_snr, ax_delta):
            ax.cla()
            self._style_ax(ax)

        # Collect union of all compare labels (preserve insertion order)
        seen: dict = {}
        for base_lbl in self._base_labels:
            roi_res = self._roi_results.get(base_lbl)
            if roi_res:
                for k in roi_res.snr_per_diff:
                    seen[k] = None
        all_compare = list(seen.keys())

        if not all_compare:
            ax_snr.text(0.5, 0.5, "No ROI data available",
                        ha='center', va='center', color=self._COL_MUT,
                        transform=ax_snr.transAxes)
            self._chart_canvas.draw()
            return

        x_pos = {lbl: i for i, lbl in enumerate(all_compare)}
        markers = ['o', 's', '^', 'D', 'v', 'P', 'X', '*']

        for b_idx, base_lbl in enumerate(self._base_labels):
            roi_res = self._roi_results.get(base_lbl)
            if roi_res is None or not roi_res.snr_per_diff:
                continue
            color   = self._MULTI_COLORS[b_idx % len(self._MULTI_COLORS)]
            marker  = markers[b_idx % len(markers)]
            snr_d   = roi_res.snr_per_diff
            xs_b    = [x_pos[k] for k in snr_d]
            snr_b   = [snr_d[k].snr for k in snr_d]
            delta_b = [(snr_d[k].mu_target - snr_d[k].mu_ref) * 255.0 for k in snr_d]
            sigma_b = [snr_d[k].sigma_ref * 255.0 for k in snr_d]

            ax_snr.plot(xs_b, snr_b, marker=marker, color=color,
                        linewidth=1.8, markersize=7, label=base_lbl, zorder=3)
            ax_delta.plot(xs_b, delta_b, marker=marker, color=color,
                          linewidth=1.8, markersize=7, label=base_lbl, zorder=3)
            ax_delta.errorbar(xs_b, delta_b, yerr=sigma_b,
                              fmt='none', color=color, capsize=4,
                              linewidth=1.2, alpha=0.6, zorder=2)

        ax_snr.axhline(0, color=self._COL_SPL, linewidth=0.8, linestyle='--')
        ax_snr.set_ylabel("SNR", color=self._COL_TXT, fontsize=9)
        ax_snr.set_title("SNR = Δ / σ_Ref   [All Bases overlay]",
                         color=self._COL_MUT, fontsize=9)
        ax_snr.legend(facecolor=self._BG_FIG, labelcolor=self._COL_TXT,
                      fontsize=8, loc='upper right')

        ax_delta.axhline(0, color=self._COL_SPL, linewidth=0.8, linestyle='--')
        ax_delta.set_ylabel("Δ = T−R  (err = ±σ_Ref)", color=self._COL_TXT, fontsize=9)
        ax_delta.set_xticks(list(range(len(all_compare))))
        ax_delta.set_xticklabels(all_compare, rotation=20, ha='right',
                                 color=self._COL_TXT, fontsize=9)
        self._chart_canvas.draw()

    # ------------------------------------------------------------------
    # Tab 2b — Pair Matrix (auto-pair only)
    # ------------------------------------------------------------------

    def _build_matrix_tab(self) -> QtWidgets.QWidget:
        """N×N SNR heatmap — rows = base, columns = compare."""
        from matplotlib.colors import LinearSegmentedColormap
        import matplotlib.ticker as mticker

        w = QtWidgets.QWidget()
        w.setStyleSheet("background: white;")
        lay = QtWidgets.QVBoxLayout(w)
        lay.setContentsMargins(16, 12, 16, 8)
        lay.setSpacing(6)

        # ── Collect ordered labels ─────────────────────────────────────
        all_labels = sorted(set(
            [r.base_label    for r in self._all_results] +
            [r.compare_label for r in self._all_results]
        ))
        n = len(all_labels)
        label_idx = {lbl: i for i, lbl in enumerate(all_labels)}

        matrix = np.full((n, n), np.nan)
        for r in self._all_results:
            roi_full = self._roi_results.get(r.base_label)
            if roi_full:
                entry = roi_full.snr_per_diff.get(r.compare_label)
                if entry is not None:
                    matrix[label_idx[r.base_label], label_idx[r.compare_label]] = entry.snr

        # ── Pastel colormap: light-red → light-amber → light-green ────
        pastel_cmap = LinearSegmentedColormap.from_list(
            'snr_pastel',
            [
                (0.00, '#FECACA'),   # pastel red   (low SNR)
                (0.35, '#FDE68A'),   # pastel amber (medium)
                (0.65, '#BBF7D0'),   # pastel green (good)
                (1.00, '#6EE7B7'),   # deeper mint  (excellent)
            ]
        )
        pastel_cmap.set_bad(color='#F1F5F9')   # diagonal / missing → very light gray

        valid = matrix[~np.isnan(matrix)]
        vmin = 0.0
        vmax = max(float(np.max(valid)) if valid.size else 0.0, self._SNR_GOOD) * 1.05

        # ── Figure ────────────────────────────────────────────────────
        cell_in = 1.1                                 # inches per cell
        margin  = 2.4                                 # left/right margin for labels + colorbar
        fig_w   = n * cell_in + margin
        fig_h   = max(4.0, n * cell_in + 1.2)
        fig = Figure(figsize=(fig_w, fig_h))
        fig.patch.set_facecolor('white')

        # Leave room: left=labels, right=colorbar
        fig.subplots_adjust(left=0.18, right=0.82, top=0.88, bottom=0.18)
        ax = fig.add_subplot(111)
        ax.set_facecolor('white')

        im = ax.imshow(matrix, cmap=pastel_cmap, aspect='equal',
                       vmin=vmin, vmax=vmax, interpolation='nearest')

        # ── Grid lines between cells ───────────────────────────────────
        ax.set_xticks(np.arange(-0.5, n, 1), minor=True)
        ax.set_yticks(np.arange(-0.5, n, 1), minor=True)
        ax.grid(which='minor', color='#CBD5E1', linewidth=0.8)
        ax.tick_params(which='minor', length=0)

        # ── Colorbar ──────────────────────────────────────────────────
        cbar = fig.colorbar(im, ax=ax, fraction=0.038, pad=0.03)
        cbar.set_label("SNR", color='#374151', fontsize=11, fontweight='bold', labelpad=8)
        cbar.ax.tick_params(labelcolor='#374151', color='#CBD5E1', labelsize=10)
        cbar.outline.set_edgecolor('#CBD5E1')
        cbar.ax.set_facecolor('white')
        # Threshold lines on colorbar
        for thresh, color in [(self._SNR_OK, '#F59E0B'), (self._SNR_GOOD, '#10B981')]:
            norm_pos = thresh / max(vmax, 1e-9)
            if 0.0 < norm_pos < 1.0:
                cbar.ax.axhline(norm_pos, color=color, linewidth=1.5, linestyle='--', alpha=0.85)
                cbar.ax.text(1.25, norm_pos, f'{thresh:.0f}',
                             transform=cbar.ax.transAxes,
                             va='center', ha='left',
                             color=color, fontsize=9, fontweight='bold')

        # ── Cell annotations ──────────────────────────────────────────
        for i in range(n):
            for j in range(n):
                if i == j:
                    # Diagonal: self-compare placeholder
                    ax.add_patch(
                        __import__('matplotlib.patches', fromlist=['FancyBboxPatch'])
                        .FancyBboxPatch(
                            (j - 0.45, i - 0.45), 0.90, 0.90,
                            boxstyle='round,pad=0.05',
                            facecolor='#E2E8F0', edgecolor='none', zorder=2,
                        )
                    )
                    ax.text(j, i, '—', ha='center', va='center',
                            fontsize=12, color='#94A3B8', zorder=3)
                elif not np.isnan(matrix[i, j]):
                    val = matrix[i, j]
                    # Text always dark on the pastel background
                    if val >= self._SNR_GOOD:
                        txt_col = '#065F46'   # dark green
                    elif val >= self._SNR_OK:
                        txt_col = '#78350F'   # dark amber
                    else:
                        txt_col = '#7F1D1D'   # dark red
                    ax.text(j, i, f"{val:.2f}", ha='center', va='center',
                            fontsize=12, color=txt_col, fontweight='bold', zorder=3)
                else:
                    ax.text(j, i, 'n/a', ha='center', va='center',
                            fontsize=10, color='#94A3B8', zorder=3)

        # ── Axes labels ───────────────────────────────────────────────
        ax.set_xticks(range(n))
        ax.set_yticks(range(n))
        ax.set_xticklabels(all_labels, rotation=35, ha='right',
                           color='#1E293B', fontsize=11, fontweight='semibold')
        ax.set_yticklabels(all_labels, color='#1E293B', fontsize=11,
                           fontweight='semibold')
        ax.set_xlabel("Compare (LE)", color='#475569', fontsize=11,
                      labelpad=10, fontweight='bold')
        ax.set_ylabel("Base (LE)", color='#475569', fontsize=11,
                      labelpad=10, fontweight='bold')
        ax.set_title("SNR Pair Matrix",
                     color='#0F172A', fontsize=14, fontweight='bold', pad=14)
        ax.tick_params(which='major', length=0, pad=6)
        for spine in ax.spines.values():
            spine.set_edgecolor('#E2E8F0')
            spine.set_linewidth(1.0)

        self._matrix_fig = fig
        canvas = FigureCanvas(fig)
        canvas.setStyleSheet("background: white;")

        # Save button row
        btn_row = QtWidgets.QHBoxLayout()
        btn_row.addStretch()
        btn_row.addWidget(self._make_save_btn(lambda: self._matrix_fig, "roi_pair_matrix"))
        lay.addWidget(canvas, stretch=1)
        lay.addLayout(btn_row)
        return w

    # ------------------------------------------------------------------
    # Tab 2c — Diff Image Matrix (auto-pair only)
    # ------------------------------------------------------------------

    def _build_diff_matrix_tab(self) -> QtWidgets.QWidget:
        """N×N grid of center-cropped (128×128) diff images — rows=base, cols=compare."""
        CROP = 128  # pixels to show per cell

        def _center_crop(img: np.ndarray, size: int) -> np.ndarray:
            """Return a square centre-crop of *img* at *size*×*size* pixels."""
            h, w = img.shape[:2]
            if h < size or w < size:
                # Pad with zeros if image is smaller than crop window
                pad_h = max(0, size - h)
                pad_w = max(0, size - w)
                img = np.pad(img,
                             ((pad_h // 2, pad_h - pad_h // 2),
                              (pad_w // 2, pad_w - pad_w // 2)),
                             mode='constant', constant_values=0)
                h, w = img.shape[:2]
            cy, cx = h // 2, w // 2
            half = size // 2
            return img[cy - half: cy - half + size, cx - half: cx - half + size]

        w = QtWidgets.QWidget()
        w.setStyleSheet("background: white;")
        lay = QtWidgets.QVBoxLayout(w)
        lay.setContentsMargins(16, 12, 16, 8)
        lay.setSpacing(6)

        # ── Collect ordered labels ─────────────────────────────────────
        all_labels = sorted(set(
            [r.base_label    for r in self._all_results] +
            [r.compare_label for r in self._all_results]
        ))
        n = len(all_labels)

        # Build lookup: (base_label, compare_label) → cropped image
        pair_image: dict = {}
        for r in self._all_results:
            if r.result_image is not None:
                pair_image[(r.base_label, r.compare_label)] = \
                    _center_crop(r.result_image, CROP)

        # ── Figure geometry — same proportions as SNR Pair Matrix ─────
        cell_in = 1.1          # inches per cell
        margin  = 2.4          # left/right margin for labels
        fig_w   = n * cell_in + margin
        fig_h   = max(4.0, n * cell_in + 1.2)

        # ── Figure: light-gray background acts as subtle grid lines ───
        # Cells butt up against each other; the tiny gap reveals the figure
        # facecolor (#CBD5E1), matching the SNR Pair Matrix grid line colour.
        GRID_CLR = '#CBD5E1'   # light blue-gray → grid line colour
        fig = Figure(figsize=(fig_w, fig_h), dpi=100)
        fig.patch.set_facecolor(GRID_CLR)

        GAP = 0.012            # gap between cells (fraction of total figure)
        fig.subplots_adjust(
            left=0.18,
            right=0.95,
            top=0.88,
            bottom=0.18,
            wspace=GAP,
            hspace=GAP,
        )

        lbl_fs  = max(7, min(11, 72 // n))    # label font size
        cell_fs = max(9, min(13, 96 // n))    # diagonal "—" size

        for i, base_lbl in enumerate(all_labels):
            for j, cmp_lbl in enumerate(all_labels):
                ax = fig.add_subplot(n, n, i * n + j + 1)
                ax.set_xticks([])
                ax.set_yticks([])
                # Hide individual spines — the figure background serves as the border
                for spine in ax.spines.values():
                    spine.set_visible(False)

                if i == j:
                    # ── Diagonal: distinct light-gray fill + clear dash ─
                    ax.set_facecolor('#E2E8F0')
                    ax.text(0.5, 0.5, '—',
                            ha='center', va='center',
                            transform=ax.transAxes,
                            fontsize=cell_fs + 2, color='#94A3B8',
                            fontweight='bold')
                else:
                    img = pair_image.get((base_lbl, cmp_lbl))
                    if img is not None:
                        ax.set_facecolor('white')
                        ax.imshow(img, cmap='gray', vmin=0, vmax=255,
                                  aspect='equal', interpolation='lanczos')
                    else:
                        ax.set_facecolor('#F1F5F9')
                        ax.text(0.5, 0.5, 'n/a', ha='center', va='center',
                                transform=ax.transAxes,
                                fontsize=8, color='#94A3B8')

                # Column label — bottom row only (rotated to match SNR matrix)
                if i == n - 1:
                    ax.set_xlabel(cmp_lbl,
                                  fontsize=lbl_fs, color='#1E293B',
                                  fontweight='semibold', labelpad=5)
                    ax.xaxis.label.set_rotation(35)
                    ax.xaxis.label.set_ha('right')
                # Row label — left column only
                if j == 0:
                    ax.set_ylabel(base_lbl,
                                  fontsize=lbl_fs, color='#1E293B',
                                  fontweight='semibold',
                                  rotation=0, ha='right',
                                  va='center', labelpad=6)

        # ── Axis direction labels — styled like SNR Pair Matrix ────────
        fig.text(0.565, 0.02, "Compare (LE)",
                 ha='center', va='bottom',
                 color='#475569', fontsize=11, fontweight='bold')
        fig.text(0.02, 0.53, "Base (LE)",
                 ha='left', va='center',
                 color='#475569', fontsize=11, fontweight='bold',
                 rotation=90)

        fig.suptitle(f"Diff Image Matrix   (center {CROP}×{CROP} px crop)",
                     color='#0F172A', fontsize=14, fontweight='bold', y=0.998)

        self._diff_matrix_fig = fig
        canvas = FigureCanvas(fig)
        canvas.setMinimumSize(400, 400)
        canvas.setStyleSheet("background: white;")

        # ── Save button row ───────────────────────────────────────────
        btn_row = QtWidgets.QHBoxLayout()
        lbl_info = QtWidgets.QLabel(
            f"Showing center {CROP}×{CROP} px  |  colormap: gray")
        lbl_info.setStyleSheet(
            "color: #64748B; font-size: 11px; font-style: italic;")
        btn_row.addWidget(lbl_info)
        btn_row.addStretch()
        btn_row.addWidget(
            self._make_save_btn(lambda: self._diff_matrix_fig, "roi_diff_matrix"))

        lay.addWidget(canvas, stretch=1)
        lay.addLayout(btn_row)
        return w

    # ------------------------------------------------------------------
    # Tab 3 — Per-ROI Mean across LE
    # ------------------------------------------------------------------

    def _build_mean_tab(self) -> QtWidgets.QWidget:
        w = QtWidgets.QWidget()
        lay = QtWidgets.QVBoxLayout(w)
        lay.setSpacing(4)

        ctrl_row, self._cmb_mean_base = self._make_base_selector()
        if self._cmb_mean_base is not None:
            self._cmb_mean_base.currentTextChanged.connect(self._on_mean_base_changed)

        ctrl_row.addWidget(QtWidgets.QLabel("ROI:"))
        self._cmb_roi = QtWidgets.QComboBox()
        self._cmb_roi.setMinimumWidth(120)
        ctrl_row.addWidget(self._cmb_roi)
        ctrl_row.addStretch()
        ctrl_row.addWidget(self._make_save_btn(lambda: self._mean_canvas.figure, "roi_per_roi_mean"))
        lay.addLayout(ctrl_row)

        fig = Figure(figsize=(5, 3), tight_layout=True)
        fig.patch.set_facecolor(self._BG_FIG)
        self._mean_ax = fig.add_subplot(111)
        self._mean_canvas = FigureCanvas(fig)
        lay.addWidget(self._mean_canvas, stretch=1)

        self._cmb_roi.currentIndexChanged.connect(self._refresh_mean_plot)
        self._on_mean_base_changed()
        return w

    def _on_mean_base_changed(self) -> None:
        """Repopulate ROI combo when base changes, then redraw."""
        roi_result = self._current_roi_result(self._cmb_mean_base)
        self._cmb_roi.blockSignals(True)
        self._cmb_roi.clear()
        if roi_result is not None:
            for roi in roi_result.roi_set.rois:
                self._cmb_roi.addItem(roi.label, roi.id)
        self._cmb_roi.blockSignals(False)
        self._refresh_mean_plot()

    def _refresh_mean_plot(self) -> None:
        roi_result = self._current_roi_result(self._cmb_mean_base)
        ax = self._mean_ax
        ax.cla()
        self._style_ax(ax)
        ax.set_ylabel("Mean intensity (norm)", color=self._COL_TXT)

        roi_id = self._cmb_roi.currentData()
        if roi_result is None or roi_id is None:
            self._mean_canvas.draw()
            return

        compare_labels = roi_result.compare_labels()
        diff_labels    = roi_result.diff_labels()
        base_layer     = roi_result.get_base_layer()
        x = range(len(compare_labels))
        le_labels = [lbl.replace('_compare', '') for lbl in compare_labels]

        if base_layer and roi_id in base_layer.roi_stats:
            ax.axhline(base_layer.roi_stats[roi_id].mean,
                       color='#60A5FA', linewidth=1.5, linestyle='--', label='Base')

        comp_means = [
            roi_result.get_layer(lbl).roi_stats[roi_id].mean
            if roi_result.get_layer(lbl) and roi_id in roi_result.get_layer(lbl).roi_stats
            else None
            for lbl in compare_labels
        ]
        vx = [i for i, v in enumerate(comp_means) if v is not None]
        vy = [v for v in comp_means if v is not None]
        if vx:
            ax.plot(vx, vy, marker='s', color='#34D399', linewidth=1.5, label='Compare')

        diff_means = [
            roi_result.get_layer(lbl).roi_stats[roi_id].mean
            if roi_result.get_layer(lbl) and roi_id in roi_result.get_layer(lbl).roi_stats
            else None
            for lbl in diff_labels
        ]
        vxd = [i for i, v in enumerate(diff_means) if v is not None]
        vyd = [v for v in diff_means if v is not None]
        if vxd:
            ax.plot(vxd, vyd, marker='^', color='#F87171', linewidth=1.5, label='Diff')

        ax.set_xticks(list(x))
        ax.set_xticklabels(le_labels, rotation=20, ha='right',
                           color=self._COL_TXT, fontsize=9)
        ax.legend(facecolor=self._BG_FIG, labelcolor=self._COL_TXT, fontsize=8)
        self._mean_canvas.draw()

    # ------------------------------------------------------------------
    # Tab 4 — Raw Table (all layer × ROI stats)
    # ------------------------------------------------------------------

    def _build_table_tab(self) -> QtWidgets.QWidget:
        w = QtWidgets.QWidget()
        lay = QtWidgets.QVBoxLayout(w)
        lay.setSpacing(4)

        ctrl_row, self._cmb_raw_base = self._make_base_selector()
        if self._cmb_raw_base is not None:
            self._cmb_raw_base.currentTextChanged.connect(self._refresh_raw_table)
        lay.addLayout(ctrl_row)

        headers = ['Layer', 'Image', 'ROI', 'Type', 'Mean', 'Std', 'P2', 'P98', 'Median', 'Pixels']
        self._raw_table = QtWidgets.QTableWidget(0, len(headers))
        self._raw_table.setHorizontalHeaderLabels(headers)
        self._raw_table.horizontalHeader().setSectionResizeMode(
            QtWidgets.QHeaderView.ResizeToContents
        )
        self._raw_table.setEditTriggers(QtWidgets.QAbstractItemView.NoEditTriggers)
        self._raw_table.setAlternatingRowColors(True)
        lay.addWidget(self._raw_table, stretch=1)

        self._refresh_raw_table()
        return w

    def _refresh_raw_table(self) -> None:
        roi_result = self._current_roi_result(self._cmb_raw_base)
        self._raw_table.setRowCount(0)
        if roi_result is None:
            return
        roi_map = {r.id: r for r in roi_result.roi_set.rois}
        for layer in roi_result.layers:
            for roi_id, stats in layer.roi_stats.items():
                roi = roi_map.get(roi_id)
                row = self._raw_table.rowCount()
                self._raw_table.insertRow(row)
                values = [
                    layer.layer_type, layer.image_label,
                    roi.label if roi else roi_id,
                    roi.roi_type if roi else '',
                    f"{stats.mean:.5f}", f"{stats.std:.5f}",
                    f"{stats.p2:.5f}", f"{stats.p98:.5f}",
                    f"{stats.median:.5f}", str(stats.pixel_count),
                ]
                for col, val in enumerate(values):
                    self._raw_table.setItem(row, col, QtWidgets.QTableWidgetItem(val))


class MultiROIManagerWidget(QtWidgets.QDialog):
    """Floating dialog for managing multiple named ROIs drawn on the base image.

    Signals
    -------
    rois_changed : Emitted whenever the MultiROISet is modified (add / remove / type change).
    """

    rois_changed = Signal()

    def __init__(self, roi_set: MultiROISet, base_widget: SyncZoomImageWidget,
                 parent=None):
        super().__init__(parent)
        self._roi_set = roi_set
        self._base_widget = base_widget
        self._img_shape: Optional[Tuple[int, int]] = None  # (H, W) of base image

        self.setWindowTitle("ROI Manager")
        self.setWindowFlags(self.windowFlags() | Qt.Tool)
        self.setMinimumWidth(480)
        self.resize(480, 540)
        self.setStyleSheet(f"""
            QDialog {{
                background: {UI_BG_WINDOW};
            }}
            QGroupBox {{
                background: {UI_BG_PANEL};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.MD};
                margin-top: 14px;
                padding-top: 12px;
                font-weight: {Typography.FONT_WEIGHT_SEMIBOLD};
            }}
            QGroupBox::title {{
                subcontrol-origin: margin;
                left: 10px;
                padding: 0 6px;
                color: {UI_TEXT};
            }}
            QListWidget {{
                background: {UI_BG_PANEL};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.SM};
            }}
        """)
        self._build_ui()
        self._connect_signals()
        self._refresh_list()

    # ------------------------------------------------------------------
    # UI construction
    # ------------------------------------------------------------------

    def _build_ui(self) -> None:
        root = QtWidgets.QVBoxLayout(self)
        root.setSpacing(8)
        root.setContentsMargins(8, 8, 8, 8)

        # ── Add Mode ────────────────────────────────────────────────
        mode_grp = QtWidgets.QGroupBox("Add Mode")
        mode_layout = QtWidgets.QVBoxLayout(mode_grp)
        mode_layout.setSpacing(4)

        btn_row = QtWidgets.QHBoxLayout()
        self._btn_drag = QtWidgets.QPushButton("Drag")
        self._btn_single = QtWidgets.QPushButton("Single Add")
        self._btn_multi = QtWidgets.QPushButton("Multi Add")
        for btn in (self._btn_drag, self._btn_single, self._btn_multi):
            btn.setCheckable(True)
            btn.setMinimumHeight(30)
            btn.setStyleSheet(
                "QPushButton { font-weight: 600; }"
                "QPushButton:checked { background-color: #F59E0B; color: #111827; border: 1px solid #D97706; }"
            )
            btn_row.addWidget(btn)
        mode_layout.addLayout(btn_row)

        size_row = QtWidgets.QHBoxLayout()
        size_row.addWidget(QtWidgets.QLabel("ROI Size  W:"))
        self._spn_w = QtWidgets.QSpinBox()
        self._spn_w.setRange(4, 2048)
        self._spn_w.setValue(64)
        self._spn_w.setSuffix(" px")
        size_row.addWidget(self._spn_w)
        size_row.addWidget(QtWidgets.QLabel("H:"))
        self._spn_h = QtWidgets.QSpinBox()
        self._spn_h.setRange(4, 2048)
        self._spn_h.setValue(64)
        self._spn_h.setSuffix(" px")
        size_row.addWidget(self._spn_h)
        mode_layout.addLayout(size_row)
        root.addWidget(mode_grp)

        # ── Multi-Add settings (collapsible) ────────────────────────
        self._multi_grp = QtWidgets.QGroupBox("Multi Add Settings")
        self._multi_grp.setVisible(False)
        multi_layout = QtWidgets.QVBoxLayout(self._multi_grp)
        multi_layout.setSpacing(4)

        grid_count_row = QtWidgets.QHBoxLayout()
        grid_count_row.addWidget(QtWidgets.QLabel("Cols:"))
        self._spn_cols = QtWidgets.QSpinBox()
        self._spn_cols.setRange(1, 100)
        self._spn_cols.setValue(7)
        self._spn_cols.setToolTip("Number of ROIs in horizontal direction")
        grid_count_row.addWidget(self._spn_cols)
        grid_count_row.addWidget(QtWidgets.QLabel("Rows:"))
        self._spn_rows = QtWidgets.QSpinBox()
        self._spn_rows.setRange(1, 100)
        self._spn_rows.setValue(3)
        self._spn_rows.setToolTip("Number of ROIs in vertical direction")
        grid_count_row.addWidget(self._spn_rows)
        multi_layout.addLayout(grid_count_row)

        self._lbl_multi_hint = QtWidgets.QLabel(
            "Click Multi Add, then click image: first click = TL, second click = BR"
        )
        self._lbl_multi_hint.setWordWrap(True)
        self._lbl_multi_hint.setStyleSheet(f"color: {UI_TEXT_SECONDARY};")
        multi_layout.addWidget(self._lbl_multi_hint)

        self._lbl_tl = QtWidgets.QLabel("TL: not set")
        self._lbl_br = QtWidgets.QLabel("BR: not set")
        multi_layout.addWidget(self._lbl_tl)
        multi_layout.addWidget(self._lbl_br)

        preview_row = QtWidgets.QHBoxLayout()
        self._btn_preview_grid = QtWidgets.QPushButton("Preview Grid")
        self._lbl_grid_count = QtWidgets.QLabel("")
        preview_row.addWidget(self._btn_preview_grid)
        preview_row.addWidget(self._lbl_grid_count)
        multi_layout.addLayout(preview_row)

        confirm_row = QtWidgets.QHBoxLayout()
        self._btn_confirm_grid = QtWidgets.QPushButton("Confirm Grid")
        self._btn_reset_anchors = QtWidgets.QPushButton("Reset Anchors")
        confirm_row.addWidget(self._btn_confirm_grid)
        confirm_row.addWidget(self._btn_reset_anchors)
        multi_layout.addLayout(confirm_row)
        root.addWidget(self._multi_grp)

        # ── ROI List ─────────────────────────────────────────────────
        list_header = QtWidgets.QHBoxLayout()
        list_header.addWidget(QtWidgets.QLabel("ROI List"))
        list_header.addStretch()
        self._btn_clear_all = QtWidgets.QPushButton("Clear All")
        self._btn_clear_all.setFixedWidth(80)
        list_header.addWidget(self._btn_clear_all)
        root.addLayout(list_header)

        self._list_widget = QtWidgets.QListWidget()
        self._list_widget.setAlternatingRowColors(True)
        self._list_widget.setSelectionMode(QtWidgets.QAbstractItemView.SingleSelection)
        self._list_widget.setSpacing(2)
        root.addWidget(self._list_widget, stretch=1)

        self._lbl_summary = QtWidgets.QLabel("Target: 0   Reference: 0")
        root.addWidget(self._lbl_summary)

        # ── Bottom buttons ────────────────────────────────────────────
        bottom_row = QtWidgets.QHBoxLayout()
        self._btn_confirm = QtWidgets.QPushButton("Confirm")
        self._btn_confirm.setMinimumHeight(34)
        self._btn_confirm.setStyleSheet(
            f"background: {UI_PRIMARY}; color: {UI_TEXT_ON_PRIMARY}; font-weight: 700;"
        )
        bottom_row.addStretch(1)
        bottom_row.addWidget(self._btn_confirm)
        root.addLayout(bottom_row)

    # ------------------------------------------------------------------
    # Signal connections
    # ------------------------------------------------------------------

    def _connect_signals(self) -> None:
        self._btn_drag.clicked.connect(lambda: self._activate_mode('drag'))
        self._btn_single.clicked.connect(lambda: self._activate_mode('single_add'))
        self._btn_multi.clicked.connect(lambda: self._toggle_multi_mode())
        self._btn_preview_grid.clicked.connect(self._on_preview_grid)
        self._btn_confirm_grid.clicked.connect(self._on_confirm_grid)
        self._btn_reset_anchors.clicked.connect(self._on_reset_anchors)
        self._btn_clear_all.clicked.connect(self._on_clear_all)
        self._btn_confirm.clicked.connect(self._on_confirm)
        # Real-time grid preview when spinbox values change
        self._spn_cols.valueChanged.connect(self._auto_preview_grid)
        self._spn_rows.valueChanged.connect(self._auto_preview_grid)
        self._spn_w.valueChanged.connect(self._auto_preview_grid)
        self._spn_h.valueChanged.connect(self._auto_preview_grid)
        # Signals from base widget
        self._base_widget.multi_roi_drawn.connect(self._on_roi_drawn)
        self._base_widget.multi_roi_anchor.connect(self._on_anchor_set)

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def set_image_shape(self, shape: Tuple[int, int]) -> None:
        """Inform the manager of the base image shape (H, W) for pixel→norm conversion."""
        self._img_shape = shape

    # ------------------------------------------------------------------
    # Mode management
    # ------------------------------------------------------------------

    def _activate_mode(self, mode: str) -> None:
        """Activate drag or single_add mode on the base widget."""
        self._btn_drag.setChecked(mode == 'drag')
        self._btn_single.setChecked(mode == 'single_add')
        self._btn_multi.setChecked(False)
        self._multi_grp.setVisible(False)
        add_size = self._get_add_size_norm()
        self._base_widget.set_multi_draw_mode(mode, add_size_norm=add_size)

    def _toggle_multi_mode(self) -> None:
        visible = self._multi_grp.isVisible()
        self._multi_grp.setVisible(not visible)
        self._btn_drag.setChecked(False)
        self._btn_single.setChecked(False)
        self._btn_multi.setChecked(not visible)
        if not visible:
            self._on_reset_anchors()
            self._base_widget.set_multi_draw_mode('multi_add')
            self._lbl_multi_hint.setText(
                "Click image twice: first point = TL, second point = BR"
            )
        else:
            self._base_widget.set_multi_draw_mode('idle')

    def _get_add_size_norm(self) -> Optional[Tuple[float, float]]:
        if self._img_shape is None:
            return None
        h, w = self._img_shape
        return self._spn_w.value() / max(w, 1), self._spn_h.value() / max(h, 1)

    # ------------------------------------------------------------------
    # ROI draw handlers (from widget signals)
    # ------------------------------------------------------------------

    def _on_roi_drawn(self, nx: float, ny: float, nw: float, nh: float) -> None:
        self._roi_set.add_roi((nx, ny, nw, nh))
        self._refresh_list()
        self._base_widget._update_display()
        self.rois_changed.emit()

    def _on_anchor_set(self, which: str, cx: float, cy: float) -> None:
        if which == 'tl':
            self._lbl_tl.setText(f"TL: ({cx:.3f}, {cy:.3f})")
            self._lbl_multi_hint.setText("TL set. Click image again to set BR.")
        else:
            self._lbl_br.setText(f"BR: ({cx:.3f}, {cy:.3f})")
            self._lbl_multi_hint.setText("TL/BR set. You can Preview Grid or Confirm Grid.")
        self._base_widget.set_grid_anchors(
            self._base_widget._grid_anchor_tl,
            self._base_widget._grid_anchor_br,
        )
        self._auto_preview_grid()

    # ------------------------------------------------------------------
    # Grid management
    # ------------------------------------------------------------------

    def _on_preview_grid(self) -> None:
        rects = self._build_grid_rects()
        self._base_widget.set_grid_preview(rects)
        c, r = self._spn_cols.value(), self._spn_rows.value()
        self._lbl_grid_count.setText(f"{c}×{r} = {len(rects)} ROIs")

    def _auto_preview_grid(self) -> None:
        """Update grid preview automatically when spinbox values change (if anchors are set)."""
        if self._base_widget._grid_anchor_tl is not None and self._base_widget._grid_anchor_br is not None:
            self._on_preview_grid()

    def _on_confirm_grid(self) -> None:
        rects = self._build_grid_rects()
        for rect in rects:
            self._roi_set.add_roi(rect)
        self._on_reset_anchors()
        self._refresh_list()
        self._base_widget._update_display()
        self.rois_changed.emit()

    def _on_reset_anchors(self) -> None:
        self._base_widget._grid_anchor_tl = None
        self._base_widget._grid_anchor_br = None
        self._base_widget.clear_grid_preview()
        self._lbl_tl.setText("TL: not set")
        self._lbl_br.setText("BR: not set")
        self._lbl_grid_count.setText("")

    def _build_grid_rects(self) -> List[Tuple]:
        if self._img_shape is None:
            return []
        tl = self._base_widget._grid_anchor_tl
        br = self._base_widget._grid_anchor_br
        if tl is None or br is None:
            return []
        return self._roi_set.generate_grid(
            anchor_tl_norm=tl,
            anchor_br_norm=br,
            cols=self._spn_cols.value(),
            rows=self._spn_rows.value(),
            roi_w_px=self._spn_w.value(),
            roi_h_px=self._spn_h.value(),
            img_shape=self._img_shape,
        )

    # ------------------------------------------------------------------
    # ROI list management
    # ------------------------------------------------------------------

    def _on_clear_all(self) -> None:
        self._roi_set.clear()
        self._base_widget.clear_grid_preview()
        self._base_widget.set_grid_anchors(None, None)
        self._base_widget._update_display()
        self._refresh_list()
        self.rois_changed.emit()

    def _refresh_list(self) -> None:
        self._list_widget.clear()
        for roi in self._roi_set.rois:
            item = QtWidgets.QListWidgetItem()
            self._list_widget.addItem(item)
            row_widget = self._make_roi_row(roi)
            item.setSizeHint(row_widget.sizeHint())
            self._list_widget.setItemWidget(item, row_widget)
        n_target = 1 if self._roi_set.get_target() else 0
        n_ref = len(self._roi_set.get_references())
        self._lbl_summary.setText(f"Target: {n_target} ★   Reference: {n_ref}")

    def _make_roi_row(self, roi: NamedROI) -> QtWidgets.QWidget:
        w = QtWidgets.QWidget()
        row = QtWidgets.QHBoxLayout(w)
        row.setContentsMargins(6, 4, 6, 4)
        row.setSpacing(6)

        # Color swatch
        swatch = QtWidgets.QLabel()
        r, g, b = roi.color_bgr[2], roi.color_bgr[1], roi.color_bgr[0]
        swatch.setFixedSize(14, 14)
        swatch.setStyleSheet(f"background: rgb({r},{g},{b}); border-radius: 3px;")
        row.addWidget(swatch)

        # Label
        lbl = QtWidgets.QLabel(roi.label)
        lbl.setMinimumWidth(82)
        lbl.setMinimumHeight(24)
        row.addWidget(lbl, stretch=1)

        # Current state badge
        state_badge = QtWidgets.QLabel("T" if roi.roi_type == 'target' else "R")
        state_badge.setFixedSize(20, 24)
        state_badge.setAlignment(Qt.AlignCenter)
        state_badge.setStyleSheet(
            "font-weight: bold; color: #ff4444;" if roi.roi_type == 'target'
            else "font-weight: bold; color: #00e5ff;"
        )
        row.addWidget(state_badge)

        # Target icon button (always sets this ROI as Target)
        btn_target = QtWidgets.QToolButton()
        btn_target.setText("★" if roi.roi_type == 'target' else "☆")
        btn_target.setToolTip("Set as Target")
        btn_target.setFixedSize(28, 24)
        btn_target.setStyleSheet(
            "QToolButton { border: 1px solid #E5E7EB; border-radius: 6px; font-size: 14px; }"
            "QToolButton:hover { border-color: #F59E0B; background: #FFF8ED; }"
        )
        btn_target.clicked.connect(lambda checked=False, rid=roi.id: self._on_promote(rid))
        row.addWidget(btn_target)

        # Delete icon button
        btn_del = QtWidgets.QToolButton()
        btn_del.setText("🗑")
        btn_del.setToolTip("Delete ROI")
        btn_del.setMinimumSize(30, 24)
        btn_del.clicked.connect(lambda checked=False, rid=roi.id: self._on_delete(rid))
        row.addWidget(btn_del)
        return w

    def _on_promote(self, roi_id: str) -> None:
        self._roi_set.set_target(roi_id)
        self._refresh_list()
        self._base_widget._update_display()
        self.rois_changed.emit()

    def _on_delete(self, roi_id: str) -> None:
        self._roi_set.remove_roi(roi_id)
        self._refresh_list()
        self._base_widget._update_display()
        self.rois_changed.emit()

    # ------------------------------------------------------------------
    # Dialog controls
    # ------------------------------------------------------------------

    def _clear_multi_add_markers(self) -> None:
        self._base_widget.set_multi_draw_mode('idle')
        self._on_reset_anchors()
        self._base_widget.clear_multi_draw_preview()

    def _on_confirm(self) -> None:
        self._clear_multi_add_markers()
        self.accept()

    def closeEvent(self, event: QtGui.QCloseEvent) -> None:
        self._clear_multi_add_markers()
        super().closeEvent(event)


def _ppt_add_roi_profile_slides(prs, roi_full_results,
                                _fill_bg, _add_text,
                                C_BG, C_CARD, C_TEXT, C_TEXT_SEC, C_PRIMARY,
                                Inches, Pt):
    """Per-ROI Intensity Profile slides — one slide per base label.

    Each slide shows a subplot grid (one subplot per ROI) with three lines:
    Base (dashed blue), Compare (green), Diff (red triangles).
    """
    from io import BytesIO
    from matplotlib.figure import Figure as _MplFig
    from matplotlib.backends.backend_agg import FigureCanvasAgg as _MplAgg

    BG_FIG = '#1F2937'
    BG_AX  = '#111827'
    COL_TXT = '#D1D5DB'
    COL_GRD = '#374151'
    COL_SPL = '#4B5563'

    def _style(ax):
        ax.set_facecolor(BG_AX)
        ax.tick_params(colors=COL_TXT, labelsize=7)
        for side in ('bottom', 'left'):
            ax.spines[side].set_color(COL_SPL)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(True, color=COL_GRD, linewidth=0.4, linestyle='--')

    for base_lbl, roi_full in roi_full_results.items():
        rois = roi_full.roi_set.rois
        if not rois:
            continue

        compare_labels = roi_full.compare_labels()
        diff_labels    = [l.image_label for l in roi_full.layers if l.layer_type == 'diff']
        base_layer     = roi_full.get_base_layer()
        le_labels      = [lbl.replace('_compare', '') for lbl in compare_labels]
        xs             = list(range(len(compare_labels)))

        # ── Compute base SNR (same formula as SNR Pair Matrix, on raw base) ──
        import numpy as _np
        base_snr: Optional[float] = None
        target_roi_b  = roi_full.roi_set.get_target()
        ref_rois_b    = roi_full.roi_set.get_references()
        if target_roi_b and ref_rois_b and base_layer:
            t_stat = base_layer.roi_stats.get(target_roi_b.id)
            ref_stats_b = [base_layer.roi_stats[r.id]
                           for r in ref_rois_b if r.id in base_layer.roi_stats]
            if t_stat and ref_stats_b:
                ref_means_b = _np.array([rs.mean for rs in ref_stats_b], dtype=_np.float32)
                mu_t_b  = float(t_stat.mean)
                mu_r_b  = float(_np.mean(ref_means_b))
                if len(ref_means_b) >= 2:
                    sigma_b = float(_np.std(ref_means_b))
                else:
                    sigma_b = float(ref_stats_b[0].std)
                base_snr = max(0.0, (mu_t_b - mu_r_b) / (sigma_b + 1e-7)) if sigma_b > 1e-7 else 0.0

        n_rois = len(rois)
        NCOLS  = min(n_rois, 3)
        NROWS  = (n_rois + NCOLS - 1) // NCOLS
        ITEMS_PER_PAGE = NCOLS * NROWS  # all on one slide for now

        # chunk by ITEMS_PER_PAGE (in case many ROIs)
        for page_start in range(0, n_rois, ITEMS_PER_PAGE):
            chunk = rois[page_start: page_start + ITEMS_PER_PAGE]
            nrows_this = (len(chunk) + NCOLS - 1) // NCOLS

            fig = _MplFig(figsize=(12, max(3.0, nrows_this * 2.8)), dpi=120)
            fig.patch.set_facecolor(BG_FIG)
            fig.subplots_adjust(
                left=0.07, right=0.97, top=0.97, bottom=0.12,
                wspace=0.32, hspace=0.48,
            )

            axes = []
            for ri in range(len(chunk)):
                axes.append(fig.add_subplot(nrows_this, NCOLS, ri + 1))

            for ax, roi in zip(axes, chunk):
                _style(ax)
                ax.set_title(roi.label, color=COL_TXT, fontsize=8, fontweight='bold', pad=3)

                # Base — horizontal dashed line
                if base_layer and roi.id in base_layer.roi_stats:
                    bv = base_layer.roi_stats[roi.id].mean
                    ax.axhline(bv, color='#60A5FA', linewidth=1.4,
                               linestyle='--', label='Base')

                # Compare
                c_means = []
                for lbl in compare_labels:
                    lay = roi_full.get_layer(lbl)
                    c_means.append(
                        lay.roi_stats[roi.id].mean
                        if lay and roi.id in lay.roi_stats else None
                    )
                vx = [i for i, v in enumerate(c_means) if v is not None]
                vy = [v for v in c_means if v is not None]
                if vx:
                    ax.plot(vx, vy, marker='s', markersize=4,
                            color='#34D399', linewidth=1.3, label='Compare')

                # Diff
                d_means = []
                for lbl in diff_labels:
                    lay = roi_full.get_layer(lbl)
                    d_means.append(
                        lay.roi_stats[roi.id].mean
                        if lay and roi.id in lay.roi_stats else None
                    )
                vxd = [i for i, v in enumerate(d_means) if v is not None]
                vyd = [v for v in d_means if v is not None]
                if vxd:
                    ax.plot(vxd, vyd, marker='^', markersize=4,
                            color='#F87171', linewidth=1.3, label='Diff')

                ax.set_xticks(xs)
                ax.set_xticklabels(le_labels, rotation=20, ha='right',
                                   color=COL_TXT, fontsize=7)
                # ROI type badge colour on y-axis label
                type_col = '#F87171' if roi.roi_type == 'target' else '#60A5FA'
                ax.set_ylabel("Mean (norm)", color=type_col, fontsize=7, labelpad=3)

                # Annotate base SNR on target ROI subplot
                if roi.roi_type == 'target' and base_snr is not None:
                    snr_txt = f"Base SNR: {base_snr:.3f}"
                    ax.annotate(snr_txt,
                                xy=(0.98, 0.97), xycoords='axes fraction',
                                ha='right', va='top', fontsize=7,
                                color='#FCD34D',
                                bbox=dict(boxstyle='round,pad=0.25',
                                          facecolor='#1F2937', edgecolor='#FCD34D',
                                          linewidth=0.8, alpha=0.85))

            # Shared legend in first subplot
            if axes:
                axes[0].legend(facecolor=BG_FIG, labelcolor=COL_TXT,
                               fontsize=7, loc='best')

            # Hide unused subplot slots
            for ax in axes[len(chunk):]:
                ax.set_visible(False)

            n_pages = (n_rois + ITEMS_PER_PAGE - 1) // ITEMS_PER_PAGE
            page_n  = page_start // ITEMS_PER_PAGE + 1
            suffix  = f"  ({page_n}/{n_pages})" if n_pages > 1 else ""
            # No suptitle — PPT slide title already covers it

            buf = BytesIO()
            _MplAgg(fig).print_figure(buf, format='png', dpi=120, facecolor=BG_FIG)
            buf.seek(0)

            sl = prs.slides.add_slide(prs.slide_layouts[6])
            _fill_bg(sl, C_BG)
            snr_subtitle = (f"  |  Base SNR = {base_snr:.3f}" if base_snr is not None else "")
            _add_text(sl, f"Per-ROI Intensity Profile — Base: {base_lbl}{suffix}",
                      Inches(0.3), Inches(0.08), Inches(9.5), Inches(0.44),
                      size=15, bold=True, color=C_PRIMARY)
            if snr_subtitle:
                _add_text(sl, f"Base SNR = {base_snr:.3f}",
                          Inches(9.9), Inches(0.10), Inches(3.2), Inches(0.40),
                          size=13, bold=True, color=C_TEXT_SEC)
            # Place chart immediately below title (0.54") to minimise gap
            sl.shapes.add_picture(buf, Inches(0.2), Inches(0.54), width=Inches(12.9))


def _ppt_add_roi_position_slides(prs, roi_full_results, images_dict,
                                  _fill_bg, _add_text,
                                  C_BG, C_CARD, C_TEXT, C_TEXT_SEC, C_PRIMARY,
                                  Inches, Pt):
    """ROI Position slides — base image with target/reference ROI boxes overlaid.

    One slide per base label. Each slide shows the full base image with coloured
    ROI rectangles and labels drawn on it, plus a small legend.
    """
    from io import BytesIO
    try:
        from pptx.dml.color import RGBColor
    except Exception:
        return

    FONT     = cv2.FONT_HERSHEY_SIMPLEX
    TGT_COL  = (60,  80, 220)   # BGR red-ish  (target)
    REF_COL  = (200, 200,  40)  # BGR cyan-ish (reference)

    for base_lbl, roi_full in roi_full_results.items():
        rois = roi_full.roi_set.rois
        if not rois:
            continue

        # Load base image (prefer in-memory, fall back gracefully)
        base_img = images_dict.get(base_lbl)
        if base_img is None:
            continue

        # Convert to BGR colour so ROI colours are visible
        if len(base_img.shape) == 2:
            vis = cv2.cvtColor(base_img, cv2.COLOR_GRAY2BGR)
        else:
            vis = base_img.copy()

        h, w = vis.shape[:2]

        # Draw each ROI
        for roi in rois:
            nx, ny, nw, nh = roi.norm_rect
            rx  = int(nx * w)
            ry  = int(ny * h)
            rpw = max(1, int(nw * w))
            rph = max(1, int(nh * h))
            col       = TGT_COL if roi.roi_type == 'target' else REF_COL
            thickness = 3       if roi.roi_type == 'target' else 2
            # Outer rectangle
            cv2.rectangle(vis, (rx, ry), (rx + rpw, ry + rph), col, thickness)
            # Inner shadow for contrast on bright backgrounds
            cv2.rectangle(vis, (rx + 1, ry + 1),
                          (rx + rpw - 1, ry + rph - 1), (0, 0, 0), 1)
            # Short label: "ROI_001" → "1", "ROI_3" → "3"
            short_lbl = roi.label.split('_')[-1].lstrip('0') or '0'
            font_scale = max(0.35, min(0.65, w / 1200))
            tx = rx + 3
            ty = max(ry - 4, 10)
            cv2.putText(vis, short_lbl, (tx + 1, ty + 1), FONT,
                        font_scale, (0, 0, 0), 2, cv2.LINE_AA)
            cv2.putText(vis, short_lbl, (tx, ty), FONT,
                        font_scale, col, 1, cv2.LINE_AA)

        # Encode to PNG via PIL
        try:
            from PIL import Image as _PilImg
            pil = _PilImg.fromarray(cv2.cvtColor(vis, cv2.COLOR_BGR2RGB))
            buf = BytesIO()
            pil.save(buf, format='PNG')
            buf.seek(0)
        except ImportError:
            ok, enc = cv2.imencode('.png', vis)
            if not ok:
                continue
            buf = BytesIO(enc.tobytes())

        # ── Slide layout ─────────────────────────────────────────────
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        _fill_bg(sl, C_BG)

        _add_text(sl, f"ROI Position Map — Base: {base_lbl}",
                  Inches(0.3), Inches(0.08), Inches(10.0), Inches(0.46),
                  size=15, bold=True, color=C_PRIMARY)

        # Image — aspect-ratio fit within available area (10.3" × 6.78")
        AVAIL_W = Inches(10.3)
        AVAIL_H = Inches(6.78)   # 7.5 - 0.62 title - 0.10 margin
        nat_h, nat_w = vis.shape[:2]
        if nat_w == 0 or nat_h == 0:
            disp_w, disp_h = AVAIL_W, AVAIL_H
        elif (nat_w / nat_h) >= (AVAIL_W / AVAIL_H):
            disp_w = AVAIL_W
            disp_h = int(AVAIL_W * nat_h / nat_w)
        else:
            disp_h = AVAIL_H
            disp_w = int(AVAIL_H * nat_w / nat_h)
        img_y = Inches(0.62) + (AVAIL_H - disp_h) // 2
        sl.shapes.add_picture(buf, Inches(0.2), img_y, width=disp_w, height=disp_h)

        # ── Legend (right side) ───────────────────────────────────────
        LEG_X = Inches(11.0)
        LEG_Y = Inches(1.5)

        def _legend_row(y_off, color_rgb, label, thick_str):
            box = sl.shapes.add_shape(
                1, LEG_X, LEG_Y + y_off, Inches(0.28), Inches(0.22))
            box.fill.solid()
            box.fill.fore_color.rgb = color_rgb
            box.line.fill.background()
            _add_text(sl, f"{label}  ({thick_str})",
                      LEG_X + Inches(0.34), LEG_Y + y_off,
                      Inches(2.0), Inches(0.28),
                      size=9, color=C_TEXT)

        _add_text(sl, "Legend", LEG_X, LEG_Y - Inches(0.32),
                  Inches(2.2), Inches(0.30),
                  size=10, bold=True, color=C_TEXT_SEC)
        _legend_row(Inches(0.00),
                    RGBColor(220, 80, 60),   "Target",    "thick border")
        _legend_row(Inches(0.35),
                    RGBColor(40, 200, 200),  "Reference", "thin border")

        # ROI count summary
        n_tgt = sum(1 for r in rois if r.roi_type == 'target')
        n_ref = sum(1 for r in rois if r.roi_type == 'reference')
        _add_text(sl,
                  f"Total: {len(rois)} ROIs\n"
                  f"  Target: {n_tgt}\n"
                  f"  Reference: {n_ref}",
                  LEG_X, LEG_Y + Inches(0.8),
                  Inches(2.2), Inches(0.85),
                  size=9, color=C_TEXT_SEC)


def _ppt_add_matrix_slide(prs, roi_full_results, roi_all_results,
                          _fill_bg, _add_text,
                          C_BG, C_CARD, C_TEXT, C_TEXT_SEC, C_PRIMARY,
                          Inches, Pt):
    """One slide: SNR Pair Matrix (left) | Diff Image Matrix (right), each with own title."""
    from io import BytesIO
    import numpy as np
    from matplotlib.figure import Figure as _MplFig
    from matplotlib.backends.backend_agg import FigureCanvasAgg as _MplAgg
    from matplotlib.colors import LinearSegmentedColormap
    try:
        from pptx.dml.color import RGBColor
    except Exception:
        return

    SNR_OK, SNR_GOOD, CROP = 1.0, 2.0, 128

    def _ccrop(img, size):
        h, w = img.shape[:2]
        if h < size or w < size:
            ph, pw = max(0, size - h), max(0, size - w)
            img = np.pad(img, ((ph // 2, ph - ph // 2), (pw // 2, pw - pw // 2)),
                         mode='constant', constant_values=0)
            h, w = img.shape[:2]
        cy, cx = h // 2, w // 2
        half = size // 2
        return img[cy - half: cy - half + size, cx - half: cx - half + size]

    all_labels = sorted(set(
        [r.base_label for r in roi_all_results] +
        [r.compare_label for r in roi_all_results]
    ))
    n = len(all_labels)
    if n == 0:
        return
    label_idx = {lbl: i for i, lbl in enumerate(all_labels)}

    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(sl, C_BG)

    # Vertical divider
    div = sl.shapes.add_shape(1, Inches(6.63), Inches(0.05), Inches(0.03), Inches(7.4))
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor(0x33, 0x41, 0x55)
    div.line.fill.background()

    # ── Left: SNR Pair Matrix ─────────────────────────────────────────
    _add_text(sl, "SNR Pair Matrix",
              Inches(0.2), Inches(0.05), Inches(6.3), Inches(0.48),
              size=16, bold=True, color=C_PRIMARY)

    matrix = np.full((n, n), np.nan)
    for r in roi_all_results:
        roi_full = roi_full_results.get(r.base_label)
        if roi_full:
            entry = roi_full.snr_per_diff.get(r.compare_label)
            if entry is not None:
                matrix[label_idx[r.base_label], label_idx[r.compare_label]] = entry.snr

    pastel_cmap = LinearSegmentedColormap.from_list(
        'snr_pastel',
        [(0.00, '#FECACA'), (0.35, '#FDE68A'), (0.65, '#BBF7D0'), (1.00, '#6EE7B7')]
    )
    pastel_cmap.set_bad(color='#F1F5F9')
    valid = matrix[~np.isnan(matrix)]
    vmax = max(float(np.max(valid)) if valid.size else 0.0, SNR_GOOD) * 1.05

    cell_in = 0.9
    fw = n * cell_in + 2.0
    fh = max(3.8, n * cell_in + 1.0)
    fig_snr = _MplFig(figsize=(fw, fh), dpi=130)
    fig_snr.patch.set_facecolor('white')
    fig_snr.subplots_adjust(left=0.20, right=0.82, top=0.92, bottom=0.20)
    ax = fig_snr.add_subplot(111)
    ax.set_facecolor('white')
    im = ax.imshow(matrix, cmap=pastel_cmap, aspect='equal',
                   vmin=0.0, vmax=vmax, interpolation='nearest')
    ax.set_xticks(np.arange(-0.5, n, 1), minor=True)
    ax.set_yticks(np.arange(-0.5, n, 1), minor=True)
    ax.grid(which='minor', color='#CBD5E1', linewidth=0.8)
    ax.tick_params(which='minor', length=0)
    cbar = fig_snr.colorbar(im, ax=ax, fraction=0.038, pad=0.03)
    cbar.set_label("SNR", color='#374151', fontsize=9, fontweight='bold', labelpad=6)
    cbar.ax.tick_params(labelcolor='#374151', labelsize=8)
    for thresh, color in [(SNR_OK, '#F59E0B'), (SNR_GOOD, '#10B981')]:
        npos = thresh / max(vmax, 1e-9)
        if 0.0 < npos < 1.0:
            cbar.ax.axhline(npos, color=color, linewidth=1.5, linestyle='--', alpha=0.85)
    lbl_fs = max(7, min(10, 72 // n))
    for i in range(n):
        for j in range(n):
            if i == j:
                ax.text(j, i, '—', ha='center', va='center',
                        fontsize=9, color='#94A3B8', zorder=3)
            elif not np.isnan(matrix[i, j]):
                v = matrix[i, j]
                c = '#065F46' if v >= SNR_GOOD else ('#78350F' if v >= SNR_OK else '#7F1D1D')
                ax.text(j, i, f"{v:.2f}", ha='center', va='center',
                        fontsize=9, color=c, fontweight='bold', zorder=3)
            else:
                ax.text(j, i, 'n/a', ha='center', va='center',
                        fontsize=7, color='#94A3B8', zorder=3)
    ax.set_xticks(range(n))
    ax.set_yticks(range(n))
    ax.set_xticklabels(all_labels, rotation=35, ha='right',
                       color='#1E293B', fontsize=lbl_fs)
    ax.set_yticklabels(all_labels, color='#1E293B', fontsize=lbl_fs)
    ax.set_xlabel("Compare (LE)", color='#475569', fontsize=9,
                  labelpad=6, fontweight='bold')
    ax.set_ylabel("Base (LE)", color='#475569', fontsize=9,
                  labelpad=6, fontweight='bold')
    ax.tick_params(which='major', length=0, pad=4)
    for spine in ax.spines.values():
        spine.set_edgecolor('#E2E8F0')
    buf_snr = BytesIO()
    _MplAgg(fig_snr).print_figure(buf_snr, format='png', dpi=130, facecolor='white')
    buf_snr.seek(0)
    sl.shapes.add_picture(buf_snr, Inches(0.2), Inches(0.6), width=Inches(6.3))

    # ── Right: Diff Image Matrix ──────────────────────────────────────
    _add_text(sl, "Diff Image Matrix",
              Inches(6.8), Inches(0.05), Inches(6.35), Inches(0.48),
              size=16, bold=True, color=C_PRIMARY)

    pair_image = {}
    for r in roi_all_results:
        if r.result_image is not None:
            pair_image[(r.base_label, r.compare_label)] = _ccrop(r.result_image, CROP)

    # Light-gray background acts as subtle grid lines, matching SNR Pair Matrix style
    GRID_CLR_DIFF = '#CBD5E1'
    GAP = 0.012
    cell_in2 = 0.9
    fw2 = n * cell_in2 + 1.8
    fh2 = max(3.8, n * cell_in2 + 1.0)
    fig_diff = _MplFig(figsize=(fw2, fh2), dpi=130)
    fig_diff.patch.set_facecolor(GRID_CLR_DIFF)
    fig_diff.subplots_adjust(
        left=0.18, right=0.95,
        top=0.88, bottom=0.20,
        wspace=GAP, hspace=GAP,
    )
    d_lbl_fs = max(6, min(9, 70 // n))
    d_cell_fs = max(8, min(12, 90 // n))
    for i, base_lbl in enumerate(all_labels):
        for j, cmp_lbl in enumerate(all_labels):
            ax2 = fig_diff.add_subplot(n, n, i * n + j + 1)
            ax2.set_xticks([])
            ax2.set_yticks([])
            for sp in ax2.spines.values():
                sp.set_visible(False)
            if i == j:
                ax2.set_facecolor('#E2E8F0')
                ax2.text(0.5, 0.5, '—', ha='center', va='center',
                         transform=ax2.transAxes, fontsize=d_cell_fs,
                         color='#94A3B8', fontweight='bold')
            else:
                img2 = pair_image.get((base_lbl, cmp_lbl))
                if img2 is not None:
                    ax2.set_facecolor('white')
                    ax2.imshow(img2, cmap='gray', vmin=0, vmax=255,
                               aspect='equal', interpolation='lanczos')
                else:
                    ax2.set_facecolor('#F1F5F9')
                    ax2.text(0.5, 0.5, 'n/a', ha='center', va='center',
                             transform=ax2.transAxes, fontsize=7, color='#94A3B8')
            if i == n - 1:
                ax2.set_xlabel(cmp_lbl, fontsize=d_lbl_fs, color='#1E293B',
                               fontweight='semibold', labelpad=3)
                ax2.xaxis.label.set_rotation(35)
                ax2.xaxis.label.set_ha('right')
            if j == 0:
                ax2.set_ylabel(base_lbl, fontsize=d_lbl_fs, color='#1E293B',
                               fontweight='semibold', rotation=0,
                               ha='right', va='center', labelpad=5)
    fig_diff.text(0.565, 0.02, "Compare (LE)",
                  ha='center', va='bottom',
                  color='#475569', fontsize=8, fontweight='bold')
    fig_diff.text(0.02, 0.54, "Base (LE)",
                  ha='left', va='center',
                  color='#475569', fontsize=8, fontweight='bold',
                  rotation=90)
    buf_diff = BytesIO()
    _MplAgg(fig_diff).print_figure(buf_diff, format='png', dpi=130, facecolor=GRID_CLR_DIFF)
    buf_diff.seek(0)
    sl.shapes.add_picture(buf_diff, Inches(6.8), Inches(0.6), width=Inches(6.35))


def _ppt_add_base_snr_gallery_slide(prs, roi_full_results,
                                     _fill_bg, _add_text,
                                     C_BG, C_CARD, C_TEXT, C_TEXT_SEC, C_PRIMARY,
                                     Inches, Pt):
    """Gallery slide showing Base SNR for each base image.

    Computes SNR from the raw base image layer using the same formula as the
    SNR Pair Matrix:  SNR = (μ_target − μ_ref) / σ_ref
    where μ_ref / σ_ref come from the reference ROIs in the base layer.
    """
    from io import BytesIO
    import numpy as _np
    from matplotlib.figure import Figure as _MplFig
    from matplotlib.backends.backend_agg import FigureCanvasAgg as _MplAgg

    BG_FIG = '#1F2937'
    BG_AX  = '#111827'
    COL_TXT = '#D1D5DB'
    COL_GRD = '#374151'
    COL_SPL = '#4B5563'
    _EPS = 1e-7

    # Build base SNR table: {base_lbl: snr_value | None}
    base_snr_map = {}
    for base_lbl, roi_full in roi_full_results.items():
        base_layer   = roi_full.get_base_layer()
        target_roi   = roi_full.roi_set.get_target()
        ref_rois     = roi_full.roi_set.get_references()
        snr_val      = None
        if target_roi and ref_rois and base_layer:
            t_stat = base_layer.roi_stats.get(target_roi.id)
            ref_stats = [base_layer.roi_stats[r.id]
                         for r in ref_rois if r.id in base_layer.roi_stats]
            if t_stat and ref_stats:
                ref_means = _np.array([rs.mean for rs in ref_stats], dtype=_np.float32)
                mu_t  = float(t_stat.mean)
                mu_r  = float(_np.mean(ref_means))
                sigma = float(_np.std(ref_means)) if len(ref_means) >= 2 else float(ref_stats[0].std)
                snr_val = max(0.0, (mu_t - mu_r) / sigma) if sigma > _EPS else 0.0
        base_snr_map[base_lbl] = snr_val

    if not base_snr_map:
        return

    labels  = list(base_snr_map.keys())
    values  = [base_snr_map[l] for l in labels]
    n       = len(labels)

    # Bar colors: green (≥2), amber (≥1), red (<1)
    SNR_OK, SNR_GOOD = 1.0, 2.0
    bar_colors = []
    for v in values:
        if v is None:
            bar_colors.append('#6B7280')
        elif v >= SNR_GOOD:
            bar_colors.append('#34D399')
        elif v >= SNR_OK:
            bar_colors.append('#FBBF24')
        else:
            bar_colors.append('#F87171')

    fw = max(5.0, n * 1.2 + 2.0)
    fig = _MplFig(figsize=(fw, 4.0), dpi=130)
    fig.patch.set_facecolor(BG_FIG)
    fig.subplots_adjust(left=0.14, right=0.96, top=0.88, bottom=0.22)
    ax = fig.add_subplot(111)
    ax.set_facecolor(BG_AX)
    ax.tick_params(colors=COL_TXT, labelsize=8)
    for side in ('bottom', 'left'):
        ax.spines[side].set_color(COL_SPL)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(True, axis='y', color=COL_GRD, linewidth=0.5, linestyle='--')

    xs_bar = list(range(n))
    bars = ax.bar(xs_bar, [v if v is not None else 0.0 for v in values],
                  color=bar_colors, width=0.6, zorder=3)

    # Value labels above bars
    for xi, v in zip(xs_bar, values):
        if v is not None:
            ax.text(xi, v + 0.04, f"{v:.3f}", ha='center', va='bottom',
                    color=COL_TXT, fontsize=8, fontweight='bold')

    # Threshold lines
    y_max = max((v for v in values if v is not None), default=SNR_GOOD)
    y_max = max(y_max, SNR_GOOD) * 1.15
    ax.set_ylim(0, y_max)
    ax.axhline(SNR_OK,   color='#FBBF24', linewidth=1.2, linestyle='--', alpha=0.8, label=f'OK (≥{SNR_OK})')
    ax.axhline(SNR_GOOD, color='#34D399', linewidth=1.2, linestyle='--', alpha=0.8, label=f'Good (≥{SNR_GOOD})')
    ax.legend(facecolor=BG_FIG, labelcolor=COL_TXT, fontsize=8, loc='upper right')

    ax.set_xticks(xs_bar)
    ax.set_xticklabels(labels, rotation=30, ha='right', color=COL_TXT, fontsize=8)
    ax.set_ylabel("Base SNR", color=COL_TXT, fontsize=9, labelpad=4)
    ax.set_title("Base Image SNR (target vs. reference ROIs)", color=COL_TXT,
                 fontsize=10, fontweight='bold', pad=6)

    buf = BytesIO()
    _MplAgg(fig).print_figure(buf, format='png', dpi=130, facecolor=BG_FIG)
    buf.seek(0)

    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(sl, C_BG)
    _add_text(sl, "Image Gallery — Base SNR Summary",
              Inches(0.3), Inches(0.08), Inches(12.8), Inches(0.44),
              size=16, bold=True, color=C_PRIMARY)
    sl.shapes.add_picture(buf, Inches(1.0), Inches(0.60), width=Inches(11.3))


def _ppt_add_diff_roi_position_slides(prs, roi_full_results, roi_all_results,
                                       _fill_bg, _add_text,
                                       C_BG, C_CARD, C_TEXT, C_TEXT_SEC, C_PRIMARY,
                                       Inches, Pt):
    """Diff Map + ROI Position slides — diff image with ROI boxes overlaid.

    For each base label, shows every diff image (one per compare LE) with the
    target/reference ROI rectangles drawn on it, mirroring the layout of the
    Base ROI Position Map slides.
    """
    from io import BytesIO
    import numpy as _np
    try:
        from pptx.dml.color import RGBColor
    except Exception:
        return

    FONT    = cv2.FONT_HERSHEY_SIMPLEX
    TGT_COL = (60,  80, 220)    # BGR blue-red (target)
    REF_COL = (200, 200,  40)   # BGR cyan-ish (reference)

    # Group roi_all_results by base label
    from collections import defaultdict as _dd
    results_by_base = _dd(list)
    for r in roi_all_results:
        results_by_base[r.base_label].append(r)

    for base_lbl, roi_full in roi_full_results.items():
        rois = roi_full.roi_set.rois
        if not rois:
            continue
        pair_results = results_by_base.get(base_lbl, [])
        if not pair_results:
            continue

        for r in pair_results:
            diff_img = getattr(r, 'result_image', None)
            if diff_img is None:
                continue

            # Convert grayscale diff to BGR for coloured ROI drawing
            if len(diff_img.shape) == 2:
                vis = cv2.cvtColor(diff_img, cv2.COLOR_GRAY2BGR)
            else:
                vis = diff_img.copy()

            h, w = vis.shape[:2]

            # Draw ROI boxes
            for roi in rois:
                nx, ny, nw, nh = roi.norm_rect
                rx  = int(nx * w)
                ry  = int(ny * h)
                rpw = max(1, int(nw * w))
                rph = max(1, int(nh * h))
                col       = TGT_COL if roi.roi_type == 'target' else REF_COL
                thickness = 3       if roi.roi_type == 'target' else 2
                cv2.rectangle(vis, (rx, ry), (rx + rpw, ry + rph), col, thickness)
                cv2.rectangle(vis, (rx + 1, ry + 1),
                              (rx + rpw - 1, ry + rph - 1), (0, 0, 0), 1)
                short_lbl  = roi.label.split('_')[-1].lstrip('0') or '0'
                font_scale = max(0.35, min(0.65, w / 1200))
                tx = rx + 3
                ty = max(ry - 4, 10)
                cv2.putText(vis, short_lbl, (tx + 1, ty + 1), FONT,
                            font_scale, (0, 0, 0), 2, cv2.LINE_AA)
                cv2.putText(vis, short_lbl, (tx, ty), FONT,
                            font_scale, col, 1, cv2.LINE_AA)

            # Encode to PNG
            try:
                from PIL import Image as _PilImg
                pil = _PilImg.fromarray(cv2.cvtColor(vis, cv2.COLOR_BGR2RGB))
                buf = BytesIO()
                pil.save(buf, format='PNG')
                buf.seek(0)
            except ImportError:
                ok, enc = cv2.imencode('.png', vis)
                if not ok:
                    continue
                buf = BytesIO(enc.tobytes())

            # Slide layout
            sl = prs.slides.add_slide(prs.slide_layouts[6])
            _fill_bg(sl, C_BG)
            _add_text(sl,
                      f"Diff Map + ROI — Base: {base_lbl}  ↔  Compare: {r.compare_label}",
                      Inches(0.3), Inches(0.08), Inches(10.5), Inches(0.46),
                      size=14, bold=True, color=C_PRIMARY)

            AVAIL_W = Inches(10.3)
            AVAIL_H = Inches(6.78)
            nat_h, nat_w = vis.shape[:2]
            if nat_w == 0 or nat_h == 0:
                disp_w, disp_h = AVAIL_W, AVAIL_H
            elif (nat_w / nat_h) >= (AVAIL_W / AVAIL_H):
                disp_w = AVAIL_W
                disp_h = int(AVAIL_W * nat_h / nat_w)
            else:
                disp_h = AVAIL_H
                disp_w = int(AVAIL_H * nat_w / nat_h)
            img_y = Inches(0.62) + (AVAIL_H - disp_h) // 2
            sl.shapes.add_picture(buf, Inches(0.2), img_y, width=disp_w, height=disp_h)

            # Legend (right side)
            LEG_X = Inches(11.0)
            LEG_Y = Inches(1.5)

            def _legend_row_d(y_off, color_rgb, label, thick_str):
                box = sl.shapes.add_shape(
                    1, LEG_X, LEG_Y + y_off, Inches(0.28), Inches(0.22))
                box.fill.solid()
                box.fill.fore_color.rgb = color_rgb
                box.line.fill.background()
                _add_text(sl, f"{label}  ({thick_str})",
                          LEG_X + Inches(0.34), LEG_Y + y_off,
                          Inches(2.0), Inches(0.28),
                          size=9, color=C_TEXT)

            _add_text(sl, "Legend", LEG_X, LEG_Y - Inches(0.32),
                      Inches(2.2), Inches(0.30),
                      size=10, bold=True, color=C_TEXT_SEC)
            _legend_row_d(Inches(0.00),
                          RGBColor(220, 80, 60),  "Target",    "thick border")
            _legend_row_d(Inches(0.35),
                          RGBColor(40, 200, 200), "Reference", "thin border")

            n_tgt = sum(1 for ro in rois if ro.roi_type == 'target')
            n_ref = sum(1 for ro in rois if ro.roi_type == 'reference')
            _add_text(sl,
                      f"Total: {len(rois)} ROIs\n"
                      f"  Target: {n_tgt}\n"
                      f"  Reference: {n_ref}",
                      LEG_X, LEG_Y + Inches(0.8),
                      Inches(2.2), Inches(0.85),
                      size=9, color=C_TEXT_SEC)


def _ppt_add_condition_gallery(prs, result_rows, crop_size,
                                _fill_bg, _add_text,
                                C_BG, C_CARD, C_TEXT, C_TEXT_SEC, C_PRIMARY,
                                Inches, Pt):
    """Final slide(s): all original images (base+compare) by condition label."""
    from io import BytesIO
    try:
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
    except Exception:
        return

    def _load_crop_png(path, size):
        """Load image, center-crop to size×size, return PNG BytesIO."""
        img = cv2.imread(path)
        if img is None:
            return None
        if len(img.shape) == 3:
            img = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        h, w = img.shape[:2]
        if h > size or w > size:
            cy, cx = h // 2, w // 2
            half = size // 2
            y1, x1 = max(0, cy - half), max(0, cx - half)
            img = img[y1: y1 + size, x1: x1 + size]
        try:
            from PIL import Image as _PilImg
            buf = BytesIO()
            _PilImg.fromarray(img, mode='L').save(buf, format='PNG')
            buf.seek(0)
            return buf
        except ImportError:
            from matplotlib.figure import Figure as _MplFig
            from matplotlib.backends.backend_agg import FigureCanvasAgg as _MplAgg
            fig_t = _MplFig(figsize=(1, 1), dpi=size)
            fig_t.patch.set_facecolor('black')
            axt = fig_t.add_axes([0, 0, 1, 1])
            axt.set_axis_off()
            axt.imshow(img, cmap='gray', vmin=0, vmax=255, aspect='auto')
            buf = BytesIO()
            _MplAgg(fig_t).print_figure(buf, format='png', dpi=size, facecolor='black')
            buf.seek(0)
            return buf

    # Collect unique (label → path), preserving first occurrence
    seen: dict = {}
    for row in result_rows:
        r = row["result"]
        p = row["paths"]
        if r.base_label not in seen and "base" in p:
            seen[r.base_label] = p["base"]
        if r.compare_label not in seen and "compare" in p:
            seen[r.compare_label] = p["compare"]
    if not seen:
        return

    ordered = sorted(seen.items(), key=lambda x: x[0])
    total = len(ordered)

    NCOLS = min(total, 5)
    NROWS_PP = 2
    PER_PAGE = NCOLS * NROWS_PP

    SLIDE_W = Inches(13.33)
    SLIDE_H = Inches(7.5)
    L_M = Inches(0.25)
    R_M = Inches(0.25)
    TITLE_H = Inches(0.52)
    T_M = Inches(0.10)
    B_M = Inches(0.12)
    GAP_X = Inches(0.10)
    GAP_Y = Inches(0.28)
    LBL_H = Inches(0.30)

    cell_w = (SLIDE_W - L_M - R_M - GAP_X * (NCOLS - 1)) / NCOLS
    avail_h = SLIDE_H - TITLE_H - T_M - B_M
    cell_img_h = (avail_h - GAP_Y * (NROWS_PP - 1)) / NROWS_PP - LBL_H

    total_pages = (total + PER_PAGE - 1) // PER_PAGE
    for slide_idx in range(0, total, PER_PAGE):
        block = ordered[slide_idx: slide_idx + PER_PAGE]
        page_n = slide_idx // PER_PAGE + 1
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        _fill_bg(sl, C_BG)

        suffix = f"  ({page_n}/{total_pages})" if total_pages > 1 else ""
        _add_text(sl, f"Image Gallery — By Condition{suffix}",
                  L_M, Inches(0.08), SLIDE_W - L_M - R_M, TITLE_H,
                  size=16, bold=True, color=C_PRIMARY)

        for ci, (lbl, path) in enumerate(block):
            row_i = ci // NCOLS
            col_i = ci % NCOLS
            x = L_M + col_i * (cell_w + GAP_X)
            y = TITLE_H + T_M + row_i * (LBL_H + cell_img_h + GAP_Y)

            # Label bar
            lb = sl.shapes.add_shape(1, x, y, cell_w, LBL_H)
            lb.fill.solid()
            lb.fill.fore_color.rgb = C_CARD
            lb.line.fill.background()
            _add_text(sl, lbl,
                      x + Inches(0.06), y + Inches(0.02),
                      cell_w - Inches(0.08), LBL_H,
                      size=9, bold=True, color=C_PRIMARY)

            # Image thumbnail — specify only width to preserve square aspect ratio
            buf = _load_crop_png(path, crop_size)
            if buf is not None:
                sl.shapes.add_picture(buf, x, y + LBL_H, width=cell_w)
            else:
                ph = sl.shapes.add_shape(1, x, y + LBL_H, cell_w, cell_img_h)
                ph.fill.solid()
                ph.fill.fore_color.rgb = C_CARD
                ph.line.fill.background()
                _add_text(sl, "[n/a]",
                          x + cell_w / 2 - Inches(0.4),
                          y + LBL_H + cell_img_h / 2 - Inches(0.15),
                          Inches(0.8), Inches(0.3),
                          size=9, color=C_TEXT_SEC, align=PP_ALIGN.CENTER)


def _ppt_add_roi_slides(prs, roi_full_results, roi_all_results,
                        _fill_bg, _add_text, _score_color,
                        C_BG, C_CARD, C_TEXT, C_TEXT_SEC,
                        C_PRIMARY, C_SUCCESS, C_WARN,
                        Inches, Pt):
    """Append ROI Analysis slides to an existing pptx Presentation object.

    Slide 1 (+ overflow): LE Summary engineering table.
    Last slide: SNR Chart rendered from matplotlib to a PNG image.
    """
    from io import BytesIO
    from matplotlib.figure import Figure as _MplFig
    from matplotlib.backends.backend_agg import FigureCanvasAgg as _MplAgg
    try:
        from pptx.dml.color import RGBColor
    except Exception:
        return

    C_SNR_GOOD   = RGBColor(0x14, 0x53, 0x2D)   # dark green-900
    C_SNR_MID    = RGBColor(0x78, 0x35, 0x07)   # dark amber-900
    C_SNR_BAD    = RGBColor(0x7F, 0x1D, 0x1D)   # dark red-900
    _STATUS_COLS = {'ok': C_SUCCESS, 'warn': C_PRIMARY, 'fail': C_WARN}

    base_labels = list(roi_full_results.keys())

    # ── Build row data ───────────────────────────────────────────────────
    roi_rows = []
    for r in roi_all_results:
        roi_full = roi_full_results.get(r.base_label)
        entry    = roi_full.snr_per_diff.get(r.compare_label) if roi_full else None
        status   = (r.alignment.status if r.alignment else '—')
        score_v  = r.alignment.final_score if r.alignment else None
        alpha    = f"{r.roi_match_alpha:.4f}" if r.roi_match_alpha is not None else '—'
        if entry is not None:
            # Scale to GLV (0-255) for display.
            mu_t_glv    = entry.mu_target * 255.0
            mu_r_glv    = entry.mu_ref    * 255.0
            sigma_glv   = entry.sigma_ref * 255.0
            delta       = mu_t_glv - mu_r_glv
            roi_rows.append({
                'base': r.base_label, 'compare': r.compare_label,
                'status': status, 'status_col': _STATUS_COLS.get(status.lower(), C_TEXT_SEC),
                'alpha': alpha,
                'score': f"{score_v:.1f}" if score_v is not None else '—', 'score_v': score_v,
                'mu_t':  f"{mu_t_glv:.2f}",
                'mu_r':  f"{mu_r_glv:.2f}",
                'sigma': f"{sigma_glv:.2f}",
                'delta': f"{delta:+.2f}",
                'snr':   f"{entry.snr:.3f}", 'snr_v': entry.snr,
            })
        else:
            roi_rows.append({
                'base': r.base_label, 'compare': r.compare_label,
                'status': status, 'status_col': _STATUS_COLS.get(status.lower(), C_TEXT_SEC),
                'alpha': alpha,
                'score': f"{score_v:.1f}" if score_v is not None else '—', 'score_v': score_v,
                'mu_t': '—', 'mu_r': '—', 'sigma': '—', 'delta': '—', 'snr': '—', 'snr_v': None,
            })

    # ── Column layout (x in inches, width in inches) ─────────────────────
    roi_cols = [
        ("Base",     0.40, 1.85),
        ("Compare",  2.30, 1.85),
        ("Status",   4.20, 0.95),
        ("α",        5.20, 0.80),
        ("Score",    6.05, 0.80),
        ("T Mean",   6.90, 1.05),
        ("R Mean",   8.00, 1.05),
        ("R Std",    9.10, 1.05),
        ("Δ",       10.20, 1.05),
        ("SNR",     11.30, 1.60),
    ]

    ROWS_PER = 20
    for block_start in range(0, max(1, len(roi_rows)), ROWS_PER):
        block = roi_rows[block_start:block_start + ROWS_PER]
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        _fill_bg(sl, C_BG)

        page_n  = block_start // ROWS_PER + 1
        n_pages = max(1, (len(roi_rows) + ROWS_PER - 1) // ROWS_PER)
        _add_text(sl, f"ROI Analysis — LE Summary  (page {page_n}/{n_pages})",
                  Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.45),
                  size=15, bold=True, color=C_PRIMARY)

        for ch, cx, cw in roi_cols:
            _add_text(sl, ch, Inches(cx), Inches(0.70), Inches(cw), Inches(0.28),
                      size=8.5, bold=True, color=C_TEXT_SEC)

        sep = sl.shapes.add_shape(1, Inches(0.4), Inches(1.01), Inches(12.9), Inches(0.02))
        sep.fill.solid()
        sep.fill.fore_color.rgb = C_TEXT_SEC
        sep.line.fill.background()

        row_h = Inches(0.275)
        for ri, d in enumerate(block):
            y = Inches(1.07) + ri * row_h
            bg = sl.shapes.add_shape(1, Inches(0.35), y, Inches(12.95), row_h)
            bg.fill.solid()
            bg.fill.fore_color.rgb = C_CARD if ri % 2 == 0 else C_BG
            bg.line.fill.background()

            snr_v = d['snr_v']
            snr_c = (C_SNR_GOOD if snr_v is not None and snr_v >= 2.0 else
                     C_SNR_MID  if snr_v is not None and snr_v >= 1.0 else
                     C_SNR_BAD  if snr_v is not None else C_TEXT_SEC)

            cell_vals = [
                (d['base'],    C_TEXT),
                (d['compare'], C_TEXT),
                (d['status'],  d['status_col']),
                (d['alpha'],   C_TEXT_SEC),
                (d['score'],   _score_color(d['score_v']) if d['score_v'] is not None else C_TEXT_SEC),
                (d['mu_t'],    C_TEXT),
                (d['mu_r'],    C_TEXT),
                (d['sigma'],   C_TEXT_SEC),
                (d['delta'],   C_TEXT),
                (d['snr'],     snr_c),
            ]
            for (val, col), (_, cx, cw) in zip(cell_vals, roi_cols):
                _add_text(sl, val, Inches(cx), y + Inches(0.02), Inches(cw), row_h,
                          size=8, color=col)

    # ── SNR Chart slide ──────────────────────────────────────────────────
    BG_FIG  = '#1F2937'
    BG_AX   = '#111827'
    COL_TXT = '#D1D5DB'
    COL_MUT = '#9CA3AF'
    COL_SPL = '#4B5563'
    COL_GRD = '#374151'
    COLORS  = ['#F59E0B', '#60A5FA', '#34D399', '#F87171', '#A78BFA', '#FB923C']
    MARKERS = ['o', 's', '^', 'D', 'v', 'P']

    def _style(ax):
        ax.set_facecolor(BG_AX)
        ax.tick_params(colors=COL_TXT)
        for sp in ('bottom', 'left'):
            ax.spines[sp].set_color(COL_SPL)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.grid(True, color=COL_GRD, linewidth=0.5)

    is_multi = len(base_labels) > 1

    if is_multi:
        fig = _MplFig(figsize=(11, 5), tight_layout=True)
        fig.patch.set_facecolor(BG_FIG)
        ax_snr, ax_delta = fig.subplots(2, 1, sharex=True)
        _style(ax_snr); _style(ax_delta)

        seen: dict = {}
        for bl in base_labels:
            rr = roi_full_results.get(bl)
            if rr:
                for k in rr.snr_per_diff:
                    seen[k] = None
        all_cmp = list(seen.keys())
        x_pos   = {lbl: i for i, lbl in enumerate(all_cmp)}

        for b_idx, bl in enumerate(base_labels):
            rr = roi_full_results.get(bl)
            if not rr or not rr.snr_per_diff:
                continue
            col = COLORS[b_idx % len(COLORS)]
            mk  = MARKERS[b_idx % len(MARKERS)]
            sd  = rr.snr_per_diff
            xs  = [x_pos[k] for k in sd]
            ax_snr.plot(xs, [sd[k].snr for k in sd],
                        marker=mk, color=col, linewidth=1.8, markersize=7, label=bl)
            db = [(sd[k].mu_target - sd[k].mu_ref) * 255.0 for k in sd]
            sb = [sd[k].sigma_ref * 255.0 for k in sd]
            ax_delta.plot(xs, db, marker=mk, color=col, linewidth=1.8, markersize=7, label=bl)
            ax_delta.errorbar(xs, db, yerr=sb, fmt='none', color=col,
                              capsize=4, linewidth=1.2, alpha=0.6)

        ax_snr.set_title("SNR = Δ / σ_Ref  — All Bases overlay", color=COL_MUT, fontsize=10)
        ax_snr.legend(facecolor=BG_FIG, labelcolor=COL_TXT, fontsize=9)
        ax_delta.set_xticks(range(len(all_cmp)))
        ax_delta.set_xticklabels(all_cmp, rotation=20, ha='right', color=COL_TXT, fontsize=9)
    else:
        bl = base_labels[0]
        rr = roi_full_results.get(bl)
        if rr is None or not rr.snr_per_diff:
            return
        fig = _MplFig(figsize=(11, 5), tight_layout=True)
        fig.patch.set_facecolor(BG_FIG)
        ax_snr, ax_delta = fig.subplots(2, 1, sharex=True)
        _style(ax_snr); _style(ax_delta)

        sd     = rr.snr_per_diff
        labels = list(sd.keys())
        xs     = list(range(len(labels)))
        sv     = [sd[k].snr for k in labels]
        dv     = [(sd[k].mu_target - sd[k].mu_ref) * 255.0 for k in labels]
        er     = [sd[k].sigma_ref * 255.0 for k in labels]

        ax_snr.bar(xs, sv, width=0.55, color='#F59E0B', alpha=0.85)
        sp = (max(sv) - min(sv)) * 0.03 if sv else 0
        for i, v in enumerate(sv):
            ax_snr.text(i, v + sp, f"{v:.2f}", ha='center', va='bottom', fontsize=8, color=COL_TXT)
        bc = ['#34D399' if d >= 0 else '#F87171' for d in dv]
        ax_delta.bar(xs, dv, width=0.55, color=bc, alpha=0.8)
        ax_delta.errorbar(xs, dv, yerr=er, fmt='none', color='#D1D5DB', capsize=5, linewidth=1.5)
        ax_snr.set_title(f"SNR = Δ / σ_Ref   [Base: {bl}]", color=COL_MUT, fontsize=10)
        ax_delta.set_xticks(xs)
        ax_delta.set_xticklabels(labels, rotation=20, ha='right', color=COL_TXT, fontsize=9)

    for ax in (ax_snr, ax_delta):
        ax.axhline(0, color=COL_SPL, linewidth=0.8, linestyle='--')
    ax_snr.set_ylabel("SNR", color=COL_TXT, fontsize=9)
    ax_delta.set_ylabel("Δ = T−R  (±σ_Ref)", color=COL_TXT, fontsize=9)

    buf = BytesIO()
    _MplAgg(fig).print_figure(buf, format='png', dpi=120, facecolor=BG_FIG)
    buf.seek(0)

    chart_sl = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(chart_sl, C_BG)
    _add_text(chart_sl, "ROI Analysis — SNR Chart",
              Inches(0.4), Inches(0.12), Inches(12.5), Inches(0.45),
              size=15, bold=True, color=C_PRIMARY)
    chart_sl.shapes.add_picture(buf, Inches(0.4), Inches(0.65), width=Inches(12.5))

    # ── Matrix overview slide (SNR Pair Matrix + Diff Image Matrix) ───────
    _ppt_add_matrix_slide(prs, roi_full_results, roi_all_results,
                          _fill_bg, _add_text,
                          C_BG, C_CARD, C_TEXT, C_TEXT_SEC, C_PRIMARY,
                          Inches, Pt)


class PerspectiveCombinationDialog(QtWidgets.QDialog):
    """Dialog for multi-image perspective combination and defect detection."""

    def __init__(self, parent=None, conditions: List[EbeamCondition] = None):
        super().__init__(parent)
        self.setWindowTitle("Fusi\u00b3 \u2014 SEM Perspective Combination Tool")
        self.setMinimumSize(1500, 900)
        self.resize(1600, 950)

        # Apply a clean sans-serif font (Arial on Windows/macOS,
        # Liberation Sans on Linux — metrically identical to Arial).
        _ui_font = QtGui.QFont()
        for candidate in ("Liberation Sans", "Arial", "Helvetica Neue", "Segoe UI"):
            _ui_font.setFamily(candidate)
            _info = QtGui.QFontInfo(_ui_font)
            if _info.family().lower().replace(" ", "") == candidate.lower().replace(" ", ""):
                break  # found a font that is actually installed
        _ui_font.setPointSize(10)
        self.setFont(_ui_font)

        self.setStyleSheet(DIALOG_STYLE)

        self._conditions = conditions or []
        self._images: Dict[str, np.ndarray] = {}
        self._results: List[SinglePairResult] = []  # Multiple results
        self._current_result_idx: int = 0  # Current result index
        self._result: Optional[CombineResult] = None  # Legacy
        self._last_settings: Dict[str, object] = {}
        self._compute_thread: Optional[QtCore.QThread] = None
        self._compute_worker: Optional[_Worker] = None
        self._display_mode = 'diff'  # right viewer: 'diff' or 'zmap'
        self._left_non_split_mode = 'base'
        self._left_view_mode = 'base'
        self._norm_compare_dialog: Optional[NormalizedCompareDialog] = None
        self._hist_range: Optional[tuple] = None  # (lo, hi) gray-level range for highlight, or None
        # Multi-ROI state
        self._multi_roi_set: MultiROISet = MultiROISet()
        self._roi_manager: Optional[MultiROIManagerWidget] = None
        self._has_computed: bool = False  # True after first successful compute
        self._roi_profile_dialog: Optional[ROIIntensityProfileDialog] = None
        # Last ROI analysis results keyed by base_label.
        # In auto-pair mode there are multiple base groups; standard mode has one.
        self._roi_full_results: Dict[str, ROIFullResult] = {}

        # The base-image label that was visible in the base viewer when the user
        # last opened / modified the ROI Manager.  ROI norm_rect values are
        # defined in this image's coordinate space and must be remapped for every
        # other base group before computing ROI stats.
        self._roi_ref_base_label: Optional[str] = None

        # Per-base remapped ROI sets (populated by _run_roi_analysis).
        # Key = base_label; value = MultiROISet with coords in that base's space.
        # Used by _apply_roi_visibility so the visual ROI overlay stays correct
        # when the user browses results with different base images.
        self._roi_remapped_sets: Dict[str, MultiROISet] = {}
        self.tutorial_overlay: Optional[WelcomeTutorialOverlay] = None
        self._tutorial_checked = False
        self._compute_abort_requested: bool = False
        self._compute_lock_effects: Dict[QtWidgets.QWidget, QtWidgets.QGraphicsOpacityEffect] = {}
        self._compute_lock_targets: Tuple[QtWidgets.QWidget, ...] = ()

        self._setup_ui()
        self._apply_toolbar_icons()
        self._connect_signals()
        self._setup_tutorial_overlay()
        self._load_images()
        # Initial badge and step-state update after all widgets exist
        QtCore.QTimer.singleShot(0, self._update_adv_badge)
        QtCore.QTimer.singleShot(0, self._update_step_states)

    def _set_button_icon(self, button: QtWidgets.QPushButton, pixmap_enum, text: str = None, size: int = 16):
        if text is not None:
            button.setText(text)
        icon = self.style().standardIcon(pixmap_enum)
        button.setIcon(icon)
        button.setIconSize(QtCore.QSize(size, size))

    def _polish_button(self, button: QtWidgets.QPushButton, variant: str | None = None):
        if variant is not None:
            button.setProperty("variant", variant)
        button.style().unpolish(button)
        button.style().polish(button)
        button.update()

    def _setup_tutorial_overlay(self):
        """Setup welcome tutorial overlay."""
        self.tutorial_overlay = WelcomeTutorialOverlay(self)
        self.tutorial_overlay.tutorial_finished.connect(self._on_tutorial_finished)
        self.tutorial_overlay.tutorial_skipped.connect(self._on_tutorial_skipped)
    
    def _show_welcome_tutorial(self):
        """Show welcome tutorial for first-time users."""
        if self.tutorial_overlay:
            self.tutorial_overlay.show_tutorial()
    
    def _on_tutorial_finished(self):
        """Handle tutorial completion."""
        mark_tutorial_completed()
        # Optional: Show a brief success message
        QtWidgets.QMessageBox.information(
            self, 
            "導覽完成", 
            "歡迎使用 Fusi³！您現在可以開始進行 SEM 影像融合分析了。"
        )
    
    def _on_tutorial_skipped(self):
        """Handle tutorial skip."""
        mark_tutorial_completed()
    
    def _apply_toolbar_icons(self):
        self._set_button_icon(self.btn_load_folder, QtWidgets.QStyle.SP_DirOpenIcon, "Load Folder")
        self._set_button_icon(self.btn_compute, QtWidgets.QStyle.SP_MediaPlay, "Compute")
        self._set_button_icon(self.btn_export, QtWidgets.QStyle.SP_DialogSaveButton, "Export")
        self._set_button_icon(self.btn_swap_base, QtWidgets.QStyle.SP_ArrowRight, "Swap Base")
        self._set_button_icon(self.btn_clear_hist_range, QtWidgets.QStyle.SP_DialogResetButton, "Clear Range")
        self._set_button_icon(self.btn_show_norm_compare, QtWidgets.QStyle.SP_FileDialogContentsView, "Normalize")
        self._set_button_icon(self.btn_qf_auto_detect, QtWidgets.QStyle.SP_FileDialogDetailedView, "Auto Detect")
        self._set_button_icon(self.btn_qf_pick_roi, QtWidgets.QStyle.SP_DialogOpenButton, "Pick ROI")
        self._set_button_icon(self.btn_qf_clear_roi, QtWidgets.QStyle.SP_DialogResetButton, "Clear ROI")
        self._set_button_icon(self.btn_pick_roi, QtWidgets.QStyle.SP_DialogOpenButton, "ROI Manager")
        self._set_button_icon(self.btn_preview_glv_mask, QtWidgets.QStyle.SP_DesktopIcon, "Preview Mask")

    def closeEvent(self, event):
        """Handle dialog close - ensure compute thread is stopped safely."""
        if self._compute_thread is not None:
            self._compute_thread.requestInterruption()
            self._compute_thread.quit()
            self._compute_worker = None
        event.accept()

    def showEvent(self, event):
        """Show tutorial overlay on first launch after dialog is visible."""
        super().showEvent(event)
        if self._tutorial_checked:
            return

        self._tutorial_checked = True
        if should_show_tutorial():
            QtCore.QTimer.singleShot(0, self._show_welcome_tutorial)

    def keyPressEvent(self, event):
        """Handle keyboard shortcuts for navigation."""
        key = event.key()

        # ← Arrow: Previous result
        if key == Qt.Key_Left:
            self._on_prev_result()
            event.accept()
            return

        # → Arrow: Next result
        if key == Qt.Key_Right:
            self._on_next_result()
            event.accept()
            return

        # Space: Compute
        if key == Qt.Key_Space and self.btn_compute.isEnabled():
            self._on_compute()
            event.accept()
            return

        # Escape: Close dialog
        if key == Qt.Key_Escape:
            self.close()
            event.accept()
            return

        super().keyPressEvent(event)

    def _setup_ui(self):
        """Build the dialog UI — target-style layout with toolbar, compact sidebar, viewer, bottom cards."""
        outer_layout = QtWidgets.QVBoxLayout(self)
        outer_layout.setContentsMargins(8, 8, 8, 8)
        outer_layout.setSpacing(8)

        # ================================================================
        # TOP TOOLBAR
        # ================================================================
        toolbar = QtWidgets.QFrame()
        toolbar.setObjectName("MainToolbar")
        self.toolbar = toolbar
        toolbar_layout = QtWidgets.QHBoxLayout(toolbar)
        toolbar_layout.setContentsMargins(8, 6, 8, 6)
        toolbar_layout.setSpacing(6)

        # Load Folder
        self.btn_load_folder = QtWidgets.QPushButton("Load Folder")
        self.btn_load_folder.setObjectName("ToolbarAction")
        toolbar_layout.addWidget(self.btn_load_folder)

        # Navigation arrows (prev/next folder image set)
        self.btn_prev_result = QtWidgets.QPushButton("\u25C0")
        self.btn_prev_result.setFixedWidth(32)
        self.btn_prev_result.setObjectName("ToolbarAction")
        self.btn_prev_result.setEnabled(False)
        self.btn_prev_result.setVisible(False)   # hidden until compute
        self.btn_next_result = QtWidgets.QPushButton("\u25B6")
        self.btn_next_result.setFixedWidth(32)
        self.btn_next_result.setObjectName("ToolbarAction")
        self.btn_next_result.setEnabled(False)
        self.btn_next_result.setVisible(False)   # hidden until compute
        toolbar_layout.addWidget(self.btn_prev_result)
        toolbar_layout.addWidget(self.btn_next_result)

        toolbar_layout.addStretch(1)

        self.btn_compute = QtWidgets.QPushButton("Compute")
        self.btn_compute.setObjectName("ToolbarPrimary")
        self.btn_compute.setMinimumHeight(44)
        # btn_compute is placed at bottom of left panel, not in toolbar

        self.btn_export = QtWidgets.QPushButton("Export")
        self.btn_export.setObjectName("ToolbarSecondary")
        self.btn_export.setEnabled(False)
        self.btn_export.setVisible(False)  # shown only in post-compute view

        self.btn_roi_manager = QtWidgets.QPushButton("ROI Manager")
        self.btn_roi_manager.setObjectName("ToolbarSecondary")
        self.btn_roi_manager.setToolTip("Open Multi-ROI Manager to define bounding-box ROIs")
        # btn_roi_manager is placed at bottom of left panel, not in toolbar

        # Right-aligned export action (replaces legacy About button slot)
        toolbar_layout.addWidget(self.btn_export)

        outer_layout.addWidget(toolbar)

        # Result info label (spans full width, below toolbar)
        self.lbl_result_info = QtWidgets.QLabel("")
        self.lbl_result_info.setAlignment(Qt.AlignCenter)
        self.lbl_result_info.setStyleSheet(f"""
            QLabel {{
                color: {UI_TEXT};
                background-color: {UI_BG_CARD};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.SM};
                padding: 4px 12px;
                font-size: {Typography.FONT_SIZE_SMALL};
            }}
        """)
        self.lbl_result_info.setFixedHeight(28)
        self.lbl_result_info.setVisible(False)
        outer_layout.addWidget(self.lbl_result_info)

        # ================================================================
        # CONTENT AREA: sidebar + viewer
        # ================================================================
        content_layout = QtWidgets.QHBoxLayout()
        content_layout.setSpacing(8)
        content_layout.setContentsMargins(0, 0, 0, 0)
        self._content_layout = content_layout  # stored for ROI slide-in panel

        # === LEFT SIDEBAR ===
        self.left_panel = QtWidgets.QWidget()
        left_panel = self.left_panel   # local alias for existing code below
        left_panel.setObjectName("LeftPanel")
        self._sidebar_pref_width = 320
        self._sidebar_min_width = 260
        self._sidebar_max_width = 360
        left_panel.setMinimumWidth(self._sidebar_min_width)
        left_panel.setMaximumWidth(self._sidebar_max_width)
        left_panel.setSizePolicy(QtWidgets.QSizePolicy.Preferred, QtWidgets.QSizePolicy.MinimumExpanding)
        left_panel.setStyleSheet(f"""
            QWidget#LeftPanel {{
                background-color: {UI_LEFT_PANEL_BG};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.MD};
            }}
            QWidget#LeftPanel > QWidget,
            QWidget#LeftPanel QWidget {{
                background-color: transparent;
                border: none;
            }}
            QWidget#LeftPanel QComboBox {{
                background-color: {UI_BG_INPUT};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.SM};
                padding: {Spacing.INPUT_PADDING};
                font-size: {Typography.FONT_SIZE_SMALL};
                min-height: 30px;
            }}
            QWidget#LeftPanel QComboBox:hover,
            QWidget#LeftPanel QComboBox:focus {{
                border-color: {UI_PRIMARY};
            }}
            QWidget#LeftPanel QComboBox::drop-down {{
                border: none;
                padding-right: 8px;
            }}
            QWidget#LeftPanel QComboBox QAbstractItemView {{
                background-color: {UI_BG_PANEL};
                color: {UI_TEXT};
                selection-background-color: {UI_PRIMARY};
                selection-color: {UI_TEXT_ON_PRIMARY};
                border: 1px solid {UI_BORDER};
            }}
            QWidget#LeftPanel QCheckBox {{
                color: {UI_TEXT};
                spacing: {Spacing.SM};
                font-size: {Typography.FONT_SIZE_SMALL};
                background: transparent;
                border: none;
                min-height: 24px;
            }}
            QWidget#LeftPanel QCheckBox::indicator {{
                width: 16px;
                height: 16px;
                border: 2px solid #D1D5DB;
                border-radius: 3px;
                background-color: {UI_BG_INPUT};
            }}
            QWidget#LeftPanel QCheckBox::indicator:checked {{
                background-color: {UI_PRIMARY};
                border-color: {UI_PRIMARY};
            }}
            QWidget#LeftPanel QSpinBox,
            QWidget#LeftPanel QDoubleSpinBox {{
                background-color: {UI_BG_INPUT};
                color: {UI_TEXT};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.SM};
                padding: 4px 22px 4px 8px;
                min-height: 30px;
                font-size: {Typography.FONT_SIZE_SMALL};
            }}
            QWidget#LeftPanel QSpinBox:hover,
            QWidget#LeftPanel QDoubleSpinBox:hover {{
                border-color: {UI_PRIMARY};
            }}
            QWidget#LeftPanel QSpinBox::up-button,
            QWidget#LeftPanel QDoubleSpinBox::up-button,
            QWidget#LeftPanel QSpinBox::down-button,
            QWidget#LeftPanel QDoubleSpinBox::down-button {{
                subcontrol-origin: border;
                width: 14px;
                border-left: 1px solid {UI_BORDER};
                background-color: {UI_BG_CARD};
            }}
            QWidget#LeftPanel QSpinBox::up-button,
            QWidget#LeftPanel QDoubleSpinBox::up-button {{
                subcontrol-position: top right;
                border-top-right-radius: {BorderRadius.SM};
            }}
            QWidget#LeftPanel QSpinBox::down-button,
            QWidget#LeftPanel QDoubleSpinBox::down-button {{
                subcontrol-position: bottom right;
                border-bottom-right-radius: {BorderRadius.SM};
                border-top: 1px solid {UI_BORDER};
            }}
            QWidget#LeftPanel QPushButton {{
                background-color: {UI_BG_PANEL};
                color: {UI_TEXT};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.SM};
                padding: 4px 10px;
                min-height: 30px;
                font-size: {Typography.FONT_SIZE_SMALL};
                font-weight: {Typography.FONT_WEIGHT_MEDIUM};
            }}
            QWidget#LeftPanel QPushButton:hover {{
                background-color: #FFF8ED;
                border-color: {UI_PRIMARY};
            }}
            QWidget#LeftPanel QPushButton:pressed {{
                background-color: #FDE7C2;
            }}
            QWidget#LeftPanel QPushButton[role="primary"] {{
                background-color: {UI_PRIMARY};
                color: {UI_TEXT_ON_PRIMARY};
                border: 1px solid {UI_PRIMARY};
                border-radius: {BorderRadius.SM};
                min-height: 34px;
                font-weight: {Typography.FONT_WEIGHT_BOLD};
            }}
            QWidget#LeftPanel QPushButton[role="primary"]:hover {{
                background-color: {UI_ACCENT_HOVER};
                border-color: {UI_ACCENT_HOVER};
            }}
            QWidget#LeftPanel QPushButton[role="secondary"] {{
                background-color: {UI_BG_PANEL};
                border: 1px solid #D1D5DB;
                color: {UI_TEXT};
            }}
            QWidget#LeftPanel QPushButton[role="utility"] {{
                background: transparent;
                border: none;
                color: {UI_PRIMARY};
                min-width: 0;
                padding: 2px 4px;
                font-weight: {Typography.FONT_WEIGHT_SEMIBOLD};
            }}
            QWidget#LeftPanel QPushButton[role="utility"]:hover {{
                color: {UI_ACCENT_HOVER};
                background: transparent;
                border: none;
            }}
            QWidget#LeftPanel QLabel {{
                color: {UI_TEXT};
                font-size: {Typography.FONT_SIZE_SMALL};
                background: transparent;
                border: none;
            }}
            QWidget#LeftPanel QLabel#SectionTitle {{
                color: {UI_TEXT_PRIMARY_STRONG};
                font-size: {Typography.FONT_SIZE_H3};
                font-weight: {Typography.FONT_WEIGHT_BOLD};
                text-transform: uppercase;
                padding: 4px 0px 2px 0px;
                letter-spacing: 0.8px;
            }}
            QWidget#LeftPanel QLabel#SectionSeparator {{
                background-color: {UI_PRIMARY};
                min-height: 2px;
                max-height: 2px;
                border: none;
                margin-top: 2px;
                margin-bottom: 2px;
            }}
            QWidget#LeftPanel QFrame[card="true"] {{
                background-color: #FFFFFF;
                border: 1px solid {UI_BORDER};
                border-radius: 6px;
            }}
            QWidget#LeftPanel QCheckBox[compareItem="true"] {{
                padding: 4px 6px;
                margin-bottom: 2px;
                border-radius: {BorderRadius.SM};
            }}
            QWidget#LeftPanel QCheckBox[compareItem="true"]:hover {{
                background-color: {UI_WINDOW_BG_LIGHT};
            }}
            QWidget#LeftPanel QCheckBox[baseItem="true"] {{
                background-color: {UI_ACCENT_LIGHT};
                border-left: 4px solid {UI_PRIMARY};
                padding-left: 4px;
                border-radius: {BorderRadius.SM};
            }}
            QWidget#LeftPanel QScrollArea {{
                border: none;
                background: transparent;
            }}
        """)
        left_layout = QtWidgets.QVBoxLayout(left_panel)
        left_layout.setContentsMargins(12, 12, 12, 12)
        left_layout.setSpacing(16)

        def _make_sidebar_card(title: str) -> tuple[QtWidgets.QFrame, QtWidgets.QVBoxLayout]:
            card = QtWidgets.QFrame()
            card.setProperty("card", True)
            card_layout = QtWidgets.QVBoxLayout(card)
            card_layout.setContentsMargins(10, 10, 10, 10)
            card_layout.setSpacing(6)

            lbl_title = QtWidgets.QLabel(title)
            lbl_title.setObjectName("SectionTitle")
            card_layout.addWidget(lbl_title)

            sep = QtWidgets.QLabel()
            sep.setObjectName("SectionSeparator")
            card_layout.addWidget(sep)
            card_layout.addSpacing(2)
            return card, card_layout

        def _make_step_header(num_text: str, title: str, optional: bool = False):
            """Create a numbered step header row.

            Returns (frame, circle_label, status_label) where:
              - circle_label  : QLabel showing the step number; tint to orange on complete
              - status_label  : QLabel showing '' / '✓' / '—'
            """
            frame = QtWidgets.QFrame()
            frame.setObjectName("StepHeader")
            frame.setContentsMargins(0, 0, 0, 0)
            h_lay = QtWidgets.QHBoxLayout(frame)
            h_lay.setContentsMargins(4, 8, 4, 2)
            h_lay.setSpacing(8)

            circle = QtWidgets.QLabel(num_text)
            circle.setFixedSize(26, 26)
            circle.setAlignment(Qt.AlignCenter)
            circle.setStyleSheet(
                "QLabel { background-color: #D1D5DB; color: #1F2937; border-radius: 13px;"
                " font-weight: 700; font-size: 12px; }"
            )
            h_lay.addWidget(circle)

            lbl = QtWidgets.QLabel(title)
            lbl.setStyleSheet(
                f"font-weight: {Typography.FONT_WEIGHT_BOLD}; font-size: 13px;"
                f" color: {UI_TEXT_PRIMARY_STRONG}; background: transparent; border: none;"
            )
            h_lay.addWidget(lbl)

            if optional:
                tag = QtWidgets.QLabel("optional")
                tag.setStyleSheet(
                    "color: #9CA3AF; font-size: 9px; border: 1px solid #D1D5DB;"
                    " border-radius: 3px; padding: 1px 4px; background: transparent;"
                )
                h_lay.addWidget(tag)

            h_lay.addStretch()

            status = QtWidgets.QLabel("")
            status.setStyleSheet("font-size: 13px; color: #6B7280; background: transparent; border: none;")
            h_lay.addWidget(status)

            return frame, circle, status

        # ─── Step 1: Images ─────────────────────────────────────────────────────
        step1_hdr, self._step1_circle, self._step1_status = _make_step_header("①", "Images")
        left_layout.addWidget(step1_hdr)

        # --- INPUT section ---
        card_input, input_layout = _make_sidebar_card("INPUT")
        left_layout.addWidget(card_input)

        self.lbl_input_empty_hint = QtWidgets.QLabel("No images loaded.\nLoad a folder to begin.")
        self.lbl_input_empty_hint.setProperty("secondary", True)
        self.lbl_input_empty_hint.setWordWrap(True)
        input_layout.addWidget(self.lbl_input_empty_hint)

        lbl_mode = QtWidgets.QLabel("Mode")
        lbl_mode.setStyleSheet(f"font-weight: {Typography.FONT_WEIGHT_SEMIBOLD}; color: {UI_TEXT_PRIMARY_STRONG};")
        input_layout.addWidget(lbl_mode)

        # Input mode (Standard only; Quadrant Fusion removed)
        self.cmb_input_mode = QtWidgets.QComboBox()
        self.cmb_input_mode.addItems(["Standard (Base/Compare)"])
        self.cmb_input_mode.setToolTip(
            "Classic Base vs Compare subtract/blend workflow."
        )
        input_layout.addWidget(self.cmb_input_mode)

        # ── Standard mode widgets (Base/Compare) ─────────────────────────────
        self.wgt_standard_select = QtWidgets.QWidget()
        std_layout = QtWidgets.QVBoxLayout(self.wgt_standard_select)
        std_layout.setContentsMargins(0, 4, 0, 0)
        std_layout.setSpacing(6)

        # Base Image label
        lbl_base = QtWidgets.QLabel("Base Image")
        lbl_base.setStyleSheet(f"font-weight: {Typography.FONT_WEIGHT_SEMIBOLD}; font-size: {Typography.FONT_SIZE_SMALL}; color: {UI_TEXT};")
        std_layout.addWidget(lbl_base)

        self.cmb_base = QtWidgets.QComboBox()
        std_layout.addWidget(self.cmb_base)

        self.lbl_base_empty_hint = QtWidgets.QLabel("Select a base image after loading a folder.")
        self.lbl_base_empty_hint.setProperty("secondary", True)
        self.lbl_base_empty_hint.setWordWrap(True)
        std_layout.addWidget(self.lbl_base_empty_hint)

        # Compare images (scrollable checkbox list)
        self.scroll_compare = QtWidgets.QScrollArea()
        self.scroll_compare.setWidgetResizable(True)
        self.scroll_compare.setMaximumHeight(140)
        self.scroll_compare.setStyleSheet("border: none; background: transparent;")
        self.compare_container = QtWidgets.QWidget()
        self.compare_layout = QtWidgets.QVBoxLayout(self.compare_container)
        self.compare_layout.setContentsMargins(2, 2, 2, 2)
        self.compare_layout.setSpacing(4)

        compare_hdr_row = QtWidgets.QHBoxLayout()
        compare_hdr_row.setContentsMargins(0, 0, 0, 0)
        compare_hdr_row.setSpacing(4)
        self.compare_hdr_row = compare_hdr_row
        self.lbl_compare_title = QtWidgets.QLabel("Compare Images (0)")
        self.lbl_compare_title.setStyleSheet(f"font-weight: {Typography.FONT_WEIGHT_SEMIBOLD}; color: {UI_TEXT_PRIMARY_STRONG};")
        compare_hdr_row.addWidget(self.lbl_compare_title)
        compare_hdr_row.addStretch()

        # Utility buttons in compare header row — give explicit style so they are
        # clearly identifiable as clickable buttons (flat=True strips all chrome)
        _util_btn_style = (
            "QPushButton {"
            "  background-color: #F3F4F6;"
            "  color: #374151;"
            "  border: 1px solid #D1D5DB;"
            "  border-radius: 4px;"
            "  padding: 2px 8px;"
            "  font-size: 11px;"
            "  font-weight: 600;"
            "}"
            "QPushButton:hover {"
            "  background-color: #E5E7EB;"
            "  border-color: #9CA3AF;"
            "}"
            "QPushButton:pressed {"
            "  background-color: #D1D5DB;"
            "}"
        )
        self.btn_select_all = QtWidgets.QPushButton("All")
        self.btn_select_all.setProperty("role", "utility")
        self.btn_select_all.setStyleSheet(_util_btn_style)
        self.btn_select_all.setFixedHeight(22)
        self.btn_select_none = QtWidgets.QPushButton("None")
        self.btn_select_none.setProperty("role", "utility")
        self.btn_select_none.setStyleSheet(_util_btn_style)
        self.btn_select_none.setFixedHeight(22)
        compare_hdr_row.addWidget(self.btn_select_all)
        compare_hdr_row.addWidget(self.btn_select_none)
        std_layout.addLayout(compare_hdr_row)

        self.scroll_compare.setWidget(self.compare_container)
        std_layout.addWidget(self.scroll_compare)

        # Auto pair checkbox moved to pairing card
        self.chk_auto_pair = QtWidgets.QCheckBox("Auto Pair")
        self.chk_auto_pair.setToolTip("Generate all unique pairs from selected images")

        # Swap Base button moved to pairing card
        self.btn_swap_base = QtWidgets.QPushButton("\u25b6 Swap Base")
        self.btn_swap_base.setProperty("role", "secondary")
        self.btn_swap_base.setFixedHeight(26)
        self.btn_swap_base.setToolTip("Swap base image with compare image selection")

        input_layout.addWidget(self.wgt_standard_select)

        # --- PAIRING section ---
        card_pairing, pairing_layout = _make_sidebar_card("PAIRING")
        pairing_layout.addWidget(self.chk_auto_pair)
        pairing_layout.addWidget(self.btn_swap_base)
        left_layout.addWidget(card_pairing)

        # Quadrant Fusion placeholders (feature removed from UI)
        self.wgt_quadrant_select = QtWidgets.QWidget()
        self.wgt_quadrant_select.setVisible(False)
        self._qf_combos: Dict[str, QtWidgets.QComboBox] = {
            name: QtWidgets.QComboBox() for name in ("Illuminator", "Top", "Bottom", "Left", "Right")
        }
        self.btn_qf_auto_detect = QtWidgets.QPushButton("Auto-detect")
        self.btn_qf_pick_roi = QtWidgets.QPushButton("Pick ROI")
        self.btn_qf_clear_roi = QtWidgets.QPushButton("Clear ROI")
        self.cmb_qf_output = QtWidgets.QComboBox()
        self.cmb_qf_output.addItems(["BSE Enhanced", "Topography", "Composite"])
        self.cmb_qf_alpha_mode = QtWidgets.QComboBox()
        self.cmb_qf_alpha_mode.addItems(["Auto", "Manual"])
        self.spn_qf_alpha = QtWidgets.QDoubleSpinBox()
        self.spn_qf_beta = QtWidgets.QDoubleSpinBox()
        self.spn_qf_sigma = QtWidgets.QDoubleSpinBox()
        self.lbl_qf_roi_info = QtWidgets.QLabel("")
        self._qf_roi_rect: Optional[tuple] = None
        self._qf_last_result: Optional[QuadrantFusionResult] = None

        # ─── Step 2: Configure ──────────────────────────────────────────────────
        step2_hdr, self._step2_circle, self._step2_status = _make_step_header("②", "Configure")
        left_layout.addWidget(step2_hdr)

        # --- OPERATION section ---
        card_operation, operation_card_layout = _make_sidebar_card("OPERATION")
        left_layout.addWidget(card_operation)

        # Operation Settings (Standard mode only)
        self.grp_op = QtWidgets.QWidget()
        op_layout = QtWidgets.QVBoxLayout(self.grp_op)
        op_layout.setContentsMargins(0, 4, 0, 0)
        op_layout.setSpacing(4)

        # Operation as radio-button-style checkboxes
        self.cmb_operation = QtWidgets.QComboBox()
        self.cmb_operation.addItems(["Subtract (|Base \u2212 Compare|)", "Blend (\u03b1\u00d7Base + \u03b2\u00d7Compare)"])
        # Combobox styled by sidebar stylesheet
        op_layout.addWidget(self.cmb_operation)

        # Blend coefficients (only visible in Blend mode)
        self.grp_blend_coef = QtWidgets.QWidget()
        blend_layout = QtWidgets.QHBoxLayout(self.grp_blend_coef)
        blend_layout.setContentsMargins(0, 0, 0, 0)
        blend_layout.addWidget(QtWidgets.QLabel("α (Base):"))
        self.spin_alpha = QtWidgets.QDoubleSpinBox()
        self.spin_alpha.setRange(0.0, 1.0)
        self.spin_alpha.setSingleStep(0.1)
        self.spin_alpha.setValue(0.5)
        self.spin_alpha.setFixedWidth(74)
        blend_layout.addWidget(self.spin_alpha)
        blend_layout.addWidget(QtWidgets.QLabel("β (Cmp):"))
        self.spin_beta = QtWidgets.QDoubleSpinBox()
        self.spin_beta.setRange(0.0, 1.0)
        self.spin_beta.setSingleStep(0.1)
        self.spin_beta.setValue(0.5)
        self.spin_beta.setFixedWidth(74)
        blend_layout.addWidget(self.spin_beta)
        self.grp_blend_coef.setVisible(False)
        op_layout.addWidget(self.grp_blend_coef)

        # ── Advanced options (collapsed by default) ─────────────────────────
        self.btn_adv_toggle = QtWidgets.QPushButton("Advanced Settings \u25b6")
        self.btn_adv_toggle.setCheckable(True)
        self.btn_adv_toggle.setChecked(False)
        self.btn_adv_toggle.setProperty("variant", "ghost")
        op_layout.addWidget(self.btn_adv_toggle)

        # P1-1: compact summary badge shown when Advanced Settings is collapsed
        self.lbl_adv_badge = QtWidgets.QLabel("")
        self.lbl_adv_badge.setWordWrap(True)
        self.lbl_adv_badge.setStyleSheet(
            "color: #9CA3AF; font-size: 10px; border: none; background: transparent;"
            " padding-left: 6px; padding-bottom: 2px;"
        )
        op_layout.addWidget(self.lbl_adv_badge)

        self.grp_advanced = QtWidgets.QWidget()
        adv_layout = QtWidgets.QVBoxLayout(self.grp_advanced)
        adv_layout.setContentsMargins(8, 0, 0, 0)
        adv_layout.setSpacing(4)
        # ── Subtract Options group ────────────────────────────────────────────
        lbl_subtract_opts = QtWidgets.QLabel("Subtract Options")
        lbl_subtract_opts.setStyleSheet(
            f"color: {UI_TEXT_PRIMARY_STRONG}; font-weight: {Typography.FONT_WEIGHT_SEMIBOLD};"
        )
        adv_layout.addWidget(lbl_subtract_opts)

        sub_mode_row = QtWidgets.QVBoxLayout()
        sub_mode_row.setContentsMargins(0, 0, 0, 0)
        sub_mode_row.setSpacing(6)
        sub_mode_row.addWidget(QtWidgets.QLabel("Subtract mode"))
        self.cmb_subtract_mode = QtWidgets.QComboBox()
        self.cmb_subtract_mode.addItems([
            "|diff| × 2  (default)",  # index 0
            "|diff|  (abs, no gain)",  # index 1
            "clip ≥ 0  (keep direction)",  # index 2
        ])
        self.cmb_subtract_mode.setToolTip(
            "|diff| × 2   : |Base−Compare|, then ×2 to enhance small differences (default)\n"
            "|diff|        : |Base−Compare|, no gain — preserves true magnitude\n"
            "clip ≥ 0      : Base−Compare, keep direction, clamp negatives to 0"
        )
        sub_mode_row.addWidget(self.cmb_subtract_mode)
        adv_layout.addLayout(sub_mode_row)

        invert_row = QtWidgets.QHBoxLayout()
        invert_row.setSpacing(8)
        self.invert_row = invert_row
        self.chk_invert_base = QtWidgets.QCheckBox("Inv Base")
        self.chk_invert_base.setToolTip("Apply 255−X to Base before operation")
        self.chk_invert_compare = QtWidgets.QCheckBox("Inv Cmp")
        self.chk_invert_compare.setToolTip("Apply 255−X to Compare before operation")
        self.chk_invert_result = QtWidgets.QCheckBox("Inv Result")
        self.chk_invert_result.setToolTip("Apply 255−X to result after operation")
        invert_row.addWidget(self.chk_invert_base)
        invert_row.addWidget(self.chk_invert_compare)
        invert_row.addWidget(self.chk_invert_result)
        adv_layout.addLayout(invert_row)

        adv_layout.addSpacing(10)

        # ── Normalization group ───────────────────────────────────────────────
        lbl_norm_opts = QtWidgets.QLabel("Normalization")
        lbl_norm_opts.setStyleSheet(
            f"color: {UI_TEXT_PRIMARY_STRONG}; font-weight: {Typography.FONT_WEIGHT_SEMIBOLD};"
        )
        adv_layout.addWidget(lbl_norm_opts)

        # ── Normalize mode drop-down ──────────────────────────────────────────
        norm_mode_row = QtWidgets.QVBoxLayout()
        norm_mode_row.setContentsMargins(0, 0, 0, 0)
        norm_mode_row.setSpacing(6)
        norm_mode_row.addWidget(QtWidgets.QLabel("Normalize method"))
        self.cmb_normalize_mode = QtWidgets.QComboBox()
        self.cmb_normalize_mode.addItems([
            "Percentile (P2–P98)",  # index 0 – default
            "GLV-Mask",  # index 1
            "Skip (raw ÷ 255)",  # index 2
            "ROI-Match (EPI Nulling)",  # index 3
        ])
        self.cmb_normalize_mode.setToolTip(
            "Percentile: each image independently mapped to [0,1] via its P2–P98 range.\n"
            "GLV-Mask:   P2/P98 computed only from pixels inside the specified GLV range\n"
            "            (e.g. MG 110–145, EPI 200–255); map applied to full image.\n"
            "Skip:       bypass normalization; divide pixels by 255 directly.\n"
            "ROI-Match:  Calibrate a scale factor from ROI Manager ROIs (bounding-box means)\n"
            "            so that reference region response cancels in subtraction, leaving residual\n"
            "            HK/Hf defect signals near inner spacer more visible."
        )
        norm_mode_row.addWidget(self.cmb_normalize_mode)
        adv_layout.addLayout(norm_mode_row)

        # GLV-Mask controls (shown only when GLV-Mask mode is selected)
        self.wgt_glv_controls = QtWidgets.QWidget()
        glv_ctrl_layout = QtWidgets.QVBoxLayout(self.wgt_glv_controls)
        glv_ctrl_layout.setContentsMargins(0, 0, 0, 0)
        glv_ctrl_layout.setSpacing(6)

        glv_ctrl_layout.addWidget(QtWidgets.QLabel("GLV Range"))

        glv_min_row = QtWidgets.QVBoxLayout()
        glv_min_row.setContentsMargins(0, 0, 0, 0)
        glv_min_row.setSpacing(4)
        glv_min_row.addWidget(QtWidgets.QLabel("Min"))
        self.spn_glv_low = QtWidgets.QSpinBox()
        self.spn_glv_low.setRange(0, 254)
        self.spn_glv_low.setValue(100)
        self.spn_glv_low.setMinimumWidth(96)
        self.spn_glv_low.setToolTip("Lower bound of GLV mask (inclusive, 0–255)")
        glv_min_row.addWidget(self.spn_glv_low)
        glv_ctrl_layout.addLayout(glv_min_row)

        glv_max_row = QtWidgets.QVBoxLayout()
        glv_max_row.setContentsMargins(0, 0, 0, 0)
        glv_max_row.setSpacing(4)
        glv_max_row.addWidget(QtWidgets.QLabel("Max"))
        self.spn_glv_high = QtWidgets.QSpinBox()
        self.spn_glv_high.setRange(1, 255)
        self.spn_glv_high.setValue(160)
        self.spn_glv_high.setMinimumWidth(96)
        self.spn_glv_high.setToolTip("Upper bound of GLV mask (inclusive, 0–255)")
        glv_max_row.addWidget(self.spn_glv_high)
        glv_ctrl_layout.addLayout(glv_max_row)

        self.btn_preview_glv_mask = QtWidgets.QPushButton("Preview Mask")
        self.btn_preview_glv_mask.setProperty("role", "secondary")
        self.btn_preview_glv_mask.setToolTip(
            "Show a preview window highlighting which pixels of the Base image\n"
            "fall within the GLV range and will be used for normalization."
        )
        self.btn_preview_glv_mask.clicked.connect(self._on_preview_glv_mask)
        glv_ctrl_layout.addWidget(self.btn_preview_glv_mask)
        adv_layout.addWidget(self.wgt_glv_controls)

        adv_layout.addSpacing(10)

        # ROI-Match controls (shown only when ROI-Match mode is selected)
        self.wgt_roi_match_controls = QtWidgets.QWidget()
        roi_match_layout = QtWidgets.QVBoxLayout(self.wgt_roi_match_controls)
        roi_match_layout.setContentsMargins(0, 0, 0, 0)
        roi_match_layout.setSpacing(4)

        roi_btn_row = QtWidgets.QHBoxLayout()
        self.btn_pick_roi = QtWidgets.QPushButton("Open ROI Manager")
        self.btn_pick_roi.setToolTip(
            "Define ROIs in ROI Manager. ROI-Match uses ROI bounding-box means\n"
            "from those ROIs to calibrate the compare scaling coefficient."
        )
        self.btn_pick_roi.setFixedWidth(140)
        roi_btn_row.addWidget(self.btn_pick_roi)
        roi_btn_row.addStretch()
        roi_match_layout.addLayout(roi_btn_row)

        self.lbl_roi_info = QtWidgets.QLabel("ROI source: ROI Manager")
        self.lbl_roi_info.setStyleSheet(
            f"color: {UI_TEXT_SECONDARY}; font-family: {Typography.FONT_FAMILY_MONO};"
            f" font-size: {Typography.FONT_SIZE_SMALL};"
        )
        roi_match_layout.addWidget(self.lbl_roi_info)

        self.lbl_roi_alpha = QtWidgets.QLabel("")
        self.lbl_roi_alpha.setStyleSheet(
            f"color: {UI_SUCCESS}; font-family: {Typography.FONT_FAMILY_MONO};"
            f" font-size: {Typography.FONT_SIZE_SMALL}; font-weight: {Typography.FONT_WEIGHT_SEMIBOLD};"
        )
        roi_match_layout.addWidget(self.lbl_roi_alpha)

        self.chk_roi_abs_diff = QtWidgets.QCheckBox("Use |diff| instead of keep-direction")
        self.chk_roi_abs_diff.setChecked(False)
        self.chk_roi_abs_diff.setToolTip(
            "Default (unchecked): Enhanced = clip(Base - α·Compare, 0) — highlights\n"
            "where Base is brighter than scaled Compare (defect path).\n"
            "Checked: Enhanced = |Base - α·Compare| — symmetric edges, legacy behavior."
        )
        roi_match_layout.addWidget(self.chk_roi_abs_diff)

        adv_layout.addWidget(self.wgt_roi_match_controls)

        self.cmb_normalize_mode.currentIndexChanged.connect(self._on_normalize_mode_changed)
        self._on_normalize_mode_changed()  # set initial visibility

        self.grp_advanced.setVisible(False)
        op_layout.addWidget(self.grp_advanced)

        operation_card_layout.addWidget(self.grp_op)

        # --- ALIGNMENT section ---
        card_alignment, alignment_card_layout = _make_sidebar_card("ALIGNMENT")
        left_layout.addWidget(card_alignment)

        self.grp_align = QtWidgets.QWidget()
        align_layout = QtWidgets.QVBoxLayout(self.grp_align)
        align_layout.setContentsMargins(0, 4, 0, 0)
        align_layout.setSpacing(4)

        align_method_row = QtWidgets.QVBoxLayout()
        align_method_row.setContentsMargins(0, 0, 0, 0)
        align_method_row.setSpacing(6)
        align_method_row.addWidget(QtWidgets.QLabel("Method"))
        self.cmb_align_method = QtWidgets.QComboBox()
        self.cmb_align_method.addItems(["Phase (robust)", "NCC (brute force)"])
        align_method_row.addWidget(self.cmb_align_method)
        align_layout.addLayout(align_method_row)

        snr_win_row = QtWidgets.QVBoxLayout()
        snr_win_row.setContentsMargins(0, 0, 0, 0)
        snr_win_row.setSpacing(6)
        snr_win_row.addWidget(QtWidgets.QLabel("SNR Window"))
        self.spn_snr_window = QtWidgets.QSpinBox()
        self.spn_snr_window.setRange(7, 127)
        self.spn_snr_window.setSingleStep(2)
        self.spn_snr_window.setValue(31)
        self.spn_snr_window.setMinimumWidth(96)
        self.spn_snr_window.setToolTip(
            "Box-filter window size for Z-Map SNR calculation (odd, \u22657).\n"
            "Larger values \u2192 smoother map.\n"
            "Recommended: 15 for ~512 px, 31 for ~1000 px, 63 for >2000 px."
        )
        snr_win_row.addWidget(self.spn_snr_window)
        align_layout.addLayout(snr_win_row)
        alignment_card_layout.addWidget(self.grp_align)

        left_layout.addStretch()

        # ─── Step 3: ROI + Compute ───────────────────────────────────────────
        step3_hdr, self._step3_circle, self._step3_status = _make_step_header(
            "③", "ROI  +  Compute", optional=True
        )
        left_layout.addWidget(step3_hdr)

        self.btn_roi_manager.setProperty("role", "secondary")
        self.btn_compute.setProperty("role", "primary")
        left_layout.addWidget(self.btn_roi_manager)

        self.lbl_roi_status = QtWidgets.QLabel("No ROIs — analysis will be skipped")
        self.lbl_roi_status.setAlignment(Qt.AlignCenter)
        self.lbl_roi_status.setWordWrap(True)
        self.lbl_roi_status.setStyleSheet(
            "color: #9CA3AF; font-size: 11px; border: none; background: transparent;"
        )
        left_layout.addWidget(self.lbl_roi_status)
        left_layout.addWidget(self.btn_compute)

        # P1-3: Space bar shortcut hint
        self.lbl_space_hint = QtWidgets.QLabel("Press  Space  to compute")
        self.lbl_space_hint.setAlignment(Qt.AlignCenter)
        self.lbl_space_hint.setStyleSheet(
            "color: #9CA3AF; font-size: 10px; border: none; background: transparent; padding-top: 2px;"
        )
        left_layout.addWidget(self.lbl_space_hint)

        self.left_panel_scroll = QtWidgets.QScrollArea()
        self.left_panel_scroll.setObjectName("LeftPanelScroll")
        self.left_panel_scroll.setWidgetResizable(True)
        self.left_panel_scroll.setHorizontalScrollBarPolicy(Qt.ScrollBarAlwaysOff)
        self.left_panel_scroll.setVerticalScrollBarPolicy(Qt.ScrollBarAsNeeded)
        self.left_panel_scroll.setFrameShape(QtWidgets.QFrame.NoFrame)
        self.left_panel_scroll.setMinimumWidth(self._sidebar_min_width)
        self.left_panel_scroll.setMaximumWidth(self._sidebar_max_width)
        self.left_panel_scroll.setSizePolicy(QtWidgets.QSizePolicy.Preferred, QtWidgets.QSizePolicy.Expanding)
        self.left_panel_scroll.setWidget(left_panel)
        content_layout.addWidget(self.left_panel_scroll, stretch=0)

        # ================================================================
        # RIGHT PANEL: Viewer + Controls
        # ================================================================
        right_panel = QtWidgets.QWidget()
        right_layout = QtWidgets.QVBoxLayout(right_panel)
        right_layout.setContentsMargins(0, 0, 0, 0)
        right_layout.setSpacing(6)

        # ── LEFT SECTION (input preview): control bar + stacked viewer ───────
        _viewer_ctrl_style = f"""
            QFrame#ViewerControlBar {{
                background-color: {UI_BG_PANEL};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.SM};
            }}
        """

        left_section = QtWidgets.QVBoxLayout()
        left_section.setSpacing(4)

        # Left control bar: Base | Compare  +  Split View toggle + slider
        left_ctrl_bar = QtWidgets.QFrame()
        left_ctrl_bar.setObjectName("ViewerControlBar")
        left_ctrl_bar.setStyleSheet(_viewer_ctrl_style)
        left_ctrl_layout = QtWidgets.QHBoxLayout(left_ctrl_bar)
        left_ctrl_layout.setContentsMargins(8, 4, 8, 4)
        left_ctrl_layout.setSpacing(0)

        self.btn_mode_base = QtWidgets.QPushButton("Base")
        self.btn_mode_base.setCheckable(True)
        self.btn_mode_base.setProperty("viewerMode", True)
        self.btn_mode_base.setObjectName("ViewerModeFirst")
        self.btn_mode_compare = QtWidgets.QPushButton("Compare")
        self.btn_mode_compare.setCheckable(True)
        self.btn_mode_compare.setProperty("viewerMode", True)
        self.btn_mode_compare.setObjectName("ViewerModeLast")

        for btn in (self.btn_mode_base, self.btn_mode_compare):
            btn.setFixedHeight(30)
            left_ctrl_layout.addWidget(btn)

        # Compare mode is not meaningful before compute (no aligned compare yet).
        self.btn_mode_compare.setVisible(False)

        left_ctrl_layout.addSpacing(12)

        # Split View toggle placed in left control bar (only affects left viewer)
        self.btn_split_view = QtWidgets.QPushButton("Split View")
        self.btn_split_view.setObjectName("ViewerToolToggle")
        self.btn_split_view.setCheckable(True)
        self.btn_split_view.setChecked(False)
        self.btn_split_view.setToolTip(
            "OFF: Magnifier mode\n"
            "ON:  Split-view Base vs Aligned Compare (drag divider or use slider)"
        )
        self.btn_split_view.setVisible(False)   # hidden until compute
        left_ctrl_layout.addWidget(self.btn_split_view)

        # Slider for split view divider
        self.blend_slider_widget = QtWidgets.QWidget()
        blend_slider_layout = QtWidgets.QHBoxLayout(self.blend_slider_widget)
        blend_slider_layout.setContentsMargins(4, 0, 0, 0)
        self.slider_blend = QtWidgets.QSlider(Qt.Horizontal)
        self.slider_blend.setRange(0, 100)
        self.slider_blend.setValue(50)
        blend_slider_layout.addWidget(self.slider_blend, 1)
        self.lbl_blend_value = QtWidgets.QLabel("50%")
        self.lbl_blend_value.setStyleSheet(
            f"font-family: {Typography.FONT_FAMILY_MONO}; font-size: {Typography.FONT_SIZE_SMALL};")
        blend_slider_layout.addWidget(self.lbl_blend_value)
        self.blend_slider_widget.setVisible(False)
        left_ctrl_layout.addWidget(self.blend_slider_widget, 1)

        left_ctrl_layout.addStretch(1)

        left_section.addWidget(left_ctrl_bar)

        # Left stacked viewer: magnifier (page 0) + split view (page 1)
        self.stk_blend = QtWidgets.QStackedWidget()
        self.stk_blend.setSizePolicy(
            QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Expanding
        )
        self.img_base_mag = SyncZoomImageWidget("Base")
        self.img_base_mag.setSizePolicy(
            QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Expanding
        )
        self.stk_blend.addWidget(self.img_base_mag)  # index 0 magnifier

        self.img_blend = SplitViewWidget()
        self.img_blend.setSizePolicy(
            QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Expanding
        )
        self.stk_blend.addWidget(self.img_blend)  # index 1 split view
        self.stk_blend.setCurrentIndex(0)
        left_section.addWidget(self.stk_blend, stretch=1)

        # ── RIGHT SECTION (result view): control bar + difference map ─────────
        # Wrapped in a QWidget so it can be hidden in pre-compute state
        self.wgt_diff_section = QtWidgets.QWidget()
        right_section = QtWidgets.QVBoxLayout(self.wgt_diff_section)
        right_section.setSpacing(4)
        right_section.setContentsMargins(0, 0, 0, 0)

        # Right control bar: Diff | Z-Map  +  Range  Normalize  Colormap
        right_ctrl_bar = QtWidgets.QFrame()
        right_ctrl_bar.setObjectName("ViewerControlBar")
        right_ctrl_bar.setStyleSheet(_viewer_ctrl_style)
        right_ctrl_layout = QtWidgets.QHBoxLayout(right_ctrl_bar)
        right_ctrl_layout.setContentsMargins(8, 4, 8, 4)
        right_ctrl_layout.setSpacing(0)

        self.btn_mode_diff = QtWidgets.QPushButton("Diff")
        self.btn_mode_diff.setCheckable(True)
        self.btn_mode_diff.setChecked(True)
        self.btn_mode_diff.setProperty("viewerMode", True)
        self.btn_mode_diff.setObjectName("ViewerModeFirst")
        self.btn_mode_topo = QtWidgets.QPushButton("Z-Map")
        self.btn_mode_topo.setCheckable(True)
        self.btn_mode_topo.setProperty("viewerMode", True)
        self.btn_mode_topo.setObjectName("ViewerModeLast")
        # Backward-compat alias
        self.btn_mode_zmap = self.btn_mode_topo

        for btn in (self.btn_mode_diff, self.btn_mode_topo):
            btn.setFixedHeight(30)
            right_ctrl_layout.addWidget(btn)

        right_ctrl_layout.addSpacing(10)

        _sep_v = QtWidgets.QFrame()
        _sep_v.setFrameShape(QtWidgets.QFrame.VLine)
        _sep_v.setStyleSheet(f"color: {UI_BORDER}; max-width: 1px; margin: 2px 8px;")
        right_ctrl_layout.addWidget(_sep_v)

        lbl_display_group = QtWidgets.QLabel("DISPLAY")
        lbl_display_group.setProperty("toolbarGroupLabel", True)
        right_ctrl_layout.addWidget(lbl_display_group)
        right_ctrl_layout.addSpacing(8)

        self.btn_show_norm_compare = QtWidgets.QPushButton("Normalize")
        self.btn_show_norm_compare.setFixedWidth(100)
        self.btn_show_norm_compare.setToolTip("Preview normalization effect")
        right_ctrl_layout.addWidget(self.btn_show_norm_compare)

        right_ctrl_layout.addSpacing(8)
        lbl_colormap = QtWidgets.QLabel("Colormap")
        lbl_colormap.setProperty("toolbarLabel", True)
        right_ctrl_layout.addWidget(lbl_colormap)
        right_ctrl_layout.addSpacing(4)
        self.cmb_colormap = QtWidgets.QComboBox()
        self.cmb_colormap.addItems(["Grayscale", "JET", "Hot", "Inferno", "Viridis"])
        self.cmb_colormap.setFixedWidth(104)
        right_ctrl_layout.addWidget(self.cmb_colormap)

        right_ctrl_layout.addSpacing(8)
        lbl_range = QtWidgets.QLabel("Range")
        lbl_range.setProperty("toolbarLabel", True)
        right_ctrl_layout.addWidget(lbl_range)
        right_ctrl_layout.addSpacing(4)
        self.cmb_range = QtWidgets.QComboBox()
        self.cmb_range.addItems(["Auto", "Zero-centered", "P1-P99", "P0.5-P99.5"])
        self.cmb_range.setFixedWidth(128)
        self.cmb_range.setToolTip("Control how difference values are scaled for display")
        right_ctrl_layout.addWidget(self.cmb_range)

        right_ctrl_layout.addSpacing(10)
        self.chk_roi_view = QtWidgets.QCheckBox("ROI View")
        self.chk_roi_view.setChecked(True)
        self.chk_roi_view.setToolTip("Show or hide ROI overlays on Base and Diff viewers")
        right_ctrl_layout.addWidget(self.chk_roi_view)

        right_ctrl_layout.addStretch(1)

        right_section.addWidget(right_ctrl_bar)

        # Right: Difference Map
        self.img_diff = SyncZoomImageWidget("Difference Map")
        self.img_diff.setSizePolicy(
            QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Expanding
        )
        right_section.addWidget(self.img_diff, stretch=1)

        # ── IMAGE AREA: left section + right section side by side ─────────────
        # Back button — only visible post-compute, restores pre-compute state
        _top_btn_row = QtWidgets.QHBoxLayout()
        self.btn_back_to_settings = QtWidgets.QPushButton("← Back")
        self.btn_back_to_settings.setToolTip("Return to Settings & pre-compute view")
        self.btn_back_to_settings.setFixedHeight(28)
        self.btn_back_to_settings.setVisible(False)
        _top_btn_row.addWidget(self.btn_back_to_settings)
        _top_btn_row.addStretch()
        right_layout.addLayout(_top_btn_row)

        image_row = QtWidgets.QHBoxLayout()
        image_row.setSpacing(8)
        image_row.addLayout(left_section, stretch=1)
        image_row.addWidget(self.wgt_diff_section, stretch=1)
        image_row.setStretch(0, 1)
        image_row.setStretch(1, 1)
        self.wgt_diff_section.setVisible(False)  # hidden until compute
        right_layout.addLayout(image_row, stretch=4)

        # Hidden backward-compat widgets (not shown)
        self.lbl_blend_info = QtWidgets.QLabel("")
        self.lbl_blend_info.setVisible(False)
        self.btn_mode_blend = QtWidgets.QPushButton("Blend")
        self.btn_mode_blend.setCheckable(True)
        self.btn_mode_blend.setVisible(False)

        right_layout.addSpacing(4)

        # === BOTTOM: Analysis (merged) | Histogram (under Difference Map) ===
        self.wgt_bottom_row = QtWidgets.QWidget()
        self.wgt_bottom_row.setVisible(False)   # hidden until compute
        bottom_row = QtWidgets.QHBoxLayout(self.wgt_bottom_row)
        bottom_row.setContentsMargins(0, 0, 0, 0)
        bottom_row.setSpacing(8)

        # Card 1 (LEFT, ratio 1): Alignment — compact alignment-quality metrics only
        self.align_panel = AlignmentPanelWidget()
        bottom_row.addWidget(self.align_panel, 1)

        # Card 2 (CENTER, ratio 2): Diff / ROI Analysis — ROI quantification summary
        self.diff_roi_panel = DiffROIAnalysisPanelWidget()
        self.diff_roi_panel.open_roi_manager_requested.connect(self._on_open_roi_manager)
        self.diff_roi_panel.roi_details_requested.connect(self._on_show_roi_profile_dialog)
        bottom_row.addWidget(self.diff_roi_panel, 2)

        # Legacy StatisticsWidget — kept in memory for backward-compat, not shown in layout
        self.stats_widget = StatisticsWidget()
        self.stats_widget.setVisible(False)

        # Card 3 (RIGHT, ratio 3): Histogram (wider, positioned under the Difference Map)
        hist_card = QtWidgets.QFrame()
        hist_card.setObjectName("BottomCard")
        hist_layout = QtWidgets.QVBoxLayout(hist_card)
        hist_layout.setContentsMargins(12, 10, 12, 8)
        hist_layout.setSpacing(3)
        lbl_hist_title = QtWidgets.QLabel("Histogram")
        lbl_hist_title.setStyleSheet(
            f"font-weight: {Typography.FONT_WEIGHT_BOLD}; font-size: {Typography.FONT_SIZE_BODY};"
            f" color: {UI_TEXT}; border: none; background: transparent;"
        )
        hist_layout.addWidget(lbl_hist_title)
        hist_hint = QtWidgets.QLabel("Click once to set low, click again to set high")
        hist_hint.setStyleSheet(
            f"color: {UI_TEXT_MUTED}; font-size: {Typography.FONT_SIZE_CAPTION};"
            f" border: none; background: transparent;"
        )
        hist_layout.addWidget(hist_hint)
        self.histogram_canvas = HistogramCanvas()
        self.histogram_canvas.setFixedHeight(186)
        hist_layout.addWidget(self.histogram_canvas)
        hist_ctrl_row = QtWidgets.QHBoxLayout()
        self.lbl_hist_range = QtWidgets.QLabel("Range: \u2014")
        self.lbl_hist_range.setStyleSheet(
            f"color: {UI_TEXT_SECONDARY}; font-family: {Typography.FONT_FAMILY_MONO};"
            f" font-size: {Typography.FONT_SIZE_SMALL}; border: none;"
        )
        hist_ctrl_row.addWidget(self.lbl_hist_range, 1)
        self.btn_clear_hist_range = QtWidgets.QPushButton("Clear Range")
        self.btn_clear_hist_range.setFixedWidth(96)
        self.btn_clear_hist_range.setFixedHeight(24)
        self.btn_clear_hist_range.setEnabled(False)
        hist_ctrl_row.addWidget(self.btn_clear_hist_range)
        hist_layout.addLayout(hist_ctrl_row)
        bottom_row.addWidget(hist_card, 3)

        # Hidden AlignmentScoreWidget — kept for backward-compat with update logic
        self.align_score_widget = AlignmentScoreWidget()
        self.align_score_widget.setVisible(False)

        right_layout.addWidget(self.wgt_bottom_row)

        # === QUADRANT FUSION RIGHT PANEL (Page 1) ===
        qf_right_panel = QtWidgets.QWidget()
        qf_right_layout = QtWidgets.QVBoxLayout(qf_right_panel)
        qf_right_layout.setContentsMargins(0, 0, 0, 0)
        qf_right_layout.setSpacing(8)

        # Top: output type toggle buttons
        qf_top_toolbar = QtWidgets.QFrame()
        qf_top_toolbar.setObjectName("TopToolbar")
        qf_top_row = QtWidgets.QHBoxLayout(qf_top_toolbar)
        qf_top_row.setContentsMargins(10, 8, 10, 8)
        qf_top_row.setSpacing(8)
        lbl_output_view = QtWidgets.QLabel("Output View")
        lbl_output_view.setProperty("toolbarLabel", True)
        qf_top_row.addWidget(lbl_output_view)
        self.btn_qf_show_bse = QtWidgets.QPushButton("BSE Enhanced")
        self.btn_qf_show_bse.setCheckable(True)
        self.btn_qf_show_bse.setChecked(True)
        self.btn_qf_show_bse.setFixedWidth(128)
        self.btn_qf_show_bse.setProperty("toolbarToggle", True)

        self.btn_qf_show_topo = QtWidgets.QPushButton("Topography")
        self.btn_qf_show_topo.setCheckable(True)
        self.btn_qf_show_topo.setFixedWidth(128)
        self.btn_qf_show_topo.setProperty("toolbarToggle", True)

        self.btn_qf_show_comp = QtWidgets.QPushButton("Composite")
        self.btn_qf_show_comp.setCheckable(True)
        self.btn_qf_show_comp.setFixedWidth(128)
        self.btn_qf_show_comp.setProperty("toolbarToggle", True)

        qf_top_row.addWidget(self.btn_qf_show_bse)
        qf_top_row.addWidget(self.btn_qf_show_topo)
        qf_top_row.addWidget(self.btn_qf_show_comp)

        qf_top_row.addSpacing(20)
        lbl_qf_colormap = QtWidgets.QLabel("Colormap")
        lbl_qf_colormap.setProperty("toolbarLabel", True)
        qf_top_row.addWidget(lbl_qf_colormap)
        self.cmb_qf_colormap = QtWidgets.QComboBox()
        self.cmb_qf_colormap.addItems(["Grayscale", "JET", "Hot", "Inferno", "Viridis"])
        self.cmb_qf_colormap.setFixedWidth(106)
        qf_top_row.addWidget(self.cmb_qf_colormap)
        qf_top_row.addStretch()
        qf_right_layout.addWidget(qf_top_toolbar)

        # 3-column image viewers: Illuminator | Main Output | Topo Reference
        qf_image_row = QtWidgets.QHBoxLayout()

        # Left: Illuminator (original reference)
        qf_illum_grp = QtWidgets.QGroupBox("Illuminator (Original)")
        qf_illum_lay = QtWidgets.QVBoxLayout(qf_illum_grp)
        qf_illum_lay.setContentsMargins(6, 12, 6, 6)
        self.img_qf_illum = SyncZoomImageWidget("Illuminator")
        self.img_qf_illum.setSizePolicy(
            QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Expanding
        )
        qf_illum_lay.addWidget(self.img_qf_illum)
        qf_image_row.addWidget(qf_illum_grp, 1)

        # Center: Main output (BSE / Topo / Composite, switched by toggle)
        self.qf_main_grp = QtWidgets.QGroupBox("BSE Enhanced (Illum − α·Topo)")
        qf_main_lay = QtWidgets.QVBoxLayout(self.qf_main_grp)
        qf_main_lay.setContentsMargins(6, 12, 6, 6)
        self.img_qf_main = SyncZoomImageWidget("Fusion Output")
        self.img_qf_main.setSizePolicy(
            QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Expanding
        )
        qf_main_lay.addWidget(self.img_qf_main)
        qf_image_row.addWidget(self.qf_main_grp, 1)

        # Right: Secondary reference (Topo when main=BSE, BSE when main=Topo, etc.)
        self.qf_ref_grp = QtWidgets.QGroupBox("Topography")
        qf_ref_lay = QtWidgets.QVBoxLayout(self.qf_ref_grp)
        qf_ref_lay.setContentsMargins(6, 12, 6, 6)
        self.img_qf_ref = SyncZoomImageWidget("Reference")
        self.img_qf_ref.setSizePolicy(
            QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Expanding
        )
        qf_ref_lay.addWidget(self.img_qf_ref)
        qf_image_row.addWidget(self.qf_ref_grp, 1)

        qf_right_layout.addLayout(qf_image_row, stretch=4)

        # Link magnifiers for synchronized zoom
        self.img_qf_illum.linkPartner(self.img_qf_main)
        self.img_qf_main.linkPartner(self.img_qf_illum)

        # Sync cursors between QF viewers
        self.img_qf_illum.cursor_moved.connect(
            lambda nx, ny: (self.img_qf_main.setCursorPos(nx, ny),
                            self.img_qf_ref.setCursorPos(nx, ny))
        )
        self.img_qf_illum.cursor_left.connect(
            lambda: (self.img_qf_main.clearCursor(), self.img_qf_ref.clearCursor())
        )
        self.img_qf_main.cursor_moved.connect(
            lambda nx, ny: (self.img_qf_illum.setCursorPos(nx, ny),
                            self.img_qf_ref.setCursorPos(nx, ny))
        )
        self.img_qf_main.cursor_left.connect(
            lambda: (self.img_qf_illum.clearCursor(), self.img_qf_ref.clearCursor())
        )
        self.img_qf_ref.cursor_moved.connect(
            lambda nx, ny: (self.img_qf_illum.setCursorPos(nx, ny),
                            self.img_qf_main.setCursorPos(nx, ny))
        )
        self.img_qf_ref.cursor_left.connect(
            lambda: (self.img_qf_illum.clearCursor(), self.img_qf_main.clearCursor())
        )

        # QF Fusion Info bar
        self.lbl_qf_result_info = QtWidgets.QLabel("No fusion result yet")
        self.lbl_qf_result_info.setAlignment(Qt.AlignCenter)
        self.lbl_qf_result_info.setStyleSheet(f"""
            QLabel {{
                color: {UI_TEXT};
                background-color: {UI_BG_CARD};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.LG};
                padding: {Spacing.BUTTON_PADDING};
                font-size: {Typography.FONT_SIZE_BODY};
            }}
        """)
        qf_right_layout.addWidget(self.lbl_qf_result_info)

        # Bottom row: Histogram + Fusion Stats
        qf_bottom_row = QtWidgets.QHBoxLayout()

        # Histogram
        qf_hist_grp = QtWidgets.QGroupBox("Fusion Histogram")
        qf_hist_lay = QtWidgets.QVBoxLayout(qf_hist_grp)
        qf_hist_lay.setSpacing(4)
        self.qf_histogram_canvas = HistogramCanvas()
        self.qf_histogram_canvas.setFixedHeight(160)
        qf_hist_lay.addWidget(self.qf_histogram_canvas)
        qf_bottom_row.addWidget(qf_hist_grp, 1)

        # Fusion Stats
        qf_stats_grp = QtWidgets.QGroupBox("Fusion Parameters")
        qf_stats_grp.setStyleSheet(f"""
            QGroupBox {{
                background-color: {UI_BG_CARD};
                border: 1px solid {UI_BORDER};
                border-radius: {BorderRadius.MD};
            }}
        """)
        qf_stats_lay = QtWidgets.QGridLayout(qf_stats_grp)
        qf_stats_lay.setContentsMargins(12, 16, 12, 12)
        qf_stats_lay.setSpacing(6)

        _mono_style = (
            f"font-family: {Typography.FONT_FAMILY_MONO};"
            f" font-size: {Typography.FONT_SIZE_SMALL};"
            f" color: {UI_TEXT}; border: none;"
        )
        _label_style = (
            f"color: {UI_TEXT_SECONDARY}; font-size: {Typography.FONT_SIZE_CAPTION}; border: none;"
        )

        self._qf_stat_labels: Dict[str, QtWidgets.QLabel] = {}
        for i, (key, label_text) in enumerate([
            ("alpha", "Alpha (α):"),
            ("beta", "Beta (β):"),
            ("sigma", "Topo Sigma:"),
            ("output", "Output Type:"),
            ("mean", "Mean:"),
            ("std", "Std Dev:"),
            ("dim", "Dimensions:"),
        ]):
            lbl = QtWidgets.QLabel(label_text)
            lbl.setStyleSheet(_label_style)
            val = QtWidgets.QLabel("--")
            val.setStyleSheet(_mono_style)
            qf_stats_lay.addWidget(lbl, i, 0)
            qf_stats_lay.addWidget(val, i, 1)
            self._qf_stat_labels[key] = val

        qf_bottom_row.addWidget(qf_stats_grp, 1)
        qf_right_layout.addLayout(qf_bottom_row)

        # === RIGHT PANEL STACKED WIDGET ===
        self.stk_right_panel = QtWidgets.QStackedWidget()
        self.stk_right_panel.addWidget(right_panel)  # index 0 = Standard
        self.stk_right_panel.addWidget(qf_right_panel)  # index 1 = Quadrant Fusion
        self.stk_right_panel.setCurrentIndex(0)

        content_layout.addWidget(self.stk_right_panel, stretch=1)

        # P1-2: ROI Manager slide-in side panel (hidden by default)
        self.roi_side_panel = QtWidgets.QFrame()
        self.roi_side_panel.setObjectName("ROISidePanel")
        self.roi_side_panel.setStyleSheet(f"""
            QFrame#ROISidePanel {{
                background-color: {UI_BG_PANEL};
                border-left: 2px solid {UI_PRIMARY};
                border-radius: 0px;
            }}
        """)
        self.roi_side_panel.setFixedWidth(480)  # match original dialog min-width
        self.roi_side_panel.setSizePolicy(
            QtWidgets.QSizePolicy.Fixed, QtWidgets.QSizePolicy.Expanding
        )
        _rsp_layout = QtWidgets.QVBoxLayout(self.roi_side_panel)
        _rsp_layout.setContentsMargins(0, 0, 0, 0)
        _rsp_layout.setSpacing(0)

        # Header row: title + close button
        _rsp_hdr = QtWidgets.QFrame()
        _rsp_hdr.setStyleSheet(
            f"QFrame {{ background-color: {UI_PRIMARY}; border: none; }}"
        )
        _rsp_hdr_lay = QtWidgets.QHBoxLayout(_rsp_hdr)
        _rsp_hdr_lay.setContentsMargins(12, 6, 8, 6)
        _rsp_hdr_lay.setSpacing(8)
        _rsp_title = QtWidgets.QLabel("📐  ROI Manager")
        _rsp_title.setStyleSheet(
            "color: #111827; font-weight: 700; font-size: 13px; background: transparent; border: none;"
        )
        _rsp_hdr_lay.addWidget(_rsp_title)
        _rsp_hdr_lay.addStretch()
        self.btn_close_roi_panel = QtWidgets.QPushButton("✕")
        self.btn_close_roi_panel.setFixedSize(24, 24)
        self.btn_close_roi_panel.setFlat(True)
        self.btn_close_roi_panel.setStyleSheet(
            "QPushButton { color: #111827; font-weight: 700; background: transparent; border: none; }"
            "QPushButton:hover { background-color: rgba(0,0,0,0.12); border-radius: 4px; }"
        )
        self.btn_close_roi_panel.clicked.connect(self._close_roi_side_panel)
        _rsp_hdr_lay.addWidget(self.btn_close_roi_panel)
        _rsp_layout.addWidget(_rsp_hdr)

        # Scroll area for ROI manager content (filled lazily on first open)
        self._roi_side_scroll = QtWidgets.QScrollArea()
        self._roi_side_scroll.setWidgetResizable(True)
        self._roi_side_scroll.setFrameShape(QtWidgets.QFrame.NoFrame)
        self._roi_side_scroll.setHorizontalScrollBarPolicy(Qt.ScrollBarAlwaysOff)
        _rsp_layout.addWidget(self._roi_side_scroll, stretch=1)

        self._roi_panel_populated = False  # filled on first open
        self.roi_side_panel.setVisible(False)
        content_layout.addWidget(self.roi_side_panel, stretch=0)

        content_layout.setStretch(0, 0)
        content_layout.setStretch(1, 1)
        content_layout.setStretch(2, 0)
        outer_layout.addLayout(content_layout, stretch=1)

        # ── Embedded progress banner (shown during compute) ───────────────
        self.wgt_progress_banner = QtWidgets.QWidget()
        self.wgt_progress_banner.setObjectName("ProgressBanner")
        self.wgt_progress_banner.setStyleSheet(
            "QWidget#ProgressBanner { background-color: #FFFBEB; border-top: 1px solid #F59E0B; }"
        )
        _pb_layout = QtWidgets.QVBoxLayout(self.wgt_progress_banner)
        _pb_layout.setContentsMargins(16, 8, 16, 8)
        _pb_layout.setSpacing(4)
        # Top row: progress text + abort button
        _pb_top_row = QtWidgets.QHBoxLayout()
        _pb_top_row.setSpacing(8)
        self.lbl_progress_text = QtWidgets.QLabel("Computing…")
        self.lbl_progress_text.setStyleSheet(
            "color: #92400E; font-size: 12px; font-weight: 600; border: none;"
        )
        _pb_top_row.addWidget(self.lbl_progress_text, stretch=1)
        self.btn_abort_compute = QtWidgets.QPushButton("■  Abort")
        self.btn_abort_compute.setFixedHeight(24)
        self.btn_abort_compute.setCursor(QtGui.QCursor(Qt.PointingHandCursor))
        self.btn_abort_compute.setStyleSheet(
            "QPushButton { background-color: #EF4444; color: #FFFFFF; border: none; "
            "border-radius: 4px; padding: 2px 12px; font-size: 11px; font-weight: 700; }"
            "QPushButton:hover { background-color: #DC2626; }"
            "QPushButton:pressed { background-color: #B91C1C; }"
        )
        _pb_top_row.addWidget(self.btn_abort_compute)
        _pb_layout.addLayout(_pb_top_row)
        self.progress_bar = QtWidgets.QProgressBar()
        self.progress_bar.setTextVisible(True)
        self.progress_bar.setFormat("%v / %m  (%p%)")
        self.progress_bar.setFixedHeight(18)
        self.progress_bar.setStyleSheet("""
            QProgressBar {
                border: 1px solid #F59E0B;
                border-radius: 4px;
                background: #FEF3C7;
                text-align: center;
                color: #92400E;
                font-size: 11px;
            }
            QProgressBar::chunk {
                background-color: #F59E0B;
                border-radius: 3px;
            }
        """)
        _pb_layout.addWidget(self.progress_bar)
        self.wgt_progress_banner.setVisible(False)
        outer_layout.addWidget(self.wgt_progress_banner)

        self._compute_lock_targets = (
            self.toolbar,
            self.lbl_result_info,
            self.left_panel_scroll,
            self.stk_right_panel,
            self.roi_side_panel,
        )
        self._apply_responsive_sidebar_layout()

    def resizeEvent(self, event: QtGui.QResizeEvent):
        super().resizeEvent(event)
        self._apply_responsive_sidebar_layout()

    def _apply_responsive_sidebar_layout(self):
        """Apply responsive sidebar sizing and control reflow based on current width."""
        if not hasattr(self, "left_panel"):
            return

        total_w = max(1, self.width())
        # Keep sidebar readable and bounded; viewer gets most extra space.
        target_w = int(total_w * 0.27)
        target_w = max(self._sidebar_min_width, min(self._sidebar_max_width, target_w))
        self.left_panel_scroll.setMinimumWidth(self._sidebar_min_width)
        self.left_panel_scroll.setMaximumWidth(target_w)
        self.left_panel.setMinimumWidth(self._sidebar_min_width)
        self.left_panel.setMaximumWidth(target_w)

        sidebar_w = self.left_panel.width()
        compact = sidebar_w < 305

        # INPUT compare header: split title and utility buttons into two rows in compact mode.
        if hasattr(self, "compare_hdr_row"):
            self.compare_hdr_row.setDirection(
                QtWidgets.QBoxLayout.TopToBottom if compact else QtWidgets.QBoxLayout.LeftToRight
            )

        # Invert options: avoid squeezing 3 checkboxes in one row.
        if hasattr(self, "invert_row"):
            self.invert_row.setDirection(
                QtWidgets.QBoxLayout.TopToBottom if compact else QtWidgets.QBoxLayout.LeftToRight
            )

    def _set_widget_locked_appearance(self, widget: QtWidgets.QWidget, locked: bool) -> None:
        if widget is None:
            return
        effect = self._compute_lock_effects.get(widget)
        if locked:
            if effect is None:
                effect = QtWidgets.QGraphicsOpacityEffect(widget)
                self._compute_lock_effects[widget] = effect
            effect.setOpacity(0.38)
            widget.setGraphicsEffect(effect)
        else:
            widget.setGraphicsEffect(None)

    def _set_compute_busy(self, busy: bool) -> None:
        if busy and self.roi_side_panel.isVisible():
            self._close_roi_side_panel()
        for widget in self._compute_lock_targets:
            widget.setEnabled(not busy)
            self._set_widget_locked_appearance(widget, busy)
        self.btn_abort_compute.setVisible(busy)
        self.btn_abort_compute.setEnabled(busy and self._compute_worker is not None)
        if not busy:
            self.btn_abort_compute.setText("Abort")

    def _on_abort_compute(self) -> None:
        if self._compute_worker is None or self._compute_abort_requested:
            return
        self._compute_abort_requested = True
        self._compute_worker.abort_requested = True
        self.btn_abort_compute.setEnabled(False)
        self.btn_abort_compute.setText("Aborting...")
        self.lbl_progress_text.setText(
            "Aborting compute after the active pair finishes..."
        )

    def _connect_signals(self):
        """Connect UI signals."""
        self.btn_load_folder.clicked.connect(self._on_load_image_folder)
        self.btn_compute.clicked.connect(self._on_compute)
        self.btn_abort_compute.clicked.connect(self._on_abort_compute)
        self.btn_export.clicked.connect(self._on_export)
        self.btn_select_all.clicked.connect(self._select_all_compare)
        self.btn_select_none.clicked.connect(self._select_none_compare)
        self.cmb_base.currentIndexChanged.connect(self._on_base_changed)
        self.cmb_base.currentIndexChanged.connect(self._update_step_states)
        self.chk_auto_pair.stateChanged.connect(self._on_auto_pair_toggle)
        self.btn_adv_toggle.toggled.connect(self._on_adv_toggle)
        # P1-1: Advanced Settings badge refresh
        self.cmb_normalize_mode.currentIndexChanged.connect(self._update_adv_badge)
        self.cmb_subtract_mode.currentIndexChanged.connect(self._update_adv_badge)
        self.cmb_align_method.currentIndexChanged.connect(self._update_adv_badge)
        self.slider_blend.valueChanged.connect(self._on_blend_change)
        self.histogram_canvas.range_changed.connect(self._on_hist_range_changed)
        self.btn_clear_hist_range.clicked.connect(self._on_clear_hist_range)

        # ROI Manager button
        self.btn_roi_manager.clicked.connect(self._on_open_roi_manager)

        # Back to Settings / Re-compute buttons (post-compute)
        self.btn_back_to_settings.clicked.connect(self._on_back_to_settings)

        # Input Mode selector
        self.cmb_input_mode.currentIndexChanged.connect(self._on_input_mode_changed)

        # Quadrant Fusion ROI, auto-detect, and view toggles
        self.btn_qf_auto_detect.clicked.connect(self._on_qf_auto_detect)
        self.btn_qf_pick_roi.clicked.connect(self._on_qf_pick_roi)
        self.btn_qf_clear_roi.clicked.connect(self._on_qf_clear_roi)
        self.img_base_mag.roi_selected.connect(self._on_roi_selected)
        self.img_qf_illum.roi_selected.connect(self._on_roi_selected)
        self.btn_qf_show_bse.clicked.connect(lambda: self._on_qf_view_toggle("bse_clean"))
        self.btn_qf_show_topo.clicked.connect(lambda: self._on_qf_view_toggle("topo"))
        self.btn_qf_show_comp.clicked.connect(lambda: self._on_qf_view_toggle("composite"))
        self.cmb_qf_colormap.currentIndexChanged.connect(lambda _: self._refresh_qf_display())

        # Operation / navigation
        self.cmb_operation.currentIndexChanged.connect(self._on_operation_changed)
        self.btn_prev_result.clicked.connect(self._on_prev_result)
        self.btn_next_result.clicked.connect(self._on_next_result)

        # Segmented display mode buttons
        self.btn_mode_base.clicked.connect(lambda: self._on_display_mode('base'))
        self.btn_mode_compare.clicked.connect(lambda: self._on_display_mode('compare'))
        self.btn_mode_diff.clicked.connect(lambda: self._on_display_mode('diff'))
        self.btn_mode_topo.clicked.connect(lambda: self._on_display_mode('zmap'))

        # Dynamic range control
        self.cmb_range.currentIndexChanged.connect(self._on_range_changed)
        self.btn_show_norm_compare.clicked.connect(self._on_show_normalized_compare)
        self.chk_roi_view.toggled.connect(self._on_roi_view_toggled)

        # Colormap selector
        self.cmb_colormap.currentIndexChanged.connect(lambda _: self._refresh_diff_display())

        # Swap Base button
        self.btn_swap_base.clicked.connect(self._on_swap_base)

        # ROI-Match: pick / clear
        self.btn_pick_roi.clicked.connect(self._on_open_roi_manager_for_match)

        # Split View toggle
        self.btn_split_view.toggled.connect(self._on_split_view_toggle)

        self.img_base_mag.clicked.connect(self._on_left_viewer_clicked)
        self.img_blend.clicked.connect(self._on_split_viewer_clicked)
        self.img_diff.clicked.connect(self._on_right_viewer_clicked)
        self.img_diff.cursor_moved.connect(self._on_diff_cursor_moved)
        self.img_diff.cursor_left.connect(self._on_diff_cursor_left)
        self.img_base_mag.cursor_moved.connect(self._on_base_mag_cursor_moved)
        self.img_base_mag.cursor_left.connect(self._on_base_mag_cursor_left)

        self._update_sidebar_empty_state(False)
        self._sync_viewer_mode_buttons()

    def _on_load_image_folder(self):
        """Select an image folder and load all supported images."""
        import os

        folder = QtWidgets.QFileDialog.getExistingDirectory(
            self,
            "Select Image Folder for Perspective Combination",
            str(Path.home())
        )
        if not folder:
            return

        extensions = (".tif", ".tiff", ".png", ".jpg", ".jpeg", ".bmp")
        image_files = [
            os.path.join(folder, fn)
            for fn in sorted(os.listdir(folder))
            if fn.lower().endswith(extensions)
        ]

        if not image_files:
            QtWidgets.QMessageBox.warning(
                self,
                "No Images Found",
                f"No image files found in the selected folder.\n\nSupported formats: {', '.join(extensions)}"
            )
            return

        if len(image_files) < 2:
            QtWidgets.QMessageBox.warning(
                self,
                "Not Enough Images",
                "Perspective Combination requires at least 2 images.\n"
                f"Found only {len(image_files)} image(s)."
            )
            return

        self._conditions = [
            EbeamCondition(
                image_path=img_path,
                kev=0.0,
                current_value=0.0,
                current_unit="nA",
                label=os.path.basename(img_path),
                defect_id="",
            )
            for img_path in image_files
        ]
        self._load_images()

    def _on_auto_pair_toggle(self, state):
        """Enable/disable base selection when auto pairing."""
        auto_pair = bool(state)
        self.cmb_base.setEnabled(not auto_pair)
        if auto_pair:
            for chk in self._compare_checkboxes:
                chk.setEnabled(True)
            # Auto pair should evaluate all directional pairs by default.
            self._select_all_compare()
        else:
            self._on_base_changed()

    # ── Input Mode switching ─────────────────────────────────────────────

    def _on_input_mode_changed(self, index: int = None):
        """Toggle between Standard (0) and Quadrant Fusion (1) UI panels."""
        is_qf = self.cmb_input_mode.currentIndex() == 1
        # Left panel: image selection
        self.wgt_standard_select.setVisible(not is_qf)
        self.wgt_quadrant_select.setVisible(is_qf)
        # Left panel: Standard-only groups
        self.grp_op.setVisible(not is_qf)
        self.grp_align.setVisible(not is_qf)
        # Keep using the existing embedded pages; never create/detach viewers here.
        self.stk_right_panel.setCurrentIndex(1 if is_qf else 0)

        if is_qf:
            self.setWindowTitle("Fusi³ — Quadrant Fusion")
        else:
            self.setWindowTitle("Fusi³ — SEM Perspective Combination Tool")
            self._refresh_standard_mode_views()

    def _refresh_standard_mode_views(self):
        """Refresh existing embedded Standard-mode viewers after mode switching."""
        if self._results:
            self._update_current_result()
        else:
            self._show_standard_empty_state()

    def _show_standard_empty_state(self):
        """Show an embedded empty state for Standard mode without extra windows."""
        self._display_mode = 'diff'
        self._left_non_split_mode = 'base'
        self._left_view_mode = 'base'
        if self.btn_split_view.isChecked():
            self.btn_split_view.setChecked(False)
        self.img_base_mag.setImage(None)
        self.img_blend.set_images(None, None)
        self.img_diff.setImage(None)
        self.histogram_canvas.plot_histogram(None, None)
        self._hist_range = None
        self.lbl_hist_range.setText("Range: —")
        self.btn_clear_hist_range.setEnabled(False)
        self.stats_widget.reset()
        self._sync_viewer_mode_buttons()

    def _on_qf_auto_detect(self):
        """Auto-detect Quadrant Fusion detector assignments from filenames."""
        file_labels = list(self._images.keys())
        if not file_labels:
            return
        patterns = {
            "Illuminator": ("illum", "central"),
            "Top": ("top",),
            "Bottom": ("bottom",),
            "Left": ("left",),
            "Right": ("right",),
        }
        for det_name, keywords in patterns.items():
            cmb = self._qf_combos[det_name]
            for i in range(cmb.count()):
                text_lower = cmb.itemText(i).lower()
                if any(kw in text_lower for kw in keywords):
                    cmb.setCurrentIndex(i)
                    break

    def _on_qf_pick_roi(self):
        """Enter ROI drawing mode for Quadrant Fusion alpha fit."""
        if not self._images:
            QtWidgets.QMessageBox.information(
                self, "Quadrant Fusion", "Please load images first."
            )
            return
        # Show Illuminator in QF left viewer and enable ROI drawing on it
        illum_label = self._qf_combos["Illuminator"].currentText()
        if illum_label and illum_label in self._images:
            self.img_qf_illum.setImage(self._images[illum_label])
        self.img_qf_illum.set_roi_mode(True)
        self.lbl_qf_roi_info.setText("ROI: draw on Illuminator image...")

    def _on_qf_clear_roi(self):
        """Clear the Quadrant Fusion ROI."""
        self._qf_roi_rect = None
        self.img_qf_illum.set_active_roi(None)
        self.img_qf_main.set_active_roi(None)
        self.btn_qf_clear_roi.setEnabled(False)
        self.lbl_qf_roi_info.setText("ROI: not set")

    def _on_blend_change(self, value):
        self.lbl_blend_value.setText(f"{value}%")
        self._update_blend_preview()

    def _on_split_view_toggle(self, checked: bool):
        """Switch left panel between magnifier (off) and split-view (on)."""
        if checked:
            self._left_view_mode = 'split'
            self.stk_blend.setCurrentIndex(1)  # show SplitViewWidget
            self.blend_slider_widget.setVisible(True)
        else:
            self._left_view_mode = self._left_non_split_mode
            self.stk_blend.setCurrentIndex(0)  # show magnifier
            self.blend_slider_widget.setVisible(False)
            self._refresh_base_compare_display(self._left_non_split_mode)
        self._sync_viewer_mode_buttons()
        self._update_blend_preview()

    def _on_diff_cursor_moved(self, norm_x: float, norm_y: float):
        """Relay Difference Map cursor position to left-panel magnifier."""
        if (not self.btn_split_view.isChecked()
                and self.img_base_mag._multi_draw_mode == 'idle'
                and not self.img_base_mag._roi_mode):
            self.img_base_mag.setCursorPos(norm_x, norm_y)

    def _on_diff_cursor_left(self):
        """Clear left-panel magnifier when cursor leaves Difference Map."""
        if not self.btn_split_view.isChecked():
            self.img_base_mag.clearCursor()

    def _on_operation_changed(self, index):
        """Show/hide blend coefficients based on operation mode."""
        is_blend = (index == 1)
        self.grp_blend_coef.setVisible(is_blend)

    def _load_images(self):
        """Load images from conditions and populate UI."""
        self._compare_checkboxes: List[QtWidgets.QCheckBox] = []
        self._images.clear()

        # Clear existing
        while self.compare_layout.count():
            item = self.compare_layout.takeAt(0)
            widget = item.widget()
            if widget is not None:
                widget.deleteLater()

        self.cmb_base.clear()
        # Clear Quadrant Fusion combos
        for cmb in self._qf_combos.values():
            cmb.clear()

        if not self._conditions:
            self._update_sidebar_empty_state(False)
            return

        # Load each image
        for cond in self._conditions:
            # Always use filename instead of label (which may have keV/nA format)
            import os
            filename = os.path.basename(cond.image_path)
            label = filename  # Use pure filename

            # Load image
            img = load_image_gray(cond.image_path)
            if img is not None:
                self._images[label] = img

                # Add to base dropdown
                self.cmb_base.addItem(label)

                # Add checkbox for compare
                chk = QtWidgets.QCheckBox(label)
                chk.setProperty("compareItem", True)
                chk.setChecked(False)
                chk.stateChanged.connect(self._update_step_states)  # P0-2
                self.compare_layout.addWidget(chk)
                self._compare_checkboxes.append(chk)

                # Populate Quadrant Fusion combos
                for cmb in self._qf_combos.values():
                    cmb.addItem(label)

        self.compare_layout.addStretch()
        self._refresh_compare_count_label()

        has_images = self.cmb_base.count() > 0
        self._update_sidebar_empty_state(has_images)

        # Select first image as base
        if has_images:
            self.cmb_base.setCurrentIndex(0)
            self._on_base_changed()

    def _update_sidebar_empty_state(self, has_images: bool):
        """Update sidebar placeholders and enablement when image set is empty/non-empty."""
        self.lbl_input_empty_hint.setVisible(not has_images)
        self.lbl_base_empty_hint.setVisible(not has_images)

        self.cmb_base.setEnabled(has_images)
        self.scroll_compare.setEnabled(has_images)
        self.btn_swap_base.setEnabled(has_images)
        self.btn_select_all.setEnabled(has_images)
        self.btn_select_none.setEnabled(has_images)
        if not has_images:
            self.lbl_compare_title.setText("Compare Images (0)")

    def _on_base_changed(self):
        """Update compare checkboxes when base changes and display base image immediately."""
        base_label = self.cmb_base.currentText()

        for chk in self._compare_checkboxes:
            # Disable checkbox for the base image
            is_base = (chk.text() == base_label)
            chk.setEnabled(not is_base)
            chk.setProperty("baseItem", is_base)
            chk.style().unpolish(chk)
            chk.style().polish(chk)
            if is_base:
                chk.setChecked(False)

        self._refresh_compare_count_label()

        # Display base image without waiting for Compute
        base_img = self._images.get(base_label)
        if base_img is not None:
            self.img_base_mag.setImage(base_img)
            if self._roi_manager is not None:
                self._roi_manager.set_image_shape(base_img.shape[:2])
                self.img_base_mag.set_multi_roi_set(self._multi_roi_set)

    def _refresh_compare_count_label(self):
        """Update compare image count in the INPUT card title."""
        base_count = 1 if self.cmb_base.currentText() else 0
        compare_count = max(0, len(getattr(self, "_compare_checkboxes", [])) - base_count)
        self.lbl_compare_title.setText(f"Compare Images ({compare_count})")

    def _select_all_compare(self):
        """Select all compare images."""
        for chk in self._compare_checkboxes:
            if chk.isEnabled():
                chk.setChecked(True)

    def _select_none_compare(self):
        """Deselect all compare images."""
        for chk in self._compare_checkboxes:
            chk.setChecked(False)

    def _on_display_mode(self, mode: str):
        """Switch left or right viewer display mode and keep button highlights in sync."""
        if mode in ('base', 'compare'):
            self._left_non_split_mode = mode
            if self.btn_split_view.isChecked():
                self.btn_split_view.setChecked(False)
            else:
                self._left_view_mode = mode
                self._refresh_base_compare_display(mode)
            self._sync_viewer_mode_buttons()
            return

        if mode in ('diff', 'zmap'):
            self._display_mode = mode
            self._sync_viewer_mode_buttons()
            self._refresh_diff_display()

    def _sync_viewer_mode_buttons(self):
        """Synchronize viewer mode button highlight states with active view."""
        is_split = self._left_view_mode == 'split'
        self.btn_split_view.setChecked(is_split)
        self.btn_mode_base.setChecked((not is_split) and self._left_non_split_mode == 'base')
        self.btn_mode_compare.setChecked((not is_split) and self._left_non_split_mode == 'compare')
        self.btn_mode_diff.setChecked(self._display_mode == 'diff')
        self.btn_mode_topo.setChecked(self._display_mode == 'zmap')

    def _on_left_viewer_clicked(self):
        """Keep left mode button highlight synced when clicking in the left magnifier."""
        if self._left_view_mode != 'split':
            self._left_view_mode = self._left_non_split_mode
        self._sync_viewer_mode_buttons()

    def _on_split_viewer_clicked(self):
        """Keep split button highlighted when user interacts in split view."""
        self._left_view_mode = 'split'
        self._sync_viewer_mode_buttons()

    def _on_right_viewer_clicked(self):
        """Keep right mode button highlight synced when clicking in the right viewer."""
        self._sync_viewer_mode_buttons()

    def _refresh_base_compare_display(self, mode: str):
        """Show base or compare image in the magnifier view."""
        if not self._results:
            return
        result = self._results[self._current_result_idx]
        if mode == 'base':
            base_img = self._images.get(result.base_label)
            if base_img is not None:
                self.img_base_mag.setImage(base_img)
        elif mode == 'compare':
            compare_img = (result.aligned_compare
                           if result.aligned_compare is not None
                           else self._images.get(result.compare_label))
            if compare_img is not None:
                self.img_base_mag.setImage(compare_img)

    def _on_range_changed(self, index: int):
        """Handle range control change."""
        self._refresh_diff_display()

    def _apply_colormap(self, gray_img: np.ndarray) -> np.ndarray:
        """Apply the selected colormap to a grayscale uint8 image → BGR uint8."""
        name = self.cmb_colormap.currentText()
        if name == "Grayscale":
            return cv2.cvtColor(gray_img, cv2.COLOR_GRAY2BGR)
        cv2_maps = {
            "JET": cv2.COLORMAP_JET,
            "Hot": cv2.COLORMAP_HOT,
            "Inferno": cv2.COLORMAP_INFERNO,
            "Viridis": cv2.COLORMAP_VIRIDIS,
        }
        if name in cv2_maps:
            return cv2.applyColorMap(gray_img, cv2_maps[name])
        return cv2.cvtColor(gray_img, cv2.COLOR_GRAY2BGR)

    def _refresh_diff_display(self):
        """Refresh the difference map display with current mode, threshold, and ROI settings."""
        if not self._results:
            return

        result = self._results[self._current_result_idx]

        if self._display_mode == 'zmap':
            display_img = colorize_snr_map(result.snr_map)
        else:
            scaled_img = self._apply_range_scaling(result.result_image)
            if scaled_img.ndim == 2:
                display_img = self._apply_colormap(scaled_img)
            else:
                display_img = scaled_img.copy()

        # ── Gray-level range highlight ─────────────────────────────────────
        if self._hist_range is not None:
            lo, hi = self._hist_range
            # Get source gray values (use original result_image, not scaled display)
            src_gray = result.result_image  # uint8
            if src_gray is not None:
                if display_img.ndim == 2:
                    display_img = cv2.cvtColor(display_img, cv2.COLOR_GRAY2BGR)
                in_range = ((src_gray >= lo) & (src_gray <= hi))
                mask3 = np.stack([in_range, in_range, in_range], axis=2)
                # Pixels outside range: dimmed to 15%
                dimmed = (display_img.astype(np.float32) * 0.15).astype(np.uint8)
                display_img = np.where(mask3, display_img, dimmed)
                # Teal tint on in-range pixels
                tint = display_img.copy()
                tint[:, :, 1] = np.where(in_range,
                                         np.clip(display_img[:, :, 1].astype(np.int16) + 30, 0, 255),
                                         display_img[:, :, 1])
                display_img = tint

        # ── Alignment failure visual overlay ──────────────────────────────
        align_score = result.alignment.final_score
        if align_score < 55:
            if display_img.ndim == 2:
                display_img = cv2.cvtColor(display_img, cv2.COLOR_GRAY2BGR)
            overlay = display_img.copy()
            h_img, w_img = display_img.shape[:2]
            # Red semi-transparent banner at top
            cv2.rectangle(overlay, (0, 0), (w_img, 36), (0, 0, 180), -1)
            cv2.addWeighted(overlay, 0.65, display_img, 0.35, 0, display_img)
            warn_txt = f"\u26a0 ALIGN FAIL  score={align_score:.0f}  diff may be unreliable"
            cv2.putText(display_img, warn_txt, (8, 24),
                        cv2.FONT_HERSHEY_SIMPLEX, 0.55, (255, 255, 255), 1, cv2.LINE_AA)

        self.img_diff.setImage(display_img)

    def _apply_range_scaling(self, img: np.ndarray) -> np.ndarray:
        """Apply dynamic range scaling based on selected mode.

        Modes:
        - Auto: Use full min-max range
        - Zero-centered: Center at 128, scale symmetric
        - P1-P99: Clip to 1st-99th percentile
        - P0.5-P99.5: Clip to 0.5th-99.5th percentile
        """
        if img is None or img.size == 0:
            return img

        range_mode = self.cmb_range.currentText()
        img_f = img.astype(np.float32)

        if range_mode == "Zero-centered":
            # Center at 128, use symmetric range
            center = 128.0
            max_dev = max(abs(img_f.max() - center), abs(img_f.min() - center))
            if max_dev > 0:
                scaled = (img_f - center) / max_dev * 127 + 128
            else:
                scaled = np.full_like(img_f, 128)
            return np.clip(scaled, 0, 255).astype(np.uint8)

        elif range_mode == "P1-P99":
            p_low, p_high = np.percentile(img_f, [1, 99])
            clipped = np.clip(img_f, p_low, p_high)
            if p_high > p_low:
                scaled = (clipped - p_low) / (p_high - p_low) * 255
            else:
                scaled = np.full_like(img_f, 128)
            return scaled.astype(np.uint8)

        elif range_mode == "P0.5-P99.5":
            p_low, p_high = np.percentile(img_f, [0.5, 99.5])
            clipped = np.clip(img_f, p_low, p_high)
            if p_high > p_low:
                scaled = (clipped - p_low) / (p_high - p_low) * 255
            else:
                scaled = np.full_like(img_f, 128)
            return scaled.astype(np.uint8)

        else:  # Auto - use full range
            return img

    def _on_compute(self):
        """Run the combination computation for all selected pairs."""
        base_label = self.cmb_base.currentText()
        is_auto_pair = self.chk_auto_pair.isChecked()

        # Get selected compare images first
        compare_labels = [
            chk.text() for chk in self._compare_checkboxes
            if chk.isChecked() and chk.isEnabled()
        ]

        # Validation
        if is_auto_pair:
            # Auto-pair mode requires at least 2 images
            if len(compare_labels) < 2:
                QtWidgets.QMessageBox.warning(self, "Error", "Please select at least two images for auto pairing.")
                return
        else:
            # Normal mode requires valid base + at least 1 compare
            if not base_label or base_label not in self._images:
                QtWidgets.QMessageBox.warning(self, "Error", "Please select a valid base image.")
                return
            if not compare_labels:
                QtWidgets.QMessageBox.warning(self, "Error", "Please select at least one compare image.")
                return

        # Get operation settings
        operation = 'subtract' if self.cmb_operation.currentIndex() == 0 else 'blend'
        alpha = self.spin_alpha.value()
        beta = self.spin_beta.value()
        invert_base = self.chk_invert_base.isChecked()
        invert_compare = self.chk_invert_compare.isChecked()
        invert_result = self.chk_invert_result.isChecked()
        alignment_method = "phase" if self.cmb_align_method.currentIndex() == 0 else "ncc"
        snr_window_size = self.spn_snr_window.value()
        # Ensure window_size is odd
        if snr_window_size % 2 == 0:
            snr_window_size += 1
        norm_mode = self.cmb_normalize_mode.currentIndex()
        # 0 = Percentile, 1 = GLV-Mask, 2 = Skip, 3 = ROI-Match
        _method_map = {0: 'percentile', 1: 'glv_mask', 2: 'skip', 3: 'roi_match'}
        normalize_method = _method_map.get(norm_mode, 'percentile')
        normalize = (normalize_method not in ('skip', 'roi_match'))
        glv_range = None
        if norm_mode == 1:
            glv_low = self.spn_glv_low.value()
            glv_high = self.spn_glv_high.value()
            if glv_low < glv_high:
                glv_range = (glv_low, glv_high)
        clahe_clip_limit = 2.0

        # ROI-Match (EPI Nulling) parameters
        # ROI is optional — if ROI-Match mode is selected but no ROIs have been
        # defined, silently fall back to Percentile normalization so Compute is
        # never blocked.  The status label updates to inform the user.
        use_roi_match = (norm_mode == 3)
        roi_rect_px = None
        if use_roi_match and len(self._multi_roi_set) == 0:
            use_roi_match = False
            normalize_method = 'percentile'
            normalize = True
            if hasattr(self, 'lbl_roi_status'):
                self.lbl_roi_status.setText(
                    "ROI-Match selected but no ROIs defined — using Percentile fallback"
                )
                self.lbl_roi_status.setStyleSheet(
                    "color: #D97706; font-size: 11px; border: none; background: transparent;"
                )

        sub_mode = self.cmb_subtract_mode.currentIndex()
        # 0 = |diff|×2 (default), 1 = |diff| no gain, 2 = clip≥0 (preserve direction)
        preserve_positive_diff = (sub_mode == 2)
        abs_no_gain = (sub_mode == 1)

        # ROI-Match override: if user wants |diff| with ROI-Match, honour the checkbox
        if use_roi_match and self.chk_roi_abs_diff.isChecked():
            abs_no_gain = True
            preserve_positive_diff = False

        self._last_settings = {
            "operation": operation,
            "alpha": alpha,
            "beta": beta,
            "invert_base": invert_base,
            "invert_compare": invert_compare,
            "invert_result": invert_result,
            "alignment_method": alignment_method,
            "subtract_mode": sub_mode,
            "preserve_positive_diff": preserve_positive_diff,
            "abs_no_gain": abs_no_gain,
            "snr_window_size": snr_window_size,
            "normalize": normalize,
            "normalize_method": normalize_method,
            "glv_range": glv_range,
            "clahe_clip_limit": clahe_clip_limit,
            "roi_match": use_roi_match,
            "roi_rect_px": roi_rect_px,
        }

        # Get images
        base_img = self._images[base_label] if base_label in self._images else None
        compare_imgs = {lbl: self._images[lbl] for lbl in compare_labels if lbl in self._images}

        self._start_compute_worker(
            base_label, base_img, compare_imgs,
            operation, alpha, beta,
            invert_base, invert_compare, invert_result,
            alignment_method,
            preserve_positive_diff,
            abs_no_gain=abs_no_gain,
            normalize=normalize,
            normalize_method=normalize_method,
            glv_range=glv_range,
            clahe_clip_limit=clahe_clip_limit,
            snr_window_size=snr_window_size,
            is_auto_pair=is_auto_pair,  # captured on main thread — do NOT read inside thread
            roi_rect=roi_rect_px,
            roi_set=self._multi_roi_set,
            roi_match=use_roi_match,
        )

    def _start_compute_worker(
            self,
            base_label: str,
            base_img: Optional[np.ndarray],
            compare_imgs: Dict[str, np.ndarray],
            operation: str,
            alpha: float,
            beta: float,
            invert_base: bool,
            invert_compare: bool,
            invert_result: bool,
            alignment_method: str,
            preserve_positive_diff: bool,
            abs_no_gain: bool = False,
            normalize: bool = True,
            normalize_method: str = 'percentile',
            glv_range: Optional[Tuple[int, int]] = None,
            clahe_clip_limit: float = 2.0,
            snr_window_size: int = 31,
            is_auto_pair: bool = False,
            roi_rect: Optional[tuple] = None,
            roi_set: Optional[MultiROISet] = None,
            roi_match: bool = False,
    ):
        if self._compute_thread is not None:
            return

        if is_auto_pair:
            pair_count = len(compare_imgs) * (len(compare_imgs) - 1)
        else:
            pair_count = len(compare_imgs)
        self._compute_abort_requested = False
        self.btn_compute.setEnabled(False)
        self.btn_compute.setText("Computing…")
        self.btn_export.setEnabled(False)
        self.btn_prev_result.setEnabled(False)
        self.btn_next_result.setEnabled(False)
        self._set_compute_busy(True)

        # ── Embedded progress banner ──────────────────────────────────────
        self.progress_bar.setRange(0, pair_count)
        self.progress_bar.setValue(0)
        self.lbl_progress_text.setText(f"Preparing {pair_count} pair(s)…")
        self.wgt_progress_banner.setVisible(True)
        QtWidgets.QApplication.processEvents()

        def _run_compute(worker):
            # IMPORTANT: do NOT access any Qt widget here — this runs on a
            # background thread.  All widget state must be captured before
            # _start_compute_worker is called (see is_auto_pair below).
            if is_auto_pair:
                results: List[SinglePairResult] = []
                labels = list(compare_imgs.keys())
                idx = 0
                for i in range(len(labels)):
                    for j in range(len(labels)):
                        if i == j:
                            continue
                        base_lbl = labels[i]
                        cmp_lbl = labels[j]
                        result = compute_single_pair(
                            base=compare_imgs[base_lbl],
                            compare=compare_imgs[cmp_lbl],
                            base_label=base_lbl,
                            compare_label=cmp_lbl,
                            operation=operation,
                            alpha=alpha,
                            beta=beta,
                            invert_base=invert_base,
                            invert_compare=invert_compare,
                            invert_result=invert_result,
                            normalize=normalize,
                            normalize_method=normalize_method,
                            glv_range=glv_range,
                            clahe_clip_limit=clahe_clip_limit,
                            search_radius=50,
                            alignment_method=alignment_method,
                            preserve_positive_diff=preserve_positive_diff,
                            abs_no_gain=abs_no_gain,
                            snr_window_size=snr_window_size,
                            roi_rect=roi_rect,
                            roi_set=roi_set,
                            roi_match=roi_match,
                        )
                        results.append(result)
                        idx += 1
                        worker.progress.emit(idx, f"{base_lbl} → {cmp_lbl}")
                        if worker.abort_requested:
                            return results
                return results
            # Standard mode: one base vs N compares
            cmp_labels = list(compare_imgs.keys())
            results = []
            for idx, cmp_lbl in enumerate(cmp_labels):
                if worker.abort_requested:
                    return results
                r = compute_single_pair(
                    base=base_img,
                    compare=compare_imgs[cmp_lbl],
                    base_label=base_label,
                    compare_label=cmp_lbl,
                    operation=operation,
                    alpha=alpha,
                    beta=beta,
                    invert_base=invert_base,
                    invert_compare=invert_compare,
                    invert_result=invert_result,
                    normalize=normalize,
                    normalize_method=normalize_method,
                    glv_range=glv_range,
                    clahe_clip_limit=clahe_clip_limit,
                    search_radius=50,
                    alignment_method=alignment_method,
                    preserve_positive_diff=preserve_positive_diff,
                    abs_no_gain=abs_no_gain,
                    snr_window_size=snr_window_size,
                    roi_rect=roi_rect,
                    roi_set=roi_set,
                    roi_match=roi_match,
                )
                results.append(r)
                worker.progress.emit(idx + 1, f"{base_label} → {cmp_lbl}")
                if worker.abort_requested:
                    return results
            return results

        def _on_progress(current, label):
            self.progress_bar.setValue(current)
            self.lbl_progress_text.setText(f"Processing: {label}")

        self._compute_thread, self._compute_worker = self._start_worker(
            _run_compute,
            on_success=self._on_compute_finished,
            on_error=self._on_compute_error,
            on_done=self._on_compute_done,
            on_progress=_on_progress,
        )
        self.btn_abort_compute.setEnabled(True)

    def _start_worker(self, fn, on_success, on_error=None, on_done=None, on_progress=None):
        """Run *fn* on a QThread; all callbacks execute on the main thread.

        Signal routing:
          worker.finished  ──QueuedConn──▶  on_success  (main thread)
          worker.error     ──QueuedConn──▶  on_error    (main thread)
          worker.progress  ──QueuedConn──▶  on_progress (main thread, optional)
          worker.finished/error ──Direct──▶ thread.quit()  (safe, thread-safe call)
          thread.finished  ──QueuedConn──▶  on_done     (main thread)
          thread.finished  ──QueuedConn──▶  deleteLater (deferred, safe)
        """
        thread = QtCore.QThread()  # no parent → avoids thread-affinity issues
        worker = _Worker(fn)
        worker.moveToThread(thread)
        thread.started.connect(worker.run)

        # on_success / on_error: worker is a QObject living in worker-thread;
        # self._on_compute_* are methods of a QObject living in main-thread
        # → AutoConnection auto-promotes to QueuedConnection cross-thread ✓
        worker.finished.connect(on_success)
        if on_error is not None:
            worker.error.connect(on_error)
        if on_progress is not None:
            worker.progress.connect(on_progress)

        # Tell the thread's event loop to exit once work is done.
        # These lambdas discard the argument and call thread.quit(),
        # which is thread-safe from any thread.
        worker.finished.connect(lambda _: thread.quit())
        worker.error.connect(lambda _: thread.quit())

        # After the event loop exits, thread.finished fires on the main thread.
        # Wire housekeeping and on_done here so they always run on main thread.
        thread.finished.connect(worker.deleteLater)
        thread.finished.connect(thread.deleteLater)
        if on_done is not None:
            thread.finished.connect(on_done)  # QueuedConnection → main thread ✓

        thread.start()
        return thread, worker

    def _on_compute_finished(self, results: List[SinglePairResult]):
        if self._compute_abort_requested and len(results) < self.progress_bar.maximum():
            return
        self._results = results
        self._current_result_idx = 0
        self._has_computed = True  # P0-3: flag for Re-compute button state
        self._update_current_result()
        self._update_navigation()
        self.btn_export.setEnabled(bool(self._results))

        # Switch to post-compute view: hide settings, show diff viewer + Back button
        self.left_panel_scroll.setVisible(False)
        self.wgt_diff_section.setVisible(True)
        self.btn_back_to_settings.setVisible(True)
        self.wgt_bottom_row.setVisible(True)
        self.btn_split_view.setVisible(True)
        self.btn_mode_compare.setVisible(True)
        self.btn_prev_result.setVisible(True)
        self.btn_next_result.setVisible(True)
        self.btn_export.setVisible(True)

        # Multi-ROI analysis — run if any ROIs are defined
        if self._multi_roi_set and results:
            self._run_roi_analysis(results)

    def _on_back_to_settings(self) -> None:
        """Restore pre-compute state: show settings panel, hide diff viewer."""
        self.left_panel_scroll.setVisible(True)
        self.wgt_diff_section.setVisible(False)
        self.btn_back_to_settings.setVisible(False)
        self.wgt_bottom_row.setVisible(False)
        self.btn_split_view.setVisible(False)
        self.btn_mode_compare.setVisible(False)
        self.btn_prev_result.setVisible(False)
        self.btn_next_result.setVisible(False)
        self.btn_export.setVisible(False)
        # Reset bottom panels to empty state
        self.align_panel.reset()
        self.diff_roi_panel.reset()
        self._roi_full_results = {}
        self._roi_remapped_sets = {}
        self._roi_ref_base_label = None

    def _run_roi_analysis(self, results: List[SinglePairResult]) -> None:
        """Compute ROI full stats, grouped by base_label to support auto-pair mode.

        In standard mode all results share the same base_label → one ROIFullResult.
        In auto-pair mode each distinct base_label gets its own ROIFullResult so
        the ROI quantification is always computed against the correct base image.
        """
        from collections import defaultdict

        norm_mode = self.cmb_normalize_mode.currentIndex()
        use_roi_match = (norm_mode == 3)

        _method_map = {0: 'percentile', 1: 'glv_mask', 2: 'skip', 3: 'skip'}
        normalize_method = _method_map.get(norm_mode, 'percentile')
        glv_range = None
        if norm_mode == 1:
            glv_range = (self.spn_glv_low.value(), self.spn_glv_high.value())
        clahe_clip = 2.0

        sub_mode = self.cmb_subtract_mode.currentIndex()
        preserve_positive = (sub_mode == 0)
        abs_diff = (sub_mode == 1)

        # ── Group results by base_label ──────────────────────────────────
        groups: Dict[str, List[SinglePairResult]] = defaultdict(list)
        for r in results:
            groups[r.base_label].append(r)

        roi_full_results: Dict[str, ROIFullResult] = {}
        roi_remapped_sets_build: Dict[str, MultiROISet] = {}

        for base_lbl, group in groups.items():
            base_img = self._images.get(base_lbl)
            if base_img is None:
                continue

            # Build aligned_compares for this base group.
            # In ROI-Match mode re-apply per-pair alpha so analysis uses the same
            # calibrated compare image that produced the displayed diff result.
            aligned_compares: Dict[str, np.ndarray] = {}
            skip_group = False
            for r in group:
                comp = r.aligned_compare
                if comp is None:
                    continue
                if use_roi_match:
                    if r.roi_match_alpha is None:
                        QtWidgets.QMessageBox.warning(
                            self,
                            "ROI Analysis",
                            f"ROI analysis is unavailable for pair "
                            f"'{r.base_label} → {r.compare_label}' because the "
                            f"ROI-match scale (α) was not found.",
                        )
                        skip_group = True
                        break
                    comp = np.clip(
                        comp.astype(np.float32) * float(r.roi_match_alpha),
                        0, 255,
                    ).astype(np.uint8)
                aligned_compares[r.compare_label] = comp

            if skip_group or not aligned_compares:
                continue

            # ── ROI coordinate remapping ─────────────────────────────────
            # If the ROI was drawn while a *different* base image was shown in
            # the base viewer, the norm_rect values are in that image's coordinate
            # space.  We use the already-computed alignment offset between the
            # ROI reference base and this base group to shift the ROI coordinates
            # so they point to the correct physical region.
            #
            # pair (roi_ref_base → base_lbl) gives the shift applied to base_lbl
            # to align it with roi_ref_base.  Under this convention:
            #   aligned_compare[y, x] = compare[y - dy, x - dx]   (warpAffine -dx/-dy)
            # → feature at (y, x) in roi_ref_base is at (y + dy, x + dx) in base_lbl
            # → the ROI must move by (+dx, +dy) ... wait: we need the same region
            #   in base_lbl that corresponds to the ROI in roi_ref_base.
            #   base_lbl feature at (y, x)  ≈ roi_ref_base feature at (y - dy, x - dx)
            #   so roi_ref_base ROI at (rx, ry) → base_lbl ROI at (rx + dx, ry + dy)
            # Remap ROI from roi_ref_base coordinate space to base_lbl.
            #
            # _apply_alignment(compare, dx, dy) uses warpAffine with
            #   M = [[1, 0, -dx], [0, 1, -dy]]  (forward mapping src→dst)
            #   which gives:  aligned[y, x] = compare[y + dy, x + dx]
            # For high NCC the pair (base=roi_ref_base, compare=base_lbl) satisfies:
            #   base[y, x] ≈ compare[y + dy, x + dx]
            # → the physical point at (rx, ry) in roi_ref_base is located at
            #   (rx + dx, ry + dy) in base_lbl's original pixel space.
            # Therefore the ROI must be shifted by (+dx, +dy).
            roi_set_for_base = self._multi_roi_set
            ref_lbl = self._roi_ref_base_label
            if (ref_lbl and ref_lbl != base_lbl
                    and len(self._multi_roi_set) > 0):
                # Find pair: base=ref_lbl, compare=base_lbl
                ref_to_base_pair = next(
                    (r for r in results
                     if r.base_label == ref_lbl and r.compare_label == base_lbl),
                    None,
                )
                if ref_to_base_pair is not None and ref_to_base_pair.alignment is not None:
                    dx = ref_to_base_pair.alignment.dx
                    dy = ref_to_base_pair.alignment.dy
                    ref_img = self._images.get(ref_lbl)
                    if ref_img is not None:
                        roi_set_for_base = self._multi_roi_set.shifted(
                            dx, dy, ref_img.shape
                        )

            # Cache the (possibly remapped) ROI set so _apply_roi_visibility
            # can display the correct overlay when this base is shown.
            roi_remapped_sets_build[base_lbl] = roi_set_for_base

            try:
                roi_result = compute_roi_full_stats(
                    base=base_img,
                    aligned_compares=aligned_compares,
                    roi_set=roi_set_for_base,
                    normalize_method=normalize_method,
                    glv_range=glv_range,
                    clahe_clip_limit=clahe_clip,
                    preserve_positive_diff=preserve_positive,
                    abs_diff=abs_diff,
                )
            except Exception as exc:
                QtWidgets.QMessageBox.warning(
                    self, "ROI Analysis",
                    f"ROI analysis failed for base '{base_lbl}':\n{exc}"
                )
                continue

            roi_full_results[base_lbl] = roi_result

        if not roi_full_results:
            return

        # ── Store and refresh UI ─────────────────────────────────────────
        self._roi_full_results = roi_full_results
        self._roi_remapped_sets = roi_remapped_sets_build

        # Refresh the center Diff/ROI panel for the currently displayed result
        if self._results:
            current = self._results[self._current_result_idx]
            self._update_diff_roi_panel(current)
            # Re-apply ROI overlay so the base viewer shows the remapped ROI
            # for the currently displayed base (not the ref-base coords).
            self._apply_roi_visibility()

        # Rebuild the ROI detail dialog with the new multi-base results.
        # Do NOT auto-show — user opens via [ROI Details…] button.
        if self._roi_profile_dialog is not None:
            self._roi_profile_dialog.close()
        self._roi_profile_dialog = ROIIntensityProfileDialog(
            roi_full_results, results, parent=self
        )

    def _on_compute_error(self, message: str):
        QtWidgets.QMessageBox.critical(self, "Error", f"Computation failed:\n{message}")

    def _on_compute_done(self):
        self._compute_thread = None
        self._compute_worker = None
        self._set_compute_busy(False)
        # Hide embedded progress banner
        self.wgt_progress_banner.setVisible(False)
        # Also close legacy progress dialog if present (e.g. from Quadrant Fusion)
        if hasattr(self, '_compute_progress') and self._compute_progress is not None:
            self._compute_progress.close()
            self._compute_progress = None
        self.btn_compute.setEnabled(True)
        # P0-3: Show "Re-compute" after first successful run
        if self._has_computed:
            self.btn_compute.setText("↺  Re-compute")
            self._set_button_icon(self.btn_compute, QtWidgets.QStyle.SP_BrowserReload, "↺  Re-compute")
        else:
            self.btn_compute.setText("Compute")
            self._set_button_icon(self.btn_compute, QtWidgets.QStyle.SP_MediaPlay, "Compute")

    # ── Quadrant Fusion compute ─────────────────────────────────────────

    def _on_compute_quadrant_fusion(self):
        """Run Quadrant Fusion computation."""
        # Validate: all 5 detectors must be assigned
        det_labels = {}
        for det_name, cmb in self._qf_combos.items():
            label = cmb.currentText()
            if not label or label not in self._images:
                QtWidgets.QMessageBox.warning(
                    self, "Quadrant Fusion",
                    f"Please assign a valid image for '{det_name}'."
                )
                return
            det_labels[det_name] = label

        # Check all images are same shape
        shapes = set()
        imgs = {}
        for det_name, label in det_labels.items():
            img = self._images[label]
            shapes.add(img.shape[:2])
            imgs[det_name] = img

        if len(shapes) > 1:
            QtWidgets.QMessageBox.warning(
                self, "Quadrant Fusion",
                "All 5 images must have the same dimensions.\n"
                f"Found different shapes: {shapes}"
            )
            return

        # Read parameters
        output_idx = self.cmb_qf_output.currentIndex()
        output_type = ["bse_clean", "topo", "composite"][output_idx]
        alpha_mode = "auto" if self.cmb_qf_alpha_mode.currentIndex() == 0 else "manual"
        alpha_manual = self.spn_qf_alpha.value()
        beta = self.spn_qf_beta.value()
        sigma = self.spn_qf_sigma.value()

        # ROI for alpha fit
        roi_rect_px = None
        if self._qf_roi_rect is not None:
            h, w = imgs["Illuminator"].shape[:2]
            nx, ny, nw, nh = self._qf_roi_rect
            roi_rect_px = (
                int(nx * w), int(ny * h),
                max(1, int(nw * w)), max(1, int(nh * h)),
            )

        # Capture images for thread
        illum = imgs["Illuminator"]
        top = imgs["Top"]
        bottom = imgs["Bottom"]
        left = imgs["Left"]
        right = imgs["Right"]

        self.btn_compute.setEnabled(False)
        self.btn_compute.setText("Computing…")
        self.btn_export.setEnabled(False)

        self._compute_progress = QtWidgets.QProgressDialog(
            "Computing Quadrant Fusion…\nProcessing 5 detector images.",
            None, 0, 0, self,
        )
        self._compute_progress.setWindowTitle("Quadrant Fusion")
        self._compute_progress.setWindowModality(Qt.NonModal)
        self._compute_progress.setMinimumDuration(0)
        self._compute_progress.setStyleSheet(DIALOG_STYLE)
        self._compute_progress.setValue(0)
        self._compute_progress.show()
        QtWidgets.QApplication.processEvents()

        def _run():
            return compute_quadrant_fusion(
                illum=illum, top=top, bottom=bottom, left=left, right=right,
                output_type=output_type,
                alpha_mode=alpha_mode, alpha_manual=alpha_manual,
                beta=beta, gaussian_sigma=sigma,
                roi_rect=roi_rect_px,
            )

        self._compute_thread, self._compute_worker = self._start_worker(
            _run,
            on_success=self._on_qf_compute_finished,
            on_error=self._on_compute_error,
            on_done=self._on_compute_done,
        )

    def _on_qf_compute_finished(self, result: QuadrantFusionResult):
        """Handle Quadrant Fusion result — display in QF-specific viewers."""
        self._qf_last_result = result
        self._result = None
        self._qf_current_view = "bse_clean"  # default main view

        # Display Illuminator in left viewer
        illum_label = self._qf_combos["Illuminator"].currentText()
        if illum_label in self._images:
            self.img_qf_illum.setImage(self._images[illum_label])

        # Display main output + reference
        self._refresh_qf_display()

        # Update toggle button states
        self.btn_qf_show_bse.setChecked(True)
        self.btn_qf_show_topo.setChecked(False)
        self.btn_qf_show_comp.setChecked(False)

        # Histogram of primary output
        counts, edges = np.histogram(result.fused_image_uint8.ravel(), bins=256, range=(0, 255))
        self.qf_histogram_canvas.plot_histogram(counts, edges)

        # Output type name
        output_names = {"bse_clean": "BSE Enhanced", "topo": "Topography", "composite": "Composite"}
        out_name = output_names.get(
            ["bse_clean", "topo", "composite"][self.cmb_qf_output.currentIndex()], "Fusion"
        )

        # Info bar
        self.lbl_qf_result_info.setText(
            f"Quadrant Fusion complete  |  Output: {out_name}  |  "
            f"alpha={result.alpha_used:.4f}  beta={result.beta_used:.2f}"
        )

        # Stats panel
        fused_f = result.fused_image_uint8.astype(np.float32)
        h, w = result.fused_image_uint8.shape[:2]
        self._qf_stat_labels["alpha"].setText(f"{result.alpha_used:.6f}")
        self._qf_stat_labels["beta"].setText(f"{result.beta_used:.4f}")
        self._qf_stat_labels["sigma"].setText(f"{self.spn_qf_sigma.value():.1f}")
        self._qf_stat_labels["output"].setText(out_name)
        self._qf_stat_labels["mean"].setText(f"{fused_f.mean():.2f}")
        self._qf_stat_labels["std"].setText(f"{fused_f.std():.2f}")
        self._qf_stat_labels["dim"].setText(f"{w} × {h}")

        self.btn_export.setEnabled(True)

    def _on_qf_view_toggle(self, view: str):
        """Switch the main QF output view (bse_clean / topo / composite)."""
        self._qf_current_view = view
        # Exclusive toggle
        self.btn_qf_show_bse.setChecked(view == "bse_clean")
        self.btn_qf_show_topo.setChecked(view == "topo")
        self.btn_qf_show_comp.setChecked(view == "composite")
        self._refresh_qf_display()

    def _refresh_qf_display(self):
        """Refresh QF image viewers based on current view selection."""
        result = self._qf_last_result
        if result is None:
            return

        view = getattr(self, '_qf_current_view', 'bse_clean')

        # Choose main + reference images and group titles
        if view == "bse_clean":
            main_img = result.bse_clean_uint8
            ref_img = result.topo_uint8
            main_title = "BSE Enhanced (Illum − α·Topo)"
            ref_title = "Topography"
        elif view == "topo":
            main_img = result.topo_uint8
            ref_img = result.bse_clean_uint8
            main_title = "Topography"
            ref_title = "BSE Enhanced"
        else:  # composite
            main_img = result.composite_uint8
            ref_img = result.topo_uint8
            main_title = "Composite (BSE + β·Topo)"
            ref_title = "Topography"

        # Apply colormap if not grayscale
        cmap_idx = self.cmb_qf_colormap.currentIndex()
        main_display = self._apply_qf_colormap(main_img, cmap_idx)
        ref_display = self._apply_qf_colormap(ref_img, cmap_idx)

        self.qf_main_grp.setTitle(main_title)
        self.qf_ref_grp.setTitle(ref_title)
        self.img_qf_main.setImage(main_display)
        self.img_qf_ref.setImage(ref_display)

        # Update histogram for main view
        counts, edges = np.histogram(main_img.ravel(), bins=256, range=(0, 255))
        self.qf_histogram_canvas.plot_histogram(counts, edges)

    @staticmethod
    def _apply_qf_colormap(gray_img: np.ndarray, cmap_idx: int) -> np.ndarray:
        """Apply colormap to a grayscale image. 0=Grayscale, 1=JET, 2=Hot, 3=Inferno, 4=Viridis."""
        if cmap_idx == 0:
            return gray_img
        cmap_map = {1: cv2.COLORMAP_JET, 2: cv2.COLORMAP_HOT,
                    3: cv2.COLORMAP_INFERNO, 4: cv2.COLORMAP_VIRIDIS}
        cm = cmap_map.get(cmap_idx, cv2.COLORMAP_JET)
        return cv2.applyColorMap(gray_img, cm)

    def _on_adv_toggle(self, checked: bool):
        """Show/hide Advanced operation options."""
        self.grp_advanced.setVisible(checked)
        self.btn_adv_toggle.setText("Advanced Settings \u25bc" if checked else "Advanced Settings \u25b6")
        # P1-1: badge is only meaningful when section is collapsed
        if hasattr(self, 'lbl_adv_badge'):
            self.lbl_adv_badge.setVisible(not checked)

    def _update_adv_badge(self):
        """P1-1: Refresh the Advanced Settings summary badge text."""
        if not hasattr(self, 'lbl_adv_badge'):
            return
        if self.btn_adv_toggle.isChecked():
            self.lbl_adv_badge.setVisible(False)
            return
        norm_labels = ["Percentile", "GLV-Mask", "Skip", "ROI-Match"]
        sub_labels   = ["|diff|×2", "|diff|", "clip≥0"]
        align_labels = ["Phase", "NCC"]
        norm  = norm_labels[min(self.cmb_normalize_mode.currentIndex(),  len(norm_labels)  - 1)]
        sub   = sub_labels  [min(self.cmb_subtract_mode.currentIndex(),  len(sub_labels)   - 1)]
        align = align_labels[min(self.cmb_align_method.currentIndex(),   len(align_labels) - 1)]
        self.lbl_adv_badge.setText(f"Norm: {norm}  ·  {sub}  ·  Align: {align}")
        self.lbl_adv_badge.setVisible(True)

    def _update_step_states(self):
        """P0-2: Refresh step header completion indicators (circle tint + ✓ status)."""
        if not hasattr(self, '_step1_circle'):
            return

        _DONE_CIRCLE  = (
            "QLabel { background-color: #F59E0B; color: #111827; border-radius: 13px;"
            " font-weight: 700; font-size: 12px; }"
        )
        _PEND_CIRCLE  = (
            "QLabel { background-color: #D1D5DB; color: #1F2937; border-radius: 13px;"
            " font-weight: 700; font-size: 12px; }"
        )
        _DONE_STATUS  = "font-size: 13px; color: #16A34A; font-weight: 700; background: transparent; border: none;"
        _SKIP_STATUS  = "font-size: 13px; color: #9CA3AF; background: transparent; border: none;"
        _PEND_STATUS  = "font-size: 13px; color: #D1D5DB; background: transparent; border: none;"

        # ── Step 1: base selected + at least one compare checked ────────────
        base_ok     = bool(self.cmb_base.currentText())
        compare_ok  = any(
            chk.isChecked() and chk.isEnabled()
            for chk in self._compare_checkboxes
        ) if hasattr(self, '_compare_checkboxes') else False
        step1_done  = base_ok and compare_ok
        self._step1_circle.setStyleSheet(_DONE_CIRCLE if step1_done else _PEND_CIRCLE)
        self._step1_status.setText("✓" if step1_done else "")
        self._step1_status.setStyleSheet(_DONE_STATUS if step1_done else _PEND_STATUS)

        # ── Step 2: ready only after Step 1 is done (encourages 1→2→3 flow) ─
        # Before images are loaded the configure section is not yet meaningful,
        # so keep it gray until the user has images + a base/compare pair ready.
        step2_done = step1_done
        self._step2_circle.setStyleSheet(_DONE_CIRCLE if step2_done else _PEND_CIRCLE)
        self._step2_status.setText("✓" if step2_done else "")
        self._step2_status.setStyleSheet(_DONE_STATUS if step2_done else _PEND_STATUS)

        # ── Step 3: optional — ✓ if ROIs defined, — if none ─────────────────
        has_rois = bool(self._multi_roi_set and len(self._multi_roi_set) > 0)
        self._step3_circle.setStyleSheet(_DONE_CIRCLE if has_rois else _PEND_CIRCLE)
        self._step3_status.setText("✓" if has_rois else "—")
        self._step3_status.setStyleSheet(
            _DONE_STATUS if has_rois else _SKIP_STATUS
        )

    def _show_about_dialog(self):
        """Show the About Fusi\u00b3 dialog."""
        dlg = QtWidgets.QMessageBox(self)
        dlg.setWindowTitle("About Fusi\u00b3")
        dlg.setIcon(QtWidgets.QMessageBox.Information)
        dlg.setText(
            "<h2>Fusi\u00b3</h2>"
            "<p><b>SEM Image Fusion & Defect Analysis</b></p>"
        )
        dlg.setInformativeText(
            "Fusi\u00b3 is an advanced SEM image analysis tool for alignment, "
            "fusion, and nanoscale defect detection across multi-detector SEM images.\n\n"
            "Version 1.1.0"
        )
        dlg.setStyleSheet(DIALOG_STYLE)
        dlg.exec()

    def _on_swap_base(self):
        """Swap base image: cycle to next image as base."""
        if self.cmb_base.count() <= 1:
            return
        next_idx = (self.cmb_base.currentIndex() + 1) % self.cmb_base.count()
        self.cmb_base.setCurrentIndex(next_idx)

    def _on_hist_range_changed(self, lo: int, hi: int):
        """Called when user finishes selecting a gray-level range on histogram."""
        self._hist_range = (lo, hi)
        self.lbl_hist_range.setText(f"Range: GL {lo} - {hi}")
        self.btn_clear_hist_range.setEnabled(True)
        self._refresh_diff_display()

    def _on_clear_hist_range(self):
        """Clear histogram range filter."""
        self._hist_range = None
        self.lbl_hist_range.setText("Range: -")
        self.btn_clear_hist_range.setEnabled(False)
        self.histogram_canvas.clear_range()
        self._refresh_diff_display()

    def _update_current_result(self):
        """Update displays for the current result."""
        if not self._results:
            return

        result = self._results[self._current_result_idx]

        # Result image - use centralized refresh method
        self._refresh_diff_display()

        # Blend preview
        self._update_blend_preview()
        if self._left_view_mode != 'split':
            self._refresh_base_compare_display(self._left_non_split_mode)
        self._sync_viewer_mode_buttons()

        # Histogram
        counts, edges = result.histogram
        self.histogram_canvas.plot_histogram(counts, edges)

        # Alignment scores - adapt SinglePairResult to CombineResult format
        self._update_alignment_display(result)

        # Statistics
        self._update_stats_display(result)
        self._refresh_normalized_compare_dialog(result)

        # ROI overlay: update to the remapped set for this result's base so the
        # overlay position is correct when the user navigates between results.
        self._apply_roi_visibility()

    def _update_blend_preview(self):
        if not self._results:
            return
        result = self._results[self._current_result_idx]
        base_img = self._images.get(result.base_label)
        # Use aligned compare so split view spatial coordinates match Diff/SNR maps
        compare_img = (result.aligned_compare
                       if result.aligned_compare is not None
                       else self._images.get(result.compare_label))

        # Page 0 – Magnifier: show Base image
        if base_img is not None:
            self.img_base_mag.setImage(base_img)

        # Page 1 – Split View: Base left, Aligned Compare right
        self.img_blend.set_images(base_img, compare_img)
        # slider 0 → Base only (divider at far right → ratio=1.0)
        # slider 100 → Compare only (divider at far left → ratio=0.0)
        ratio = 1.0 - self.slider_blend.value() / 100.0
        self.img_blend.set_divider(ratio)

    def _update_alignment_display(self, result: SinglePairResult):
        """Update alignment metrics in the LEFT Alignment panel."""
        a = result.alignment

        if a.final_score >= 75:
            status_text = "\u2713 OK"
            status_color = UI_SUCCESS
        elif a.final_score >= 55:
            status_text = "\u26a0 WARN"
            status_color = UI_PRIMARY_HOVER
        else:
            status_text = "\u2717 FAIL"
            status_color = UI_WARNING

        self.align_panel.update_alignment(
            phase=f"{a.score_phase:.3f}",
            ncc=f"{a.score_ncc:.3f}",
            residual=f"{a.score_residual:.3f}",
            final=f"{a.final_score:.1f}",
            shift=f"({a.dx:+d}, {a.dy:+d})",
            status=status_text,
            status_color=status_color,
        )
        # Keep legacy stats_widget in sync (not visible, but used by some export paths)
        self.stats_widget.update_alignment(
            phase=f"{a.score_phase:.3f}",
            ncc=f"{a.score_ncc:.3f}",
            residual=f"{a.score_residual:.3f}",
            final=f"{a.final_score:.1f}",
            shift=f"({a.dx:+d}, {a.dy:+d})",
            status=status_text,
            status_color=status_color,
        )

    def _update_stats_display(self, result: SinglePairResult):
        """Update statistics widget for single result."""
        s = result.stats
        self.stats_widget.stats_labels["diff_mean"].setText(f"{s.get('diff_mean', 0):.4f}")
        self.stats_widget.stats_labels["diff_std"].setText(f"{s.get('diff_std', 0):.4f}")
        self.stats_widget.stats_labels["hot_pixels"].setText(f"{s.get('hot_pixels', 0)}")
        # Display normalize coefficients: a*I + b
        method_display = _NORMALIZE_METHOD_LABELS.get(result.normalize_method, result.normalize_method)
        is_linear = result.normalize_method in ('percentile', 'glv_mask')
        if is_linear:
            self.stats_widget.stats_labels["norm_coeff"].setText(
                f"{method_display}  a={result.norm_a:.4f}"
            )
        else:
            self.stats_widget.stats_labels["norm_coeff"].setText(method_display)
        # Sub-pixel shift
        dx_f = s.get('dx_subpixel', float(result.alignment.dx))
        dy_f = s.get('dy_subpixel', float(result.alignment.dy))
        self.stats_widget.stats_labels["subpixel_shift"].setText(f"({dx_f:+.2f}, {dy_f:+.2f})")

        # ROI-Match alpha readout (STATE A sidebar label)
        if result.roi_match_alpha is not None:
            self.lbl_roi_alpha.setText(f"ROI-match \u03b1 = {result.roi_match_alpha:.4f}")
        else:
            self.lbl_roi_alpha.setText("")

        # Refresh the CENTER Diff/ROI Analysis panel
        self._update_diff_roi_panel(result)

    def _update_diff_roi_panel(self, result: SinglePairResult) -> None:
        """Refresh the DiffROIAnalysisPanelWidget for the given result.

        Looks up the matching ROIFullResult by result.base_label so that
        auto-pair results are always shown against the correct base group.
        """
        n_target = len([r for r in self._multi_roi_set.rois if r.roi_type == 'target'])
        n_ref = len(self._multi_roi_set.get_references())
        if n_target == 0 and n_ref == 0:
            self.diff_roi_panel.show_no_roi()
        else:
            roi_full = self._roi_full_results.get(result.base_label)
            self.diff_roi_panel.update_result(result, roi_full, n_target, n_ref)

    def _on_show_roi_profile_dialog(self) -> None:
        """Open or raise the ROI Intensity Profile dialog (from [ROI Details…] button)."""
        if self._roi_profile_dialog is None:
            QtWidgets.QMessageBox.information(
                self, "ROI Details",
                "ROI analysis has not been run yet.\n"
                "Compute a result with ROI defined to enable the details view."
            )
            return
        self._roi_profile_dialog.show()
        self._roi_profile_dialog.raise_()
        self._roi_profile_dialog.activateWindow()

    def _update_navigation(self):
        """Update navigation buttons and label."""
        n = len(self._results)
        if n == 0:
            self.lbl_result_info.setText("")
            self.lbl_result_info.setVisible(False)
            self.btn_prev_result.setEnabled(False)
            self.btn_next_result.setEnabled(False)
            return

        idx = self._current_result_idx
        result = self._results[idx]

        # Update label with pair info
        op_name = "Subtract" if result.operation == 'subtract' else "Blend"
        self.lbl_result_info.setText(
            f"{idx + 1}/{n}: {result.base_label} \u2192 {result.compare_label} ({op_name})"
        )
        self.lbl_result_info.setVisible(True)
        self.setWindowTitle(
            f"Fusi\u00b3 \u2014 {result.base_label} \u2192 {result.compare_label}"
        )

        # Enable/disable buttons
        self.btn_prev_result.setEnabled(idx > 0)
        self.btn_next_result.setEnabled(idx < n - 1)

    def _reset_hist_range(self):
        """Clear histogram range filter when switching to a different pair."""
        self._hist_range = None
        self.lbl_hist_range.setText("Range: —")
        self.btn_clear_hist_range.setEnabled(False)
        self.histogram_canvas.clear_range()

    def _on_prev_result(self):
        """Go to previous result."""
        if self._current_result_idx > 0:
            self._reset_hist_range()
            self._current_result_idx -= 1
            self._update_current_result()
            self._update_navigation()

    def _on_next_result(self):
        """Go to next result."""
        if self._current_result_idx < len(self._results) - 1:
            self._reset_hist_range()
            self._current_result_idx += 1
            self._update_current_result()
            self._update_navigation()

    # Legacy methods for backward compatibility
    def _update_results(self):
        """Legacy: Update result displays."""
        if self._result is None:
            return
        self.img_diff.setImage(self._result.diff_image)
        counts, edges = self._result.histogram
        self.histogram_canvas.plot_histogram(counts, edges)
        self.align_score_widget.update_scores(self._result)
        self.stats_widget.update_stats(self._result)

    def _update_hint_overlay(self):
        """Legacy: no-op (hint overlay removed)."""
        pass

    def _on_base_mag_cursor_moved(self, norm_x: float, norm_y: float):
        """Relay left-panel magnifier cursor position to Difference Map (bidirectional zoom)."""
        if (not self.btn_split_view.isChecked()
                and self.img_base_mag._multi_draw_mode == 'idle'
                and not self.img_base_mag._roi_mode):
            self.img_diff.setCursorPos(norm_x, norm_y)

    def _on_base_mag_cursor_left(self):
        """Clear Difference Map magnifier when cursor leaves left-panel magnifier."""
        if not self.btn_split_view.isChecked():
            self.img_diff.clearCursor()

    # ── Multi-ROI Manager ────────────────────────────────────────────────

    def _on_open_roi_manager(self) -> None:
        """Toggle the ROI Manager slide-in side panel (P1-2)."""
        # If already visible, toggle-close it (resets draw mode too)
        if self.roi_side_panel.isVisible():
            self._close_roi_side_panel()
            return

        # Lazily create and embed MultiROIManagerWidget on first open
        if not self._roi_panel_populated:
            self._roi_manager = MultiROIManagerWidget(
                roi_set=self._multi_roi_set,
                base_widget=self.img_base_mag,
                parent=None,  # no dialog parent — we embed it as a plain widget
            )
            # Remove dialog chrome — make it behave like a regular embedded widget
            self._roi_manager.setWindowFlags(Qt.Widget)
            self._roi_manager.rois_changed.connect(self._on_multi_rois_changed)
            self._roi_side_scroll.setWidget(self._roi_manager)
            self._roi_panel_populated = True

            # --- Fix: _on_confirm calls accept() which hides the embedded widget.
            # Disconnect the default confirm signal and replace with a handler that
            # (a) clears markers, (b) refreshes ROI state, (c) closes the side panel
            # — without hiding the _roi_manager widget itself.
            try:
                self._roi_manager._btn_confirm.clicked.disconnect()
            except RuntimeError:
                pass  # no connection to disconnect
            self._roi_manager._btn_confirm.clicked.connect(self._on_roi_panel_confirm)

        # Ensure the embedded widget is always visible (accept() may have hidden it)
        self._roi_manager.setVisible(True)

        # Provide current base image shape for pixel→norm conversion
        base_label = self.cmb_base.currentText()
        base_img = self._images.get(base_label)
        if base_img is not None:
            self._roi_manager.set_image_shape(base_img.shape[:2])
            self._apply_roi_visibility()

        # In auto-pair mode before any compute, img_base_mag may show stale or
        # empty content because _on_base_changed() is not called when auto-pair
        # is toggled on.  Force-display the cmb_base image so the user draws ROI
        # on a known, labeled image, ensuring _roi_ref_base_label is correct.
        if self.chk_auto_pair.isChecked() and not self._results and base_img is not None:
            self.img_base_mag.setImage(base_img)

        # Record which base image the ROI is being drawn on.
        # In auto-pair mode, use the currently displayed result's base_label
        # (that image is what img_base_mag shows).  In standard mode, use cmb_base.
        self._capture_roi_ref_base()

        self.roi_side_panel.setVisible(True)

    def _close_roi_side_panel(self) -> None:
        """Close the ROI side panel and reset the image viewer's draw mode to idle.

        Called by the ✕ button and by the toggle-close path so that the ghost
        bounding-box cursor preview stops rendering as soon as the panel closes.
        """
        if self._roi_manager is not None:
            self._roi_manager._clear_multi_add_markers()  # sets draw mode → idle
        self.img_diff.clearCursor()
        self.roi_side_panel.setVisible(False)

    def _on_roi_panel_confirm(self) -> None:
        """Embedded-panel version of ROI confirm.

        Replaces MultiROIManagerWidget._on_confirm so that accept() is never
        called on the embedded widget (which would hide it and break re-opens).
        Instead we clean up multi-add state, refresh ROI overlays, then close
        the side panel — leaving _roi_manager itself visible inside the scroll.
        """
        if self._roi_manager is not None:
            self._roi_manager._clear_multi_add_markers()
        self._on_multi_rois_changed()
        self.img_diff.clearCursor()
        self.roi_side_panel.setVisible(False)

    def _capture_roi_ref_base(self) -> None:
        """Record the base-image label currently shown in the base viewer."""
        if self._results and 0 <= self._current_result_idx < len(self._results):
            self._roi_ref_base_label = self._results[self._current_result_idx].base_label
        elif self.cmb_base.currentText():
            self._roi_ref_base_label = self.cmb_base.currentText()

    def _on_multi_rois_changed(self) -> None:
        """Refresh all image widgets when ROI set changes."""
        self._apply_roi_visibility()
        self._update_roi_status_label()
        # Keep the ref-base up to date when the user adds/modifies ROIs
        self._capture_roi_ref_base()
        self._update_step_states()  # P0-2: update Step 3 ✓ indicator

    def _on_roi_view_toggled(self, _checked: bool) -> None:
        self._apply_roi_visibility()

    def _apply_roi_visibility(self) -> None:
        show = getattr(self, 'chk_roi_view', None) is None or self.chk_roi_view.isChecked()
        if not show:
            self.img_base_mag.set_multi_roi_set(None)
            self.img_diff.set_multi_roi_set(None)
            return

        # In auto-pair mode, use the remapped ROI set for the currently
        # displayed base so the overlay stays spatially correct on that image.
        roi_set = self._multi_roi_set  # default (ref-base coords or no results)
        if self._roi_remapped_sets and self._results:
            cur = self._results[self._current_result_idx]
            roi_set = self._roi_remapped_sets.get(cur.base_label, self._multi_roi_set)

        self.img_base_mag.set_multi_roi_set(roi_set)
        self.img_diff.set_multi_roi_set(roi_set)

    def _update_roi_status_label(self) -> None:
        """Update the ROI count status label below the ROI Manager button."""
        n = len(self._multi_roi_set)
        if n == 0:
            self.lbl_roi_status.setText("No ROIs — analysis will be skipped")
            self.lbl_roi_status.setStyleSheet(
                "color: #9CA3AF; font-size: 11px; border: none; background: transparent;"
            )
        else:
            t = sum(1 for r in self._multi_roi_set.rois if r.roi_type == 'target')
            ref = n - t
            self.lbl_roi_status.setText(f"{n} ROIs  (T:{t}  R:{ref}) — analysis enabled ✓")
            self.lbl_roi_status.setStyleSheet(
                "color: #16A34A; font-size: 11px; font-weight: 600; border: none; background: transparent;"
            )

    # ── ROI-Match (EPI Nulling) handlers ─────────────────────────────────

    def _on_open_roi_manager_for_match(self):
        """Open ROI Manager so ROI-Match can use ROIs defined there."""
        self._on_open_roi_manager()
        n = len(self._multi_roi_set)
        self.lbl_roi_info.setText(f"ROI source: ROI Manager ({n} ROI{'s' if n != 1 else ''})")

    def _on_roi_selected(self, norm_x: float, norm_y: float, norm_w: float, norm_h: float):
        """Handle ROI drawn on an image widget (shared by Standard & QF modes)."""
        roi = (norm_x, norm_y, norm_w, norm_h)

        self.img_base_mag.set_active_roi(roi)
        self.img_diff.set_active_roi(roi)
        self.lbl_roi_alpha.setText("")

    def _on_normalize_mode_changed(self):
        """Show/hide GLV-Mask / ROI-Match controls depending on the selected normalize mode."""
        idx = self.cmb_normalize_mode.currentIndex()
        self.wgt_glv_controls.setVisible(idx == 1)  # GLV-Mask
        self.wgt_roi_match_controls.setVisible(idx == 3)  # ROI-Match (EPI Nulling)

    def _on_preview_glv_mask(self):
        """Open the interactive GLV Mask Preview dialog."""
        base_label = self.cmb_base.currentText() if hasattr(self, "cmb_base") else ""
        base_img = self._images.get(base_label) if base_label else None
        if base_img is None:
            QtWidgets.QMessageBox.information(
                self, "GLV Mask Preview", "Please load a Base image first."
            )
            return
        glv_low = self.spn_glv_low.value()
        glv_high = self.spn_glv_high.value()

        def _apply_to_main(low: int, high: int):
            """Sync values from the preview dialog back to the main dialog spinboxes."""
            self.spn_glv_low.setValue(low)
            self.spn_glv_high.setValue(high)

        dlg = GLVMaskPreviewDialog(
            base_img, glv_low, glv_high,
            parent=self,
            apply_callback=_apply_to_main,
        )
        dlg.exec()

    def _on_show_normalized_compare(self):
        """Show normalized compare image dialog."""
        if not self._results:
            QtWidgets.QMessageBox.information(self, "Info", "Please run Compute first.")
            return
        result = self._results[self._current_result_idx]
        if result.normalized_compare is None or result.aligned_compare is None:
            QtWidgets.QMessageBox.information(self, "Info", "Normalized compare image is not available.")
            return
        base_img = self._images.get(result.base_label)
        if self._norm_compare_dialog is None:
            self._norm_compare_dialog = NormalizedCompareDialog(self)
        self._norm_compare_dialog.set_images(
            result.aligned_compare,
            result.normalized_compare,
            base_img,
            result.base_label,
            result.compare_label,
            result.norm_a,
            result.norm_b,
            method_name=result.normalize_method,
        )
        self._norm_compare_dialog.show()
        self._norm_compare_dialog.raise_()
        self._norm_compare_dialog.activateWindow()

    def _refresh_normalized_compare_dialog(self, result: SinglePairResult):
        if self._norm_compare_dialog is None:
            return
        if not self._norm_compare_dialog.isVisible():
            return
        if result.normalized_compare is None or result.aligned_compare is None:
            return
        base_img = self._images.get(result.base_label)
        self._norm_compare_dialog.set_images(
            result.aligned_compare,
            result.normalized_compare,
            base_img,
            result.base_label,
            result.compare_label,
            result.norm_a,
            result.norm_b,
            method_name=result.normalize_method,
        )

    def _center_crop_image(self, image: Optional[np.ndarray], crop_size: int) -> Optional[np.ndarray]:
        """Center-crop an image to square crop_size; if too small, crop to max possible square."""
        if image is None or image.size == 0:
            return image
        h, w = image.shape[:2]
        side = max(1, min(int(crop_size), h, w))
        cy, cx = h // 2, w // 2
        half = side // 2
        y1 = max(0, cy - half)
        x1 = max(0, cx - half)
        y2 = y1 + side
        x2 = x1 + side
        if y2 > h:
            y2 = h
            y1 = h - side
        if x2 > w:
            x2 = w
            x1 = w - side
        return image[y1:y2, x1:x2]

    def _save_export_images_for_result(
            self,
            result: SinglePairResult,
            export_dirs: Dict[str, str],
            prefix: str,
            center_crop: bool = False,
            crop_size: int = 512,
    ) -> Dict[str, str]:
        """Save per-result images and return saved path mapping."""
        import os

        def _maybe_crop(img: Optional[np.ndarray]) -> Optional[np.ndarray]:
            if not center_crop:
                return img
            return self._center_crop_image(img, crop_size)

        paths: Dict[str, str] = {}

        base_img = _maybe_crop(self._images.get(result.base_label))
        # Use aligned compare so it shares the same coordinate system as diff / SNR maps
        compare_img = _maybe_crop(
            result.aligned_compare
            if result.aligned_compare is not None
            else self._images.get(result.compare_label)
        )
        diff_img = _maybe_crop(result.result_image)
        # snr_map is computed from result_float (same spatial shape as result_image),
        # so colorized SNR map is pixel-perfectly aligned with diff_img — same coord system.
        snr_img = _maybe_crop(colorize_snr_map(result.snr_map))
        norm_img = _maybe_crop(result.normalized_compare)
        aligned_img = _maybe_crop(result.aligned_compare)

        if base_img is not None:
            base_path = os.path.join(export_dirs["base"], f"{prefix}_base.png")
            cv2.imwrite(base_path, base_img)
            paths["base"] = base_path
        if compare_img is not None:
            compare_path = os.path.join(export_dirs["compare"], f"{prefix}_compare.png")
            cv2.imwrite(compare_path, compare_img)
            paths["compare"] = compare_path

        diff_path = os.path.join(export_dirs["diff"], f"{prefix}_diff.png")
        cv2.imwrite(diff_path, diff_img)
        paths["diff"] = diff_path

        snr_path = os.path.join(export_dirs["snr"], f"{prefix}_snr.png")
        cv2.imwrite(snr_path, snr_img)
        paths["snr"] = snr_path

        if norm_img is not None:
            norm_path = os.path.join(export_dirs["normalized"], f"{prefix}_normalized.png")
            cv2.imwrite(norm_path, norm_img)
            paths["normalized"] = norm_path

        if aligned_img is not None:
            aligned_path = os.path.join(export_dirs["aligned"], f"{prefix}_aligned.png")
            cv2.imwrite(aligned_path, aligned_img)
            paths["aligned"] = aligned_path

        return paths

    def _ask_export_options(self) -> tuple:
        """Show the Export Options dialog.

        Returns
        -------
        (do_center_crop: bool, crop_size: int, export_gif: bool,
         ppt_sections: dict)  or  (False, 512, False, None) on cancel.
        """
        has_roi = bool(self._roi_full_results)

        dlg = QtWidgets.QDialog(self)
        dlg.setWindowTitle("Export Options")
        dlg.setStyleSheet(DIALOG_STYLE)
        dlg.resize(480, 0)
        root = QtWidgets.QVBoxLayout(dlg)
        root.setSpacing(12)
        root.setContentsMargins(16, 16, 16, 12)

        def _section_label(text: str) -> QtWidgets.QLabel:
            lbl = QtWidgets.QLabel(text)
            lbl.setStyleSheet(
                "font-weight: bold; color: #111827; font-size: 12px;"
                "padding-bottom: 2px; border-bottom: 1px solid #E5E7EB;"
            )
            return lbl

        # ── Image / File options ──────────────────────────────────────
        root.addWidget(_section_label("Image & File Options"))

        chk_crop = QtWidgets.QCheckBox("Center-crop images  (useful when defect is centered)")
        chk_crop.setChecked(False)
        root.addWidget(chk_crop)

        crop_row = QtWidgets.QHBoxLayout()
        crop_row.setContentsMargins(20, 0, 0, 0)
        crop_row.addWidget(QtWidgets.QLabel("Crop size (px):"))
        spn = QtWidgets.QSpinBox()
        spn.setRange(64, 4096)
        spn.setSingleStep(64)
        spn.setValue(512)
        spn.setEnabled(False)
        crop_row.addWidget(spn)
        crop_row.addStretch()
        root.addLayout(crop_row)
        chk_crop.toggled.connect(spn.setEnabled)

        chk_gif = QtWidgets.QCheckBox(
            "Export animated GIF  (Base → Normalized Compare → Diff loop)"
        )
        chk_gif.setToolTip(
            "Generates a looping GIF per pair — handy for slides/email.\n"
            "Requires Pillow (pip install Pillow)."
        )
        chk_gif.setChecked(False)
        root.addWidget(chk_gif)

        # ── PPT Report sections ───────────────────────────────────────
        root.addWidget(_section_label("PowerPoint Report  —  Sections to Include"))

        # (label, key, default_checked, tooltip)
        _PPT_SECTION_DEFS = [
            ("Overview Table  — summary of all pairs",
             "overview_table",    True,
             "One or more slides listing all pairs with alignment/diff metrics."),
            ("Detail Slides  — Base / Compare / Diff images per pair",
             "detail_slides",     True,
             "3×2 image grid for every pair: base, aligned compare, normalized compare,\n"
             "diff map, JET SNR map, and GIF animation frame."),
            ("ROI Pair Summary Table  — LE metrics & SNR per pair",
             "roi_summary_table", True,
             "Multi-page table with T/R mean, Δ, Pair SNR columns (requires ROI data)."),
            ("SNR Pair Matrix  — heatmap of SNR across all conditions",
             "snr_matrix",        True,
             "N×N heatmap where each cell shows the Pair SNR for that (base, compare) pair."),
            ("Intensity Profile Charts  — per-ROI mean across LE",
             "intensity_profiles", True,
             "Line chart per base image: base mean (dashed), compare mean, diff mean per ROI."),
            ("Image Gallery  — all images organized by condition",
             "image_gallery",     True,
             "Thumbnail grid showing all acquired images grouped by LE condition."),
            ("ROI Position Maps  — base image with ROI overlays",
             "roi_position_maps", False,
             "Annotated base image showing target and reference ROI bounding boxes."),
            ("Diff Map + ROI Overlay  — diff image with ROI boxes",
             "diff_roi_overlay",  False,
             "Diff image for each pair with target/reference ROI boxes drawn on top."),
            ("Base SNR Summary Chart  — bar chart per base image",
             "base_snr_chart",    False,
             "Bar chart comparing Base SNR across all base images (color-coded by threshold)."),
        ]

        section_checks: Dict[str, QtWidgets.QCheckBox] = {}
        grid = QtWidgets.QGridLayout()
        grid.setHorizontalSpacing(8)
        grid.setVerticalSpacing(4)
        grid.setContentsMargins(4, 0, 0, 0)

        for row_i, (label, key, default, tip) in enumerate(_PPT_SECTION_DEFS):
            # Disable ROI-related sections when no ROI data loaded
            roi_keys = {"roi_summary_table", "snr_matrix", "intensity_profiles",
                        "roi_position_maps", "diff_roi_overlay", "base_snr_chart"}
            chk_s = QtWidgets.QCheckBox(label)
            chk_s.setChecked(default)
            chk_s.setToolTip(tip)
            if key in roi_keys and not has_roi:
                chk_s.setChecked(False)
                chk_s.setEnabled(False)
                chk_s.setToolTip("No ROI data loaded — run ROI analysis first.")
            section_checks[key] = chk_s
            grid.addWidget(chk_s, row_i, 0)

        root.addLayout(grid)

        # ── Dialog buttons ────────────────────────────────────────────
        btn_row = QtWidgets.QDialogButtonBox(
            QtWidgets.QDialogButtonBox.Ok | QtWidgets.QDialogButtonBox.Cancel
        )
        btn_row.accepted.connect(dlg.accept)
        btn_row.rejected.connect(dlg.reject)
        root.addWidget(btn_row)

        dlg.adjustSize()

        if dlg.exec() == QtWidgets.QDialog.Accepted:
            ppt_sections = {k: chk.isChecked() for k, chk in section_checks.items()}
            return chk_crop.isChecked(), spn.value(), chk_gif.isChecked(), ppt_sections
        return False, 512, False, None

    # ------------------------------------------------------------------
    # PPT ROI helper  (module-level, called from _export_ppt_report)
    # ------------------------------------------------------------------

    def _export_ppt_report(self, out_dir: str, result_rows: List[Dict[str, object]],
                           settings: Dict[str, object],
                           do_center_crop: bool = False,
                           crop_size: int = 512,
                           roi_full_results: Optional[Dict[str, object]] = None,
                           roi_all_results: Optional[List[object]] = None,
                           ppt_sections: Optional[Dict[str, bool]] = None,
                           progress_callback=None) -> Optional[str]:
        """Build a dark-themed PPT report for all computed image pairs.

        ppt_sections keys: overview_table, detail_slides, roi_summary_table,
          snr_matrix, roi_position_maps, diff_roi_overlay, intensity_profiles,
          image_gallery, base_snr_chart.
        All sections default to True when ppt_sections is None.
        """
        # Default: include every section
        _sec: Dict[str, bool] = ppt_sections or {}
        def _s(key: str, default: bool = True) -> bool:
            return _sec.get(key, default)

        def _prog(label: str) -> None:
            if progress_callback:
                progress_callback(label)

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            from pptx.oxml.ns import qn
            from lxml import etree
        except Exception:
            return None

        import datetime

        # ── Color palette ────────────────────────────────────────────────────
        C_BG = RGBColor(0x0F, 0x14, 0x1A)  # #0F141A  almost-black bg
        C_CARD = RGBColor(0x16, 0x1E, 0x2C)  # #161E2C  card bg
        C_PRIMARY = RGBColor(0xF5, 0x9E, 0x0B)  # #F59E0B  amber accent
        C_TEXT = RGBColor(0xE2, 0xE8, 0xF0)  # #E2E8F0  light text
        C_TEXT_SEC = RGBColor(0x8A, 0x99, 0xAA)  # muted secondary
        C_SUCCESS = RGBColor(0x26, 0xD7, 0xAE)  # teal / highlight
        C_WARN = RGBColor(0xEF, 0x44, 0x44)  # red for bad score

        SLIDE_W = Inches(13.33)
        SLIDE_H = Inches(7.5)

        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        def _fill_bg(slide, color: RGBColor):
            """Fill slide background with a solid color via XML."""
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = color

        def _add_text(slide, text, left, top, w, h, size=11, bold=False,
                      color=None, align=PP_ALIGN.LEFT, italic=False):
            txb = slide.shapes.add_textbox(left, top, w, h)
            tf = txb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color or C_TEXT
            return txb

        def _score_color(score: float) -> RGBColor:
            if score >= 75:
                return C_SUCCESS
            if score >= 55:
                return C_PRIMARY
            return C_WARN

        # ── Title slide (always included) ────────────────────────────────────
        _prog("Generating Title slide…")
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        _fill_bg(title_slide, C_BG)

        # Amber bar on left
        bar = title_slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), Inches(0.18), SLIDE_H
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = C_PRIMARY
        bar.line.fill.background()

        _add_text(title_slide, "Perspective Combination",
                  Inches(0.4), Inches(2.6), Inches(9), Inches(0.9),
                  size=36, bold=True, color=C_TEXT)
        _add_text(title_slide, "Multi-Image Comparison Report",
                  Inches(0.4), Inches(3.5), Inches(9), Inches(0.55),
                  size=22, color=C_PRIMARY)
        _add_text(title_slide,
                  f"Generated: {datetime.date.today()}   |   "
                  f"Pairs: {len(result_rows)}   |   "
                  f"Mode: {settings.get('operation', '?').capitalize()}   |   "
                  f"Alignment: {settings.get('alignment_method', '?').capitalize()}",
                  Inches(0.4), Inches(4.3), Inches(12.5), Inches(0.4),
                  size=11, color=C_TEXT_SEC)

        # ── Overview table slide(s) ─────────────────────────────────────────
        if _s('overview_table'):
            _prog("Generating Overview Table…")
        ROWS_PER = 18
        for block_start in range(0, len(result_rows), ROWS_PER):
            if not _s('overview_table'):
                break
            block = result_rows[block_start:block_start + ROWS_PER]
            ov = prs.slides.add_slide(prs.slide_layouts[6])
            _fill_bg(ov, C_BG)

            page_n = block_start // ROWS_PER + 1
            total_pages = (len(result_rows) + ROWS_PER - 1) // ROWS_PER
            _add_text(ov, f"Overview  —  All Pairs  (page {page_n}/{total_pages})",
                      Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.45),
                      size=16, bold=True, color=C_PRIMARY)

            # Column headers
            cols = ["#", "Base → Compare", "Mode", "dx / dy", "Align", "Diff mean", "Diff std", "Hot px"]
            col_x = [0.4, 0.8, 4.5, 5.5, 6.5, 7.6, 9.0, 10.4]
            col_w = [0.38, 3.65, 0.95, 0.95, 1.05, 1.35, 1.35, 1.9]
            for ci, (ch, cx, cw) in enumerate(zip(cols, col_x, col_w)):
                _add_text(ov, ch, Inches(cx), Inches(0.72), Inches(cw), Inches(0.28),
                          size=9, bold=True, color=C_TEXT_SEC)

            # Separator line
            line = ov.shapes.add_shape(1, Inches(0.4), Inches(1.02), Inches(12.5), Inches(0.02))
            line.fill.solid()
            line.fill.fore_color.rgb = C_TEXT_SEC
            line.line.fill.background()

            row_h = Inches(0.33)
            for ri, row in enumerate(block):
                result: SinglePairResult = row["result"]
                s = result.stats
                y = Inches(1.1) + ri * row_h
                bg_col = C_CARD if ri % 2 == 0 else C_BG
                bg_rect = ov.shapes.add_shape(1, Inches(0.35), y, Inches(12.6), row_h)
                bg_rect.fill.solid()
                bg_rect.fill.fore_color.rgb = bg_col
                bg_rect.line.fill.background()

                global_idx = block_start + ri + 1
                score = result.stats.get('alignment_score', result.alignment.final_score)
                vals = [
                    f"{global_idx:02d}",
                    f"{result.base_label[:18]} → {result.compare_label[:18]}",
                    result.operation,
                    f"({result.alignment.dx:+d}, {result.alignment.dy:+d})",
                    f"{score:.1f}",
                    f"{s.get('diff_mean', 0):.5f}",
                    f"{s.get('diff_std', 0):.5f}",
                    f"{s.get('hot_pixels', 0)}",
                ]
                for ci, (val, cx, cw) in enumerate(zip(vals, col_x, col_w)):
                    col_color = _score_color(score) if ci == 4 else C_TEXT
                    _add_text(ov, val, Inches(cx), y + Inches(0.04), Inches(cw), row_h,
                              size=9, color=col_color)

        # ── Detail slides  ─  3 columns × 2 rows  ────────────────────────────
        #  Row 1: Base Image | Compare Image | Normalized Compare
        #  Row 2: Difference Map | JET Z-Map  | GIF Animation frame
        LABEL_H = Inches(0.26)
        HEADER_H = Inches(0.85)
        LEFT_M = Inches(0.22)
        COL_GAP = Inches(0.09)
        ROW_GAP = Inches(0.09)
        BOT_M = Inches(0.08)

        # 3-column widths (equal)
        IMG_W = (SLIDE_W - 2 * LEFT_M - 2 * COL_GAP) / 3  # ≈ 4.24″
        col_x = [
            LEFT_M,
            LEFT_M + IMG_W + COL_GAP,
            LEFT_M + 2 * (IMG_W + COL_GAP),
        ]

        # 2-row heights (equal)
        total_img_h = (SLIDE_H - HEADER_H - LABEL_H * 2 - ROW_GAP - BOT_M)
        IMG_H = total_img_h / 2
        row_y = [
            HEADER_H + Inches(0.04),
            HEADER_H + Inches(0.04) + LABEL_H + IMG_H + ROW_GAP,
        ]

        # Operation badge colors
        C_SUBTRACT = RGBColor(0x26, 0xD7, 0xAE)  # teal  → Subtract
        C_BLEND = RGBColor(0x60, 0x9C, 0xFF)  # blue  → Blend

        for idx, row in enumerate(result_rows, start=1):
            if not _s('detail_slides'):
                break
            result: SinglePairResult = row["result"]
            paths: Dict[str, str] = row["paths"]
            s = result.stats
            score = s.get('alignment_score', result.alignment.final_score)
            _prog(f"Detail slide  [{result.base_label} → {result.compare_label}]  ({idx}/{len(result_rows)})")

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _fill_bg(slide, C_BG)

            # ── Header bar ──────────────────────────────────────────────────
            hdr_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, HEADER_H)
            hdr_bar.fill.solid()
            hdr_bar.fill.fore_color.rgb = C_CARD
            hdr_bar.line.fill.background()

            # Pair number badge (amber)
            badge = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.55), HEADER_H)
            badge.fill.solid()
            badge.fill.fore_color.rgb = C_PRIMARY
            badge.line.fill.background()
            _add_text(slide, f"{idx:02d}",
                      Inches(0.0), Inches(0.18), Inches(0.55), Inches(0.5),
                      size=18, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

            _add_text(slide,
                      f"{result.base_label}  →  {result.compare_label}",
                      Inches(0.65), Inches(0.06), Inches(8.0), Inches(0.38),
                      size=14, bold=True, color=C_TEXT)

            crop_note = f"  crop={crop_size}×{crop_size}px" if do_center_crop else ""
            _add_text(slide,
                      f"align={result.alignment.method}  "
                      f"shift=({result.alignment.dx:+d},{result.alignment.dy:+d})  "
                      f"score={score:.1f}  "
                      f"norm a={result.norm_a:.3f} b={result.norm_b:.2f}  "
                      f"diff μ={s.get('diff_mean', 0):.4f} σ={s.get('diff_std', 0):.4f}  "
                      f"hot={s.get('hot_pixels', 0)}{crop_note}",
                      Inches(0.65), Inches(0.48), Inches(9.6), Inches(0.32),
                      size=8.5, color=C_TEXT_SEC)

            # ── Operation mode badge ─────────────────────────────────────────
            op_label = result.operation.capitalize()
            op_color = C_SUBTRACT if result.operation == 'subtract' else C_BLEND
            op_pill = slide.shapes.add_shape(1, Inches(10.35), Inches(0.08),
                                             Inches(1.25), Inches(0.36))
            op_pill.fill.solid()
            op_pill.fill.fore_color.rgb = op_color
            op_pill.line.fill.background()
            _add_text(slide, op_label,
                      Inches(10.35), Inches(0.10), Inches(1.25), Inches(0.32),
                      size=10, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

            # ── Alignment score colored pill ─────────────────────────────────
            score_c = _score_color(score)
            spill = slide.shapes.add_shape(1, Inches(11.7), Inches(0.08),
                                           Inches(1.45), Inches(0.36))
            spill.fill.solid()
            spill.fill.fore_color.rgb = score_c
            spill.line.fill.background()
            align_status = ("OK" if score >= 75
                            else ("WARN" if score >= 55 else "FAIL"))
            _add_text(slide, f"Align {align_status}  {score:.0f}",
                      Inches(11.7), Inches(0.10), Inches(1.45), Inches(0.32),
                      size=9, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

            # ── 6-image grid  (3 × 2) ────────────────────────────────────────
            gif_label = "GIF Animation (1st frame)" if "gif" in paths else "GIF (not exported)"
            layout_slots = [
                # Top row
                ("base", "Base Image", col_x[0], row_y[0]),
                ("compare", "Aligned Compare", col_x[1], row_y[0]),
                ("normalized", "Normalized Compare", col_x[2], row_y[0]),
                # Bottom row
                ("diff", "Difference Map", col_x[0], row_y[1]),
                ("snr", "JET Z-Map  (SNR)", col_x[1], row_y[1]),
                ("gif", gif_label, col_x[2], row_y[1]),
            ]

            def _add_picture_aspect(slide, path, lft, img_top, avail_w, avail_h):
                """Insert picture centered within (avail_w × avail_h), preserving aspect ratio."""
                probe = cv2.imread(path)
                if probe is None:
                    slide.shapes.add_picture(path, lft, img_top,
                                             width=avail_w, height=avail_h)
                    return
                nat_h, nat_w = probe.shape[:2]
                if nat_h == 0 or nat_w == 0:
                    slide.shapes.add_picture(path, lft, img_top,
                                             width=avail_w, height=avail_h)
                    return
                ar_img = nat_w / nat_h
                ar_box = avail_w / avail_h
                if ar_img >= ar_box:
                    disp_w = avail_w
                    disp_h = int(avail_w / ar_img)
                else:
                    disp_h = avail_h
                    disp_w = int(avail_h * ar_img)
                offset_x = (avail_w - disp_w) // 2
                offset_y = (avail_h - disp_h) // 2
                slide.shapes.add_picture(
                    path,
                    lft + offset_x,
                    img_top + offset_y,
                    width=disp_w,
                    height=disp_h,
                )

            for key, label, lft, top_y in layout_slots:
                if key not in paths:
                    # Empty placeholder with label
                    ph = slide.shapes.add_shape(1, lft, top_y, IMG_W, LABEL_H + IMG_H)
                    ph.fill.solid()
                    ph.fill.fore_color.rgb = C_CARD
                    ph.line.fill.background()
                    _add_text(slide, f"{label}  [n/a]",
                              lft + Inches(0.08), top_y + Inches(0.05),
                              IMG_W - Inches(0.1), LABEL_H,
                              size=9, italic=True, color=C_TEXT_SEC)
                    continue

                # Label bar
                lbl_bg = slide.shapes.add_shape(1, lft, top_y, IMG_W, LABEL_H)
                lbl_bg.fill.solid()
                lbl_bg.fill.fore_color.rgb = C_CARD
                lbl_bg.line.fill.background()
                _add_text(slide, label,
                          lft + Inches(0.08), top_y + Inches(0.01),
                          IMG_W - Inches(0.1), LABEL_H,
                          size=9, bold=True, color=C_PRIMARY)
                # Image — inserted with aspect-ratio preservation
                _add_picture_aspect(slide, paths[key], lft, top_y + LABEL_H, IMG_W, IMG_H)

        # ── ROI Analysis slides ──────────────────────────────────────────────
        if roi_full_results and roi_all_results:
            if _s('roi_summary_table'):
                _prog("Generating ROI Pair Summary Table…")
                _ppt_add_roi_slides(prs, roi_full_results, roi_all_results,
                                    _fill_bg, _add_text, _score_color,
                                    C_BG, C_CARD, C_TEXT, C_TEXT_SEC,
                                    C_PRIMARY, C_SUCCESS, C_WARN,
                                    Inches, Pt)

            if _s('snr_matrix'):
                _prog("Generating SNR Pair Matrix…")
                _ppt_add_matrix_slide(prs, roi_full_results, roi_all_results,
                                      _fill_bg, _add_text,
                                      C_BG, C_CARD, C_TEXT, C_TEXT_SEC, C_PRIMARY,
                                      Inches, Pt)

            if _s('roi_position_maps'):
                _prog("Generating ROI Position Maps…")
                _ppt_add_roi_position_slides(prs, roi_full_results, self._images,
                                             _fill_bg, _add_text,
                                             C_BG, C_CARD, C_TEXT, C_TEXT_SEC, C_PRIMARY,
                                             Inches, Pt)

            if _s('diff_roi_overlay'):
                _prog("Generating Diff Map + ROI Overlay…")
                _ppt_add_diff_roi_position_slides(prs, roi_full_results, roi_all_results,
                                                  _fill_bg, _add_text,
                                                  C_BG, C_CARD, C_TEXT, C_TEXT_SEC, C_PRIMARY,
                                                  Inches, Pt)

            if _s('intensity_profiles'):
                _prog("Generating Intensity Profile Charts…")
                _ppt_add_roi_profile_slides(prs, roi_full_results,
                                            _fill_bg, _add_text,
                                            C_BG, C_CARD, C_TEXT, C_TEXT_SEC, C_PRIMARY,
                                            Inches, Pt)

        # ── Image Gallery — By Condition ──────────────────────────────────────
        if _s('image_gallery'):
            _prog("Generating Image Gallery…")
            _ppt_add_condition_gallery(prs, result_rows, crop_size,
                                       _fill_bg, _add_text,
                                       C_BG, C_CARD, C_TEXT, C_TEXT_SEC, C_PRIMARY,
                                       Inches, Pt)

        # ── Image Gallery — Base SNR Summary ─────────────────────────────────
        if roi_full_results and _s('base_snr_chart'):
            _prog("Generating Base SNR Summary Chart…")
            _ppt_add_base_snr_gallery_slide(prs, roi_full_results,
                                            _fill_bg, _add_text,
                                            C_BG, C_CARD, C_TEXT, C_TEXT_SEC, C_PRIMARY,
                                            Inches, Pt)

        _prog("Saving PowerPoint file…")
        ppt_path = str(Path(out_dir) / "perspective_report.pptx")
        prs.save(ppt_path)
        return ppt_path

    def _on_export(self):
        """Export results to PNG files."""
        # ── Quadrant Fusion export ────────────────────────────────────────
        if self._qf_last_result is not None and self.cmb_input_mode.currentIndex() == 1:
            self._on_export_quadrant_fusion()
            return

        if not self._results and self._result is None:
            return

        import csv
        import os

        def _safe_name(text: str) -> str:
            return "".join(c if c.isalnum() or c in ("-", "_") else "_" for c in text)

        if self._results:
            out_dir = QtWidgets.QFileDialog.getExistingDirectory(
                self, "Export Results Folder"
            )
            if not out_dir:
                return
            try:
                summary_path = os.path.join(out_dir, "perspective_summary.csv")
                export_dirs = {
                    "base": os.path.join(out_dir, "base"),
                    "compare": os.path.join(out_dir, "compare"),
                    "diff": os.path.join(out_dir, "diff"),
                    "snr": os.path.join(out_dir, "snr"),
                    "normalized": os.path.join(out_dir, "normalized"),
                    "aligned": os.path.join(out_dir, "aligned"),
                }
                for folder_path in export_dirs.values():
                    os.makedirs(folder_path, exist_ok=True)

                settings = self._last_settings or {}
                # Read export options (includes PPT section selections)
                do_center_crop, crop_size, export_gif, ppt_sections = \
                    self._ask_export_options()
                if ppt_sections is None:
                    return   # user cancelled

                # ── Progress dialog ───────────────────────────────────
                n_pairs = len(self._results)
                has_roi = bool(self._roi_full_results)
                _sec = ppt_sections or {}
                ppt_step_count = (
                    1                                          # title
                    + (1 if _sec.get('overview_table')   else 0)
                    + (n_pairs if _sec.get('detail_slides') else 0)
                    + (1 if _sec.get('roi_summary_table') and has_roi else 0)
                    + (1 if _sec.get('snr_matrix')        and has_roi else 0)
                    + (1 if _sec.get('roi_position_maps') and has_roi else 0)
                    + (1 if _sec.get('diff_roi_overlay')  and has_roi else 0)
                    + (1 if _sec.get('intensity_profiles') and has_roi else 0)
                    + (1 if _sec.get('image_gallery')     else 0)
                    + (1 if _sec.get('base_snr_chart')    and has_roi else 0)
                    + 1                                        # save file
                )
                total_steps = (
                    1                        # folder setup
                    + n_pairs                # image saving
                    + (n_pairs if export_gif else 0)
                    + ppt_step_count
                )
                prog = QtWidgets.QProgressDialog(
                    "Preparing export…", "Cancel", 0, total_steps, self
                )
                prog.setWindowTitle("Exporting…")
                prog.setWindowModality(Qt.WindowModal)
                prog.setMinimumDuration(0)
                prog.setStyleSheet(DIALOG_STYLE)
                prog.setValue(0)

                def _step(label: str) -> bool:
                    """Advance progress by 1 step. Returns True if cancelled."""
                    prog.setLabelText(label)
                    prog.setValue(prog.value() + 1)
                    QtWidgets.QApplication.processEvents()
                    return prog.wasCanceled()

                _step("Setting up export folders…")

                csv_headers = [
                    "base", "compare", "operation", "align_method", "dx", "dy",
                    "align_score", "diff_mean", "diff_std", "hot_pixels",
                    "invert_base", "invert_compare", "invert_result",
                    "alpha", "beta", "normalize", "glv_range", "norm_a", "norm_b",
                    "subtract_mode", "preserve_positive_diff", "abs_no_gain",
                    "center_crop", "crop_size",
                ]
                result_rows: List[Dict[str, object]] = []
                with open(summary_path, "w", encoding="utf-8", newline="") as f:
                    writer = csv.writer(f)
                    writer.writerow(csv_headers)
                    for result in self._results:
                        if _step(f"Saving images  [{result.base_label} → {result.compare_label}]"):
                            return
                        base_name = _safe_name(result.base_label)
                        cmp_name = _safe_name(result.compare_label)
                        op_name = _safe_name(result.operation)
                        prefix = f"{base_name}__{op_name}__{cmp_name}"

                        saved_paths = self._save_export_images_for_result(
                            result, export_dirs, prefix,
                            center_crop=do_center_crop,
                            crop_size=crop_size,
                        )
                        result_rows.append({"result": result, "paths": saved_paths})

                        s = result.stats
                        writer.writerow([
                            result.base_label, result.compare_label, result.operation,
                            result.alignment.method, result.alignment.dx, result.alignment.dy,
                            f"{s.get('alignment_score', 0):.2f}",
                            f"{s.get('diff_mean', 0):.6f}", f"{s.get('diff_std', 0):.6f}",
                            s.get("hot_pixels", 0),
                            settings.get("invert_base", False),
                            settings.get("invert_compare", False),
                            settings.get("invert_result", False),
                            settings.get("alpha", result.blend_alpha),
                            settings.get("beta", result.blend_beta),
                            bool(settings.get("normalize", True)),
                            str(settings.get("glv_range", "")),
                            f"{result.norm_a:.6f}", f"{result.norm_b:.6f}",
                            settings.get("subtract_mode", 0),
                            bool(settings.get("preserve_positive_diff", False)),
                            bool(settings.get("abs_no_gain", False)),
                            do_center_crop, crop_size,
                        ])

                # ── Animated GIF export (before PPT so frames can be embedded) ─
                gif_note = ""
                gif_dir = os.path.join(out_dir, "gif")
                if export_gif:
                    os.makedirs(gif_dir, exist_ok=True)
                    gif_count = 0
                    try:
                        from PIL import Image as _PilImage

                        def _maybe_crop_gif(img):
                            if do_center_crop and img is not None:
                                return self._center_crop_image(img, crop_size)
                            return img

                        for row in result_rows:
                            result = row["result"]
                            if _step(f"Creating GIF  [{result.base_label} → {result.compare_label}]"):
                                break
                            base_img_g = self._images.get(result.base_label)
                            norm_img_g = result.normalized_compare
                            diff_img_g = result.result_image

                            frames_np = [
                                _maybe_crop_gif(base_img_g),
                                _maybe_crop_gif(norm_img_g),
                                _maybe_crop_gif(diff_img_g),
                            ]
                            pil_frames = []
                            for fr in frames_np:
                                if fr is None:
                                    continue
                                if fr.dtype != np.uint8:
                                    fr = np.clip(fr, 0, 255).astype(np.uint8)
                                if fr.ndim == 2:
                                    pil_frames.append(
                                        _PilImage.fromarray(fr, mode="L").convert("RGB"))
                                else:
                                    pil_frames.append(_PilImage.fromarray(
                                        cv2.cvtColor(fr, cv2.COLOR_BGR2RGB)))

                            if len(pil_frames) >= 2:
                                base_n = _safe_name(result.base_label)
                                cmp_n = _safe_name(result.compare_label)
                                gif_path = os.path.join(gif_dir, f"{base_n}__{cmp_n}.gif")
                                pil_frames[0].save(
                                    gif_path,
                                    save_all=True,
                                    append_images=pil_frames[1:],
                                    loop=0,
                                    duration=700,
                                    optimize=False,
                                )
                                # Store gif path so PPT can embed it
                                row["paths"]["gif"] = gif_path
                                gif_count += 1

                        gif_note = f"\n• {gif_count} GIF(s) → {gif_dir}"
                    except ImportError:
                        gif_note = "\n⚠ GIF skipped — Pillow not installed (pip install Pillow)"
                    except Exception as gif_err:
                        gif_note = f"\n⚠ GIF error: {gif_err}"

                ppt_path = self._export_ppt_report(
                    out_dir, result_rows, settings,
                    do_center_crop=do_center_crop,
                    crop_size=crop_size,
                    roi_full_results=self._roi_full_results or {},
                    roi_all_results=self._results or [],
                    ppt_sections=ppt_sections,
                    progress_callback=_step,
                )
                prog.setValue(total_steps)
                if ppt_path:
                    QtWidgets.QMessageBox.information(
                        self,
                        "Export Complete",
                        f"Saved:\n• {out_dir}\n• {summary_path}\n• {ppt_path}{gif_note}",
                    )
                else:
                    QtWidgets.QMessageBox.information(
                        self,
                        "Export Complete",
                        f"Saved:\n• {out_dir}\n• {summary_path}\n\n"
                        f"PPT report was skipped because python-pptx is not available.{gif_note}",
                    )
            except Exception as e:
                QtWidgets.QMessageBox.critical(self, "Export Error", str(e))
            return

        path, _ = QtWidgets.QFileDialog.getSaveFileName(
            self, "Export Difference Map", "perspective_diff.png",
            "PNG (*.png);;TIFF (*.tif)"
        )

        if not path:
            return

        try:
            # Save difference map
            cv2.imwrite(path, self._result.diff_image)

            # Save SNR map
            snr_path = path.replace(".png", "_snr.png").replace(".tif", "_snr.tif")
            snr_color = colorize_snr_map(self._result.snr_map)
            cv2.imwrite(snr_path, snr_color)

            # Save normalized compare image
            norm_path = path.replace(".png", "_normalized.png").replace(".tif", "_normalized.tif")
            saved_paths = [path, snr_path]
            if self._result.normalized_compare is not None:
                cv2.imwrite(norm_path, self._result.normalized_compare)
                saved_paths.append(norm_path)

            # Save aligned compare image
            aligned_path = path.replace(".png", "_aligned.png").replace(".tif", "_aligned.tif")
            if self._result.aligned_compare is not None:
                cv2.imwrite(aligned_path, self._result.aligned_compare)
                saved_paths.append(aligned_path)

            QtWidgets.QMessageBox.information(
                self, "Export Complete",
                "Saved:\n" + "\n".join(f"• {saved}" for saved in saved_paths)
            )
        except Exception as e:
            QtWidgets.QMessageBox.critical(self, "Export Error", str(e))

    def _on_export_quadrant_fusion(self):
        """Export Quadrant Fusion result images."""
        import os

        result = self._qf_last_result
        if result is None:
            return

        out_dir = QtWidgets.QFileDialog.getExistingDirectory(
            self, "Export Quadrant Fusion Results"
        )
        if not out_dir:
            return

        try:
            paths = []
            failed = []
            for name, img in [
                ("bse_enhanced", result.bse_clean_uint8),
                ("topography", result.topo_uint8),
                ("composite", result.composite_uint8),
                ("fused_output", result.fused_image_uint8),
            ]:
                if img is None or img.size == 0:
                    failed.append(f"{name}: image is empty")
                    continue
                # Ensure contiguous uint8
                out_img = np.ascontiguousarray(img)
                if out_img.dtype != np.uint8:
                    out_img = np.clip(out_img, 0, 255).astype(np.uint8)
                p = os.path.join(out_dir, f"quadrant_fusion_{name}.png")
                ok = cv2.imwrite(p, out_img)
                if ok:
                    paths.append(p)
                else:
                    failed.append(f"{name}: cv2.imwrite failed for {p}")

            # Write a small metadata text file
            meta_path = os.path.join(out_dir, "quadrant_fusion_info.txt")
            with open(meta_path, "w", encoding="utf-8") as f:
                f.write("Quadrant Fusion Result\n")
                f.write(f"Alpha used: {result.alpha_used:.6f}\n")
                f.write(f"Beta used: {result.beta_used:.4f}\n")
                f.write(f"Notes: {result.notes}\n")
                for det_name, cmb in self._qf_combos.items():
                    f.write(f"{det_name}: {cmb.currentText()}\n")
            paths.append(meta_path)

            msg = "Saved:\n" + "\n".join(f"  {p}" for p in paths)
            if failed:
                msg += "\n\nWarnings:\n" + "\n".join(f"  {f}" for f in failed)
            QtWidgets.QMessageBox.information(self, "Export Complete", msg)
        except Exception as e:
            QtWidgets.QMessageBox.critical(self, "Export Error", str(e))


__all__ = ['PerspectiveCombinationDialog']
